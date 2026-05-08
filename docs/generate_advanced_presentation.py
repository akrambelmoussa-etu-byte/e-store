# -*- coding: utf-8 -*-
"""
Genere une presentation E-Store inspiree du style NotebookLM :
  - presentation-avancee.pptx        (30 slides riches en schemas)
  - script-notebooklm-part-1.pdf     (slides 1-10 : Intro & Architecture)
  - script-notebooklm-part-2.pdf     (slides 11-20 : Backend & Securite)
  - script-notebooklm-part-3.pdf     (slides 21-30 : Frontend, Demo & Conclusion)

Le script PDF est concu pour etre injecte dans NotebookLM, qui ne peut
generer que 10 slides a la fois - d'ou le decoupage en 3 parties.
"""
import os
import sys
import io

if sys.stdout.encoding and sys.stdout.encoding.lower() != "utf-8":
    try:
        sys.stdout.reconfigure(encoding="utf-8", errors="replace")
    except Exception:
        sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import cm
from reportlab.lib.colors import HexColor, black, white, grey
from reportlab.lib.enums import TA_LEFT, TA_CENTER, TA_JUSTIFY
from reportlab.platypus import (
    SimpleDocTemplate, Paragraph, Spacer, PageBreak,
    Table, TableStyle, Preformatted, KeepTogether
)

OUT_DIR = os.path.dirname(os.path.abspath(__file__))

# =====================================================================
# PALETTE NOTEBOOKLM
# =====================================================================
BG          = RGBColor(0xE8, 0xEE, 0xF2)   # arriere-plan tres clair
DARK        = RGBColor(0x1F, 0x29, 0x37)   # texte titre
GRAY        = RGBColor(0x6B, 0x72, 0x80)
LIGHT_GRAY  = RGBColor(0xF5, 0xF7, 0xFA)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)

# Accents (extraits du NotebookLM)
MINT       = RGBColor(0x6F, 0xC1, 0x8E)
TEAL       = RGBColor(0x4D, 0xB1, 0xC4)
SKY        = RGBColor(0x5B, 0xA3, 0xCB)
LAVENDER   = RGBColor(0x9B, 0x82, 0xB7)
PEACH      = RGBColor(0xE8, 0x94, 0x6C)
CORAL      = RGBColor(0xD8, 0x6A, 0x6A)
GOLD       = RGBColor(0xE6, 0xB7, 0x4F)

MINT_HEX     = "#6FC18E"
TEAL_HEX     = "#4DB1C4"
SKY_HEX      = "#5BA3CB"
LAVENDER_HEX = "#9B82B7"
PEACH_HEX    = "#E8946C"
CORAL_HEX    = "#D86A6A"
GOLD_HEX     = "#E6B74F"
DARK_HEX     = "#1F2937"
BG_HEX       = "#E8EEF2"


# =====================================================================
# HELPERS PPTX
# =====================================================================
def set_slide_bg(slide, color=BG):
    """Definit la couleur de fond de la slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_grid_pattern(slide, prs):
    """Dessine un quadrillage tres leger pour evoquer un blueprint."""
    cols, rows = 30, 18
    cw = prs.slide_width / cols
    rh = prs.slide_height / rows
    grid_color = RGBColor(0xD0, 0xD8, 0xE0)
    for c in range(cols + 1):
        x = int(c * cw)
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x, 0, x, prs.slide_height)
        line.line.color.rgb = grid_color
        line.line.width = Emu(2540)  # ~0.2pt
    for r in range(rows + 1):
        y = int(r * rh)
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, 0, y, prs.slide_width, y)
        line.line.color.rgb = grid_color
        line.line.width = Emu(2540)


def add_text(slide, x, y, w, h, text, *, size=14, bold=False, italic=False,
             color=DARK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font="Calibri"):
    """Ajoute un textbox avec mise en forme."""
    tx = slide.shapes.add_textbox(x, y, w, h)
    tf = tx.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(0); tf.margin_right = Pt(0)
    tf.margin_top = Pt(0); tf.margin_bottom = Pt(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.text = text
    for r in p.runs:
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
    return tx


def pill(slide, x, y, w, h, text, *, fill=WHITE, accent=None, accent_pos="left",
         text_color=DARK, size=12, bold=False, font="Calibri", align=PP_ALIGN.CENTER):
    """Boite en pilule (rounded rectangle) avec accent de couleur optionnel."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.adjustments[0] = 0.35
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = RGBColor(0xC8, 0xCF, 0xD6)
    shape.line.width = Pt(0.75)

    if accent:
        if accent_pos == "left":
            bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                         x, y, Emu(80000), h)
            bar.adjustments[0] = 0.5
        else:  # top
            bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                         x + Emu(50000), y,
                                         w - Emu(100000), Emu(50000))
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent
        bar.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Pt(8); tf.margin_right = Pt(8)
    tf.margin_top = Pt(4); tf.margin_bottom = Pt(4)
    p = tf.paragraphs[0]
    p.alignment = align
    p.text = text
    for r in p.runs:
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = text_color
    return shape


def card(slide, x, y, w, h, title, body, *, accent=TEAL, accent_pos="left",
         icon=None):
    """Carte avec titre + paragraphe."""
    bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    bg.adjustments[0] = 0.08
    bg.fill.solid(); bg.fill.fore_color.rgb = WHITE
    bg.line.color.rgb = RGBColor(0xD0, 0xD8, 0xE0)
    bg.line.width = Pt(0.75)

    # Bande de couleur a gauche
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 x, y + Emu(80000), Emu(80000), h - Emu(160000))
    bar.fill.solid(); bar.fill.fore_color.rgb = accent
    bar.line.fill.background()

    # Icone (texte unicode dans un petit cercle)
    icon_w = Emu(450000)
    icon_x = x + Emu(180000)
    icon_y = y + (h - icon_w) // 2
    if icon:
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, icon_x, icon_y, icon_w, icon_w)
        circle.fill.solid(); circle.fill.fore_color.rgb = LIGHT_GRAY
        circle.line.color.rgb = accent
        circle.line.width = Pt(1.5)
        tf = circle.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        p.text = icon
        for r in p.runs:
            r.font.size = Pt(18); r.font.color.rgb = accent

    # Titre + texte
    tx_x = icon_x + icon_w + Emu(150000) if icon else x + Emu(250000)
    tx_w = x + w - tx_x - Emu(150000)
    tx = slide.shapes.add_textbox(tx_x, y + Emu(120000), tx_w, h - Emu(240000))
    tf = tx.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(0); tf.margin_right = Pt(0)
    tf.margin_top = Pt(0); tf.margin_bottom = Pt(0)

    p1 = tf.paragraphs[0]
    p1.text = title
    for r in p1.runs:
        r.font.name = "Calibri"; r.font.size = Pt(13); r.font.bold = True
        r.font.color.rgb = DARK

    p2 = tf.add_paragraph()
    p2.text = body
    p2.space_before = Pt(2)
    for r in p2.runs:
        r.font.name = "Calibri"; r.font.size = Pt(11); r.font.color.rgb = DARK


def slide_title(slide, text, size=36, color=DARK, x=Inches(0.6), y=Inches(0.4),
                w=Inches(12), h=Inches(1)):
    add_text(slide, x, y, w, h, text, size=size, bold=True,
             color=color, font="Calibri")


def add_arrow(slide, x1, y1, x2, y2, color=GRAY, width=1.25):
    """Ajoute une fleche directe entre deux points."""
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    conn.line.color.rgb = color
    conn.line.width = Pt(width)
    # Tete de fleche
    line = conn.line._get_or_add_ln()
    head = etree.SubElement(line, qn('a:tailEnd'))
    head.set('type', 'triangle')
    head.set('w', 'med'); head.set('len', 'med')
    return conn


def numbered_circle(slide, x, y, n, color=TEAL, size_emu=Emu(380000)):
    """Cercle numerote."""
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, size_emu, size_emu)
    c.fill.solid(); c.fill.fore_color.rgb = color
    c.line.color.rgb = WHITE; c.line.width = Pt(2)
    tf = c.text_frame
    tf.margin_left = Pt(0); tf.margin_right = Pt(0)
    tf.margin_top = Pt(0); tf.margin_bottom = Pt(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    p.text = str(n)
    for r in p.runs:
        r.font.bold = True; r.font.size = Pt(14); r.font.color.rgb = WHITE
    return c


def add_footer(slide, prs, page, total):
    add_text(slide, Inches(0.5), prs.slide_height - Inches(0.32),
             prs.slide_width - Inches(1), Inches(0.3),
             f"E-STORE  ·  A. Belmoussa & N. Ben Soumane  ·  Pr. O. Zahour  ·  {page}/{total}",
             size=9, color=GRAY, align=PP_ALIGN.CENTER)


def blank_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s, BG)
    return s


# =====================================================================
# DEFINITION DES 30 SLIDES
# =====================================================================
TOTAL_SLIDES = 30


# ── PART 1 — Introduction & Architecture (slides 1-10) ─────────────
def slide_01_cover(prs, page):
    s = blank_slide(prs)
    add_grid_pattern(s, prs)

    # Titre geant
    add_text(s, Inches(0.5), Inches(2.4), Inches(12.3), Inches(2),
             "E-STORE", size=140, bold=True, color=DARK,
             align=PP_ALIGN.CENTER, font="Calibri")

    # Sous-titre
    add_text(s, Inches(0.5), Inches(4.4), Inches(12.3), Inches(0.5),
             "Ingenierie d'une Plateforme E-Commerce Full-Stack :",
             size=22, color=DARK, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(4.85), Inches(12.3), Inches(0.5),
             "Conception, Persistance Hybride et Securite",
             size=22, color=DARK, align=PP_ALIGN.CENTER)

    # 3 pilules en bas
    pw = Inches(3.8); ph = Inches(0.95); spacing = Inches(0.3)
    total_w = pw * 3 + spacing * 2
    start_x = (prs.slide_width - total_w) // 2
    y = Inches(6.2)

    pill(s, start_x, y, pw, ph,
         "Auteurs :\nAkram Belmoussa & Nouhaila Ben Soumane",
         fill=WHITE, size=11, font="Consolas", align=PP_ALIGN.CENTER)
    add_text(s, start_x + Inches(0.4), y + Inches(0.1), pw - Inches(0.8), Inches(0.3),
             "Auteurs :", size=10, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

    pill(s, start_x + pw + spacing, y, pw, ph,
         "Supervision academique :\nPr. Omar Zahour",
         fill=WHITE, size=11, font="Consolas", align=PP_ALIGN.CENTER)

    pill(s, start_x + 2 * (pw + spacing), y, pw, ph,
         "Etablissement :\nFSBM - Universite Hassan II\n(Annee 2025-2026)",
         fill=WHITE, size=10, font="Consolas", align=PP_ALIGN.CENTER)

    add_footer(s, prs, page, TOTAL_SLIDES)


def slide_02_mandat(prs, page):
    s = blank_slide(prs)

    add_text(s, Inches(0.5), Inches(0.5), Inches(5), Inches(1.6),
             "Mandat\nd'Ingenierie", size=44, bold=True, color=DARK)

    # Carte gauche : Contexte projet
    cx, cy, cw, ch = Inches(0.5), Inches(2.4), Inches(5.5), Inches(4.5)
    bg = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, cy, cw, ch)
    bg.adjustments[0] = 0.04
    bg.fill.solid(); bg.fill.fore_color.rgb = WHITE
    bg.line.color.rgb = RGBColor(0xD0, 0xD8, 0xE0); bg.line.width = Pt(0.75)

    add_text(s, cx + Inches(0.3), cy + Inches(0.25), cw - Inches(0.6), Inches(0.5),
             "Contexte Projet", size=22, bold=True, color=DARK)
    add_text(s, cx + Inches(0.3), cy + Inches(0.95), cw - Inches(0.6), Inches(2.5),
             "Mini-projet pedagogique e-commerce.\n"
             "Perimetre d'application complet :\n"
             "Inscription, catalogue, panier, commandes,\n"
             "et gestion des avis.",
             size=13, color=DARK)

    # Petit schema isometrique en bas
    icons_x = cx + Inches(0.6); icons_y = cy + Inches(2.5)
    items = ["Perimetre d'application", "Inscription", "Catalogue",
             "Panier", "Commandes", "Gestion des avis"]
    for i, it in enumerate(items):
        bx = icons_x + Inches(0.3) * i
        by = icons_y + Inches(0.18) * i
        sq = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, bx, by,
                                 Inches(0.32), Inches(0.32))
        sq.fill.solid(); sq.fill.fore_color.rgb = LIGHT_GRAY
        sq.line.color.rgb = GRAY; sq.line.width = Pt(0.75)
        add_text(s, bx + Inches(0.4), by + Inches(0.05),
                 Inches(2.6), Inches(0.3), it, size=9, color=DARK)

    # Cartes objectifs a droite
    objs = [
        ("Objectif 1", "Maitriser une architecture full-stack moderne.", MINT),
        ("Objectif 2", "Illustrer la puissance d'une persistance hybride SQL + NoSQL.", LAVENDER),
        ("Objectif 3", "Appliquer les bonnes pratiques (DTOs, @Transactional, tests unitaires).", GRAY),
        ("Objectif 4", "Verrouiller les acces via JWT et Spring Security 6.", PEACH),
        ("Objectif 5", "Produire une UI reactive avec Angular 17 en mode standalone.", TEAL),
    ]
    ox = Inches(6.5); oy = Inches(1.0); ow = Inches(6.4); oh = Inches(0.95)
    for i, (lbl, txt, col) in enumerate(objs):
        y = oy + (oh + Inches(0.18)) * i
        # Carte blanche
        bg = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, ox, y, ow, oh)
        bg.adjustments[0] = 0.18
        bg.fill.solid(); bg.fill.fore_color.rgb = WHITE
        bg.line.color.rgb = RGBColor(0xD0, 0xD8, 0xE0); bg.line.width = Pt(0.5)
        # Bande couleur a gauche
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  ox, y + Inches(0.1), Inches(0.1), oh - Inches(0.2))
        bar.fill.solid(); bar.fill.fore_color.rgb = col
        bar.line.fill.background()
        # Icone (cercle)
        ic = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                 ox + Inches(0.3), y + Inches(0.22),
                                 Inches(0.5), Inches(0.5))
        ic.fill.solid(); ic.fill.fore_color.rgb = LIGHT_GRAY
        ic.line.color.rgb = col; ic.line.width = Pt(1.2)
        # Texte
        add_text(s, ox + Inches(1), y + Inches(0.18),
                 ow - Inches(1.2), Inches(0.4),
                 lbl + " :", size=12, bold=True, color=DARK)
        add_text(s, ox + Inches(1), y + Inches(0.45),
                 ow - Inches(1.2), Inches(0.5),
                 txt, size=11, color=DARK)

    add_footer(s, prs, page, TOTAL_SLIDES)


def slide_03_plan(prs, page):
    s = blank_slide(prs)
    slide_title(s, "Plan de la Presentation", size=36)

    parts = [
        ("PARTIE 1 - Fondations", MINT, [
            "1. Couverture", "2. Mandat d'ingenierie", "3. Plan",
            "4. Cahier des charges", "5. Matrice architecturale",
            "6. Couches techniques", "7. Domaines DDD",
            "8. Ecosysteme tech", "9. Stack frontend", "10. Stack backend",
        ]),
        ("PARTIE 2 - Backend & Securite", LAVENDER, [
            "11. Modele de donnees", "12. Domaine customer",
            "13. Domaine catalog", "14. Domaine inventory",
            "15. Domaine shopping", "16. Domaine billing",
            "17. Persistance hybride", "18. Domaine review (Mongo)",
            "19. Cycle JWT", "20. Endpoints REST",
        ]),
        ("PARTIE 3 - Frontend, Demo, Conclusion", PEACH, [
            "21. Facade Angular 17", "22. Architecture frontend",
            "23. Signals & services", "24. Interceptors & guards",
            "25. Parcours utilisateur", "26. Tests unitaires",
            "27. Difficulte Docker", "28. Bilan",
            "29. Feuille de route", "30. Q & R",
        ]),
    ]
    pw = Inches(4.0); ph = Inches(5.5); spacing = Inches(0.25)
    total_w = pw * 3 + spacing * 2
    start_x = (prs.slide_width - total_w) // 2
    y = Inches(1.5)

    for i, (title, color, items) in enumerate(parts):
        x = start_x + (pw + spacing) * i
        # Carte
        bg = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, pw, ph)
        bg.adjustments[0] = 0.04
        bg.fill.solid(); bg.fill.fore_color.rgb = WHITE
        bg.line.color.rgb = RGBColor(0xD0, 0xD8, 0xE0); bg.line.width = Pt(0.5)
        # Bandeau colore en haut
        top_bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      x + Inches(0.15), y, pw - Inches(0.3), Inches(0.08))
        top_bar.fill.solid(); top_bar.fill.fore_color.rgb = color
        top_bar.line.fill.background()
        # Titre
        add_text(s, x + Inches(0.3), y + Inches(0.3), pw - Inches(0.6), Inches(0.5),
                 title, size=14, bold=True, color=DARK, align=PP_ALIGN.CENTER)
        # Liste
        for j, it in enumerate(items):
            add_text(s, x + Inches(0.4), y + Inches(0.95) + Inches(0.4) * j,
                     pw - Inches(0.8), Inches(0.4),
                     it, size=11, color=DARK)

    add_footer(s, prs, page, TOTAL_SLIDES)


def slide_04_specifications(prs, page):
    s = blank_slide(prs)
    slide_title(s, "Cahier des charges fonctionnel")

    items = [
        ("Inscription & connexion",
         "Compte utilisateur avec email unique, BCrypt, role USER/ADMIN.", MINT),
        ("Catalogue produits",
         "Categories, recherche par mot-cle, filtre, pagination, fiche detaillee.", TEAL),
        ("Panier",
         "Ajout/modification/suppression d'articles, prix fige, controle stock.", SKY),
        ("Commandes",
         "Validation transactionnelle, decrement stock, historique consultable.", LAVENDER),
        ("Avis (Reviews)",
         "Note 1-5 + commentaire stocke en MongoDB, lecture publique.", PEACH),
        ("Administration",
         "CRUD produits/categories reserve aux ADMIN via @PreAuthorize.", CORAL),
    ]
    cw = Inches(6.2); ch = Inches(1.4)
    for i, (title, body, color) in enumerate(items):
        col = i % 2; row = i // 2
        x = Inches(0.4) + col * (cw + Inches(0.3))
        y = Inches(1.5) + row * (ch + Inches(0.25))
        card(s, x, y, cw, ch, title, body, accent=color)

    add_footer(s, prs, page, TOTAL_SLIDES)


def slide_05_matrix(prs, page):
    s = blank_slide(prs)
    slide_title(s, "Matrice Architecturale : 3 Couches x 6 Domaines",
                size=32)

    # Headers en pilules (en haut)
    headers = [
        ("PRESENTATION (Angular)", TEAL),
        ("LOGIQUE METIER (Spring Boot)", MINT),
        ("ACCES AUX DONNEES", LAVENDER),
    ]
    col_x = [Inches(2.5), Inches(6.5), Inches(10.5)]
    cw = Inches(3.5)
    for (text, color), x in zip(headers, col_x):
        # Bandeau couleur
        bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  x, Inches(1.4), cw, Inches(0.1))
        bar.adjustments[0] = 0.5
        bar.fill.solid(); bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        # Texte
        add_text(s, x, Inches(1.5), cw, Inches(0.4), text,
                 size=11, bold=True, color=DARK, align=PP_ALIGN.CENTER)

    # Lignes
    rows = [
        ("customer",  "Composants Auth",  "User, Profile, JWT",          "JPA"),
        ("catalog",   "Services HTTP",     "Categorie, Produit",          "JPA"),
        ("inventory", "UI Stock",          "Stock par produit",           "JPA"),
        ("shopping",  "Etat Signals",      "Cart, CartItem",              "JPA"),
        ("billing",   "UI Commandes",      "Order (@Transactional)",      "JPA"),
        ("review",    "Affichage Avis",    "Avis Metier",                  "MongoDB"),
    ]
    rh = Inches(0.55)
    for i, (dom, p, l, d) in enumerate(rows):
        y = Inches(2.1) + rh * i
        # Label domaine (gauche)
        add_text(s, Inches(0.3), y + Inches(0.12), Inches(2), Inches(0.4),
                 dom, size=12, color=GRAY, font="Consolas", italic=True)
        # 3 cellules pilules
        for j, txt in enumerate([p, l, d]):
            color = WHITE if (d != "MongoDB" or j != 2) else LAVENDER
            text_color = DARK if color == WHITE else WHITE
            pill(s, col_x[j] + Inches(0.2), y + Inches(0.05),
                 cw - Inches(0.4), Inches(0.45),
                 txt, fill=color, text_color=text_color,
                 size=10, font="Consolas", align=PP_ALIGN.CENTER)

    add_footer(s, prs, page, TOTAL_SLIDES)


def slide_06_layers(prs, page):
    s = blank_slide(prs)
    slide_title(s, "Justification : Les 3 Couches Techniques")

    # 3 grandes pilules verticales
    layers = [
        ("PRESENTATION",
         "Angular 17 standalone\nBootstrap 5 + signals\nResponsabilites :\nAffichage UI - validation client", TEAL),
        ("LOGIQUE METIER",
         "Spring Boot 3.3\nControllers + Services + DTOs\nResponsabilites :\nRegles metier - orchestration\n@Transactional - securite", MINT),
        ("ACCES AUX DONNEES",
         "Spring Data JPA\nSpring Data MongoDB\nResponsabilites :\nPersistance - requetes\nAbstraction du SGBD", LAVENDER),
    ]
    cw = Inches(4); ch = Inches(4.2)
    spacing = Inches(0.3)
    total_w = cw * 3 + spacing * 2
    start_x = (prs.slide_width - total_w) // 2
    y = Inches(1.7)
    for i, (t, body, col) in enumerate(layers):
        x = start_x + (cw + spacing) * i
        # Pilule verticale
        shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, cw, ch)
        shape.adjustments[0] = 0.06
        shape.fill.solid(); shape.fill.fore_color.rgb = col
        shape.line.fill.background()
        # Titre (blanc)
        add_text(s, x, y + Inches(0.3), cw, Inches(0.6),
                 t, size=18, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, font="Calibri")
        # Corps (blanc)
        add_text(s, x + Inches(0.3), y + Inches(1.1), cw - Inches(0.6),
                 ch - Inches(1.4), body, size=12, color=WHITE,
                 align=PP_ALIGN.CENTER)

    # Fleches verticales entre couches
    add_text(s, Inches(0.5), Inches(6.1), Inches(12.3), Inches(0.3),
             "Communication HTTP/JSON  →  appels methodes  →  requetes SQL/Mongo",
             size=12, italic=True, color=GRAY, align=PP_ALIGN.CENTER)

    add_footer(s, prs, page, TOTAL_SLIDES)


def slide_07_domains(prs, page):
    s = blank_slide(prs)
    slide_title(s, "Justification : Les Domaines (DDD)")

    add_text(s, Inches(0.5), Inches(1.3), Inches(12), Inches(0.5),
             "Chaque domaine est un sous-package autonome (entity, dto, repository, service, controller)",
             size=14, italic=True, color=GRAY, align=PP_ALIGN.CENTER)

    domains = [
        ("customer",  "User, Profile, JWT, Auth",                MINT),
        ("catalog",   "Categorie, Produit, recherche",            TEAL),
        ("inventory", "Stock par produit",                         SKY),
        ("shopping",  "Cart, CartItem, gestion panier",            LAVENDER),
        ("billing",   "Order, OrderItem, checkout transactionnel", PEACH),
        ("review",    "Avis (NoSQL MongoDB)",                       CORAL),
    ]
    cw = Inches(6.0); ch = Inches(1.4)
    for i, (name, body, col) in enumerate(domains):
        c = i % 2; r = i // 2
        x = Inches(0.4) + c * (cw + Inches(0.4))
        y = Inches(2.1) + r * (ch + Inches(0.2))
        card(s, x, y, cw, ch, name, body, accent=col, icon="◆")

    add_footer(s, prs, page, TOTAL_SLIDES)


def slide_08_ecosystem(prs, page):
    s = blank_slide(prs)
    slide_title(s, "Cartographie de l'Ecosysteme Technologique", size=32)

    # Noyau central
    cx = prs.slide_width // 2 - Inches(1.4)
    cy = Inches(3.3)
    nucleus = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  cx, cy, Inches(2.8), Inches(1.8))
    nucleus.adjustments[0] = 0.1
    nucleus.fill.solid(); nucleus.fill.fore_color.rgb = DARK
    nucleus.line.fill.background()
    add_text(s, cx, cy + Inches(0.4), Inches(2.8), Inches(1),
             "Noyau\nApplication\nE-STORE",
             size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # 6 satellites
    sats = [
        ("Frontend", "Angular 17 standalone\nBootstrap 5\nRxJS\nTypeScript 5", TEAL,
         Inches(5.5), Inches(0.9)),
        ("Backend", "Spring Boot 3.3\nSpring Security 6\nJWT (jjwt 0.12)", MINT,
         Inches(0.5), Inches(2.7)),
        ("Persistance", "Spring Data JPA + Hibernate\nSpring Data MongoDB\nMySQL 8 / H2 / MongoDB 7", LAVENDER,
         Inches(5.5), Inches(5.5)),
        ("Build", "Maven 3.9\nnpm\nAngular CLI", SKY,
         Inches(10.5), Inches(0.9)),
        ("Tests", "JUnit 5 + Mockito\n9 tests, BUILD SUCCESS", MINT,
         Inches(10.5), Inches(5.5)),
        ("DevOps", "Docker Compose\nMySQL + Mongo\nphpMyAdmin\nmongo-express", LAVENDER,
         Inches(0.5), Inches(5)),
    ]
    nucleus_cx = cx + Inches(1.4); nucleus_cy = cy + Inches(0.9)
    for title, body, color, x, y in sats:
        # Carte
        cw = Inches(2.3); ch = Inches(1.4)
        bg = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, cw, ch)
        bg.adjustments[0] = 0.12
        bg.fill.solid(); bg.fill.fore_color.rgb = WHITE
        bg.line.color.rgb = color; bg.line.width = Pt(1.5)
        add_text(s, x, y + Inches(0.1), cw, Inches(0.4),
                 title, size=14, bold=True, color=DARK, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.15), y + Inches(0.5), cw - Inches(0.3), ch - Inches(0.5),
                 body, size=9, color=DARK, align=PP_ALIGN.CENTER, font="Consolas")
        # Connecteur vers le noyau
        sat_cx = x + cw // 2; sat_cy = y + ch // 2
        add_arrow(s, sat_cx, sat_cy, nucleus_cx, nucleus_cy, color=color, width=1)

    add_footer(s, prs, page, TOTAL_SLIDES)


def slide_09_frontend_stack(prs, page):
    s = blank_slide(prs)
    slide_title(s, "Stack Frontend en detail")

    items = [
        ("Angular 17", "Framework SPA, version standalone (pas de NgModule).\nTypeScript 5, RxJS 7, signals (nouveaute v17).", TEAL),
        ("Bootstrap 5", "Framework CSS responsive.\nGrille 12 col, composants (cards, modals, navbar).\nThematique cohesive sans JS Bootstrap.", LAVENDER),
        ("RxJS + Signals", "RxJS pour les flux HTTP (Observable).\nSignals pour l'etat reactif local (cart.itemCount).", MINT),
        ("Tooling", "Angular CLI + npm.\nLazy loading par route -> bundle initial 13 ko.\nAOT + tree-shaking en production.", SKY),
    ]
    cw = Inches(6); ch = Inches(2.4)
    for i, (t, b, col) in enumerate(items):
        c = i % 2; r = i // 2
        x = Inches(0.4) + c * (cw + Inches(0.4))
        y = Inches(1.5) + r * (ch + Inches(0.2))
        card(s, x, y, cw, ch, t, b, accent=col, icon="●")

    add_footer(s, prs, page, TOTAL_SLIDES)


def slide_10_backend_stack(prs, page):
    s = blank_slide(prs)
    slide_title(s, "Stack Backend & DevOps en detail")

    items = [
        ("Spring Boot 3.3", "Web (REST), Validation (Bean Validation),\nData JPA + Data MongoDB,\nSecurity 6, Lombok.", MINT),
        ("Persistance", "Hibernate (ORM JPA)\nMySQL 8 prod / H2 in-memory dev\nMongoDB 7 (avis)", LAVENDER),
        ("Securite", "Spring Security 6 stateless\nJWT HS256 (jjwt 0.12, 24h)\nBCrypt sur les mots de passe", PEACH),
        ("Build & Test", "Maven 3.9\nJUnit 5 + Mockito (9 tests)\nDocker Compose (4 services)", SKY),
    ]
    cw = Inches(6); ch = Inches(2.4)
    for i, (t, b, col) in enumerate(items):
        c = i % 2; r = i // 2
        x = Inches(0.4) + c * (cw + Inches(0.4))
        y = Inches(1.5) + r * (ch + Inches(0.2))
        card(s, x, y, cw, ch, t, b, accent=col, icon="◆")

    add_footer(s, prs, page, TOTAL_SLIDES)


# ── PART 2 — Backend & Securite (slides 11-20) ─────────────────────
def slide_11_data_model(prs, page):
    s = blank_slide(prs)
    slide_title(s, "Modele de donnees relationnel")

    # Boites entites (carte par entite)
    entities = [
        ("User",      ["+ id : Long (PK)", "+ email : String UK", "+ password : String", "+ role : Role"], Inches(0.5), Inches(1.5), MINT),
        ("Profile",   ["+ id : Long (PK)", "+ phone, address", "+ city, country"], Inches(4.5), Inches(1.5), MINT),
        ("Cart",      ["+ id : Long (PK)", "+ createdAt", "+ updatedAt"], Inches(8.5), Inches(1.5), TEAL),
        ("CartItem",  ["+ id : Long (PK)", "+ quantity", "+ unitPrice"], Inches(11.0), Inches(1.5), TEAL),
        ("Category",  ["+ id : Long (PK)", "+ name", "+ description"], Inches(0.5), Inches(4.3), LAVENDER),
        ("Product",   ["+ id : Long (PK)", "+ name, description", "+ price : BigDecimal", "+ imageUrl"], Inches(4), Inches(4.3), LAVENDER),
        ("Inventory", ["+ id : Long (PK)", "+ quantity : int"], Inches(7.7), Inches(4.3), LAVENDER),
        ("Order",     ["+ id : Long (PK)", "+ orderDate", "+ totalAmount", "+ status : OrderStatus"], Inches(10.2), Inches(4.3), PEACH),
        ("OrderItem", ["+ id : Long (PK)", "+ quantity", "+ unitPrice"], Inches(0.5), Inches(6.1), PEACH),
    ]
    for (name, fields, x, y, col) in entities:
        cw = Inches(2.5); ch = Inches(1.6)
        bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, cw, ch)
        bg.fill.solid(); bg.fill.fore_color.rgb = WHITE
        bg.line.color.rgb = col; bg.line.width = Pt(1.5)
        # Header
        hd = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, cw, Inches(0.32))
        hd.fill.solid(); hd.fill.fore_color.rgb = col
        hd.line.fill.background()
        add_text(s, x, y + Inches(0.04), cw, Inches(0.3),
                 name, size=11, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, font="Consolas")
        # Champs
        for j, f in enumerate(fields):
            add_text(s, x + Inches(0.1), y + Inches(0.36) + Inches(0.22) * j,
                     cw - Inches(0.2), Inches(0.22), f, size=8.5,
                     color=DARK, font="Consolas")

    add_text(s, Inches(0.5), Inches(7.0), Inches(12), Inches(0.3),
             "Relations : User 1-1 Profile / 1-1 Cart / 1-N Order / Cart 1-N CartItem / Order 1-N OrderItem / Product 1-1 Inventory / Category 1-N Product",
             size=9, italic=True, color=GRAY, align=PP_ALIGN.CENTER)

    add_footer(s, prs, page, TOTAL_SLIDES)


def slide_12_customer(prs, page):
    s = blank_slide(prs)
    slide_title(s, "Domaine customer : User & Auth")

    add_text(s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.4),
             "Inscription, connexion, gestion du profil. C'est ce domaine qui heberge la securite (JWT, BCrypt).",
             size=12, italic=True, color=GRAY)

    # Bloc code Java
    code_x = Inches(0.5); code_y = Inches(1.9); code_w = Inches(7.5); code_h = Inches(4.8)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, code_x, code_y, code_w, code_h)
    bg.fill.solid(); bg.fill.fore_color.rgb = RGBColor(0x1E, 0x1E, 0x1E)
    bg.line.fill.background()
    add_text(s, code_x + Inches(0.2), code_y + Inches(0.1), code_w - Inches(0.4), Inches(0.3),
             "AuthService.register()", size=11, color=RGBColor(0x80, 0xC0, 0xFF),
             font="Consolas", italic=True)
    code = """@Transactional
public AuthResponseDto register(RegisterDto dto) {
    if (userRepository.existsByEmail(dto.getEmail()))
        throw new BusinessException(
            "Cet email est deja utilise");
    User user = User.builder()
        .firstName(dto.getFirstName())
        .lastName(dto.getLastName())
        .email(dto.getEmail())
        .password(passwordEncoder.encode(
            dto.getPassword()))
        .role(Role.USER).build();
    Profile profile = Profile.builder()
        .user(user).build();
    user.setProfile(profile);
    User saved = userRepository.save(user);
    return AuthResponseDto.builder()
        .token(jwtService.generateToken(saved))
        .user(UserDto.from(saved)).build();
}"""
    add_text(s, code_x + Inches(0.3), code_y + Inches(0.5),
             code_w - Inches(0.6), code_h - Inches(0.6),
             code, size=10, color=RGBColor(0xDC, 0xDC, 0xDC), font="Consolas")

    # Cartes a droite
    cards_data = [
        ("BCrypt", "Hash one-way avec sel aleatoire (cout 10).\nImpossible de retrouver le mot de passe.", MINT),
        ("@JsonIgnore", "Le password n'est JAMAIS serialise dans les reponses JSON.", LAVENDER),
        ("Email unique", "@UniqueConstraint + check existsByEmail.\nRetour HTTP 409 Conflict si deja pris.", PEACH),
        ("Profile cascade", "@OneToOne mappedBy + cascade ALL.\nLe profil est cree avec l'utilisateur.", TEAL),
    ]
    cw = Inches(4.7); ch = Inches(1.1)
    cx = Inches(8.2)
    for i, (t, b, col) in enumerate(cards_data):
        y = Inches(1.9) + i * (ch + Inches(0.1))
        card(s, cx, y, cw, ch, t, b, accent=col)

    add_footer(s, prs, page, TOTAL_SLIDES)


def slide_13_catalog(prs, page):
    s = blank_slide(prs)
    slide_title(s, "Domaine catalog : recherche paginee")

    add_text(s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.4),
             "Categories + produits avec recherche combinee (mot-cle + filtre categorie) et pagination.",
             size=12, italic=True, color=GRAY)

    # Code JPQL
    code_x = Inches(0.5); code_y = Inches(1.9); code_w = Inches(8.5); code_h = Inches(4.8)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, code_x, code_y, code_w, code_h)
    bg.fill.solid(); bg.fill.fore_color.rgb = RGBColor(0x1E, 0x1E, 0x1E)
    bg.line.fill.background()
    add_text(s, code_x + Inches(0.2), code_y + Inches(0.1),
             code_w - Inches(0.4), Inches(0.3),
             "ProductRepository.search() - JPQL", size=11,
             color=RGBColor(0x80, 0xC0, 0xFF), font="Consolas", italic=True)
    code = """public interface ProductRepository
        extends JpaRepository<Product, Long> {

    @Query(\"\"\"
        SELECT p FROM Product p
        WHERE (:categoryId IS NULL
                OR p.category.id = :categoryId)
          AND (:q IS NULL
                OR LOWER(p.name)
                    LIKE LOWER(CONCAT('%', :q, '%'))
                OR LOWER(p.description)
                    LIKE LOWER(CONCAT('%', :q, '%')))
        \"\"\")
    Page<Product> search(
        @Param("categoryId") Long categoryId,
        @Param("q") String q,
        Pageable pageable);
}"""
    add_text(s, code_x + Inches(0.3), code_y + Inches(0.5),
             code_w - Inches(0.6), code_h - Inches(0.6),
             code, size=11, color=RGBColor(0xDC, 0xDC, 0xDC), font="Consolas")

    # Cartes avantages a droite
    items = [
        ("Requete unique", "Filtre + recherche en une seule requete SQL.", MINT),
        ("Parametres optionnels", "IS NULL OR ... permet de tout combiner.", TEAL),
        ("Pagination", "Pageable pris en charge par Spring Data.", LAVENDER),
        ("Insensible casse", "LOWER + LIKE %% : la recherche capte 'XPS' / 'xps'.", PEACH),
    ]
    cw = Inches(4); ch = Inches(1.1)
    cx = Inches(9.2)
    for i, (t, b, col) in enumerate(items):
        y = Inches(1.9) + i * (ch + Inches(0.1))
        card(s, cx, y, cw, ch, t, b, accent=col)

    add_footer(s, prs, page, TOTAL_SLIDES)


def slide_14_inventory(prs, page):
    s = blank_slide(prs)
    slide_title(s, "Domaine inventory : gestion du stock")

    add_text(s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.4),
             "Stock isole dans une entite separee (Inventory) reliee a Product par @OneToOne.",
             size=12, italic=True, color=GRAY)

    # Schema visuel : Product → Inventory
    bx, by = Inches(1), Inches(2.2)
    pcard = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                bx, by, Inches(2.8), Inches(1.6))
    pcard.adjustments[0] = 0.08
    pcard.fill.solid(); pcard.fill.fore_color.rgb = LAVENDER
    pcard.line.fill.background()
    add_text(s, bx, by + Inches(0.4), Inches(2.8), Inches(0.5),
             "Product", size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, bx + Inches(0.2), by + Inches(0.95), Inches(2.4), Inches(0.6),
             "id, name, price\nimageUrl, category", size=10, color=WHITE,
             align=PP_ALIGN.CENTER, font="Consolas")

    # Fleche 1-1
    add_arrow(s, bx + Inches(2.8), by + Inches(0.8),
              bx + Inches(4.4), by + Inches(0.8), color=DARK, width=2)
    add_text(s, bx + Inches(2.9), by + Inches(0.4), Inches(1.5), Inches(0.4),
             "@OneToOne", size=10, italic=True, color=DARK, font="Consolas")
    add_text(s, bx + Inches(2.9), by + Inches(1.0), Inches(1.5), Inches(0.4),
             "1 ↔ 1", size=11, bold=True, color=DARK, align=PP_ALIGN.CENTER)

    bx2 = bx + Inches(4.4)
    icard = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                bx2, by, Inches(2.8), Inches(1.6))
    icard.adjustments[0] = 0.08
    icard.fill.solid(); icard.fill.fore_color.rgb = SKY
    icard.line.fill.background()
    add_text(s, bx2, by + Inches(0.4), Inches(2.8), Inches(0.5),
             "Inventory", size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, bx2 + Inches(0.2), by + Inches(0.95), Inches(2.4), Inches(0.6),
             "id, quantity\nproduct (FK)", size=10, color=WHITE,
             align=PP_ALIGN.CENTER, font="Consolas")

    # 3 cartes operations
    ops = [
        ("checkAvailability", "Verifie sans modifier : leve BusinessException si quantity insuffisante.", MINT),
        ("decrement", "Verifie ET decrement : utilisee pendant le checkout.", PEACH),
        ("update (admin)", "Permet a l'admin de reapprovisionner via PUT /api/inventory/{id}.", TEAL),
    ]
    ow = Inches(4.0); oh = Inches(1.4)
    for i, (t, b, col) in enumerate(ops):
        x = Inches(0.4) + i * (ow + Inches(0.3))
        y = Inches(4.5)
        card(s, x, y, ow, oh, t, b, accent=col, icon="◇")

    add_text(s, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.5),
             "Cas d'usage : appele depuis CartService (avant ajout) et OrderService (avant validation).",
             size=11, italic=True, color=GRAY, align=PP_ALIGN.CENTER)

    add_footer(s, prs, page, TOTAL_SLIDES)


def slide_15_shopping(prs, page):
    s = blank_slide(prs)
    slide_title(s, "Domaine shopping : gestion du panier")

    # Diagramme : User → Cart → CartItem → Product
    elems = [
        ("User", MINT, Inches(0.7)),
        ("Cart", TEAL, Inches(3.7)),
        ("CartItem", LAVENDER, Inches(6.7)),
        ("Product", PEACH, Inches(9.7)),
    ]
    y = Inches(2)
    for name, col, x in elems:
        c = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                x, y, Inches(2.5), Inches(1.4))
        c.adjustments[0] = 0.1
        c.fill.solid(); c.fill.fore_color.rgb = col
        c.line.fill.background()
        add_text(s, x, y + Inches(0.45), Inches(2.5), Inches(0.5),
                 name, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Cardinalites
    rels = [
        ("1 - 1", Inches(3.2)),
        ("1 - N", Inches(6.2)),
        ("N - 1", Inches(9.2)),
    ]
    for txt, x in rels:
        add_text(s, x, y - Inches(0.2), Inches(0.5), Inches(0.3),
                 txt, size=10, bold=True, color=DARK,
                 align=PP_ALIGN.CENTER, font="Consolas")

    # Fleches
    for i in range(3):
        x1 = Inches(0.7) + Inches(2.5) + Inches(3) * i
        x2 = Inches(0.7) + Inches(3) * (i + 1)
        add_arrow(s, x1, y + Inches(0.7), x2, y + Inches(0.7),
                  color=DARK, width=1.5)

    # 3 cartes operations (en bas)
    ops = [
        ("addItem", "Verifie le stock global (qty existante + ajoutee) avant insertion.\nFige le prix unitaire au moment de l'ajout.", MINT),
        ("updateItem", "Met a jour la quantite avec checkAvailability.\nRequiert que l'item appartienne au panier de l'utilisateur courant.", LAVENDER),
        ("clear", "Vide le panier.\nUtilise apres validation de commande.", PEACH),
    ]
    ow = Inches(4.0); oh = Inches(1.6)
    for i, (t, b, col) in enumerate(ops):
        x = Inches(0.4) + i * (ow + Inches(0.3))
        y = Inches(4.3)
        card(s, x, y, ow, oh, t, b, accent=col, icon="●")

    add_footer(s, prs, page, TOTAL_SLIDES)


def slide_16_billing(prs, page):
    s = blank_slide(prs)
    slide_title(s, "Garantie d'Integrite : Le Compteur Atomique", size=30)

    # Bloc code @Transactional a gauche
    code_x = Inches(0.4); code_y = Inches(1.5); code_w = Inches(7); code_h = Inches(5.5)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, code_x, code_y, code_w, code_h)
    bg.fill.solid(); bg.fill.fore_color.rgb = RGBColor(0x1E, 0x1E, 0x1E)
    bg.line.fill.background()
    add_text(s, code_x + Inches(0.2), code_y + Inches(0.1), code_w - Inches(0.4), Inches(0.3),
             "CODE BLOCK : Transactional Checkout", size=11,
             color=RGBColor(0x80, 0xC0, 0xFF), font="Consolas")
    code = """@Transactional
public OrderDto checkout() {
    Cart cart = cartService.getOrCreateCart(user);
    if (cart.getItems().isEmpty())
        throw new BusinessException(
            "Votre panier est vide");
    for (CartItem ci : cart.getItems())
        inventoryService.checkAvailability(...);
    Order order = Order.builder()
        .user(user).status(PENDING).build();
    for (CartItem ci : cart.getItems()) {
        order.getItems().add(
            OrderItem.builder()...build());
        inventoryService.decrement(...);
    }
    order.setStatus(CONFIRMED);
    Order saved = orderRepository.save(order);
    cartService.clearCart(cart);
    return OrderDto.from(saved);
}"""
    add_text(s, code_x + Inches(0.3), code_y + Inches(0.5),
             code_w - Inches(0.6), code_h - Inches(0.6),
             code, size=10.5, color=RGBColor(0xDC, 0xDC, 0xDC), font="Consolas")

    # Boite ACID a droite
    ax = Inches(7.7); ay = Inches(1.5); aw = Inches(5.3); ah = Inches(5.5)
    abg = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, ax, ay, aw, ah)
    abg.adjustments[0] = 0.04
    abg.fill.solid(); abg.fill.fore_color.rgb = WHITE
    abg.line.color.rgb = MINT; abg.line.width = Pt(1.5)

    # 3 engrenages numerotes
    gear_y = ay + Inches(0.5)
    gear_centers = [
        ("Verification simultanee\ndes stocks", MINT),
        ("Creation commande\n+ Decrement stock", TEAL),
        ("Vidage du panier", PEACH),
    ]
    for i, (txt, col) in enumerate(gear_centers):
        gy = gear_y + Inches(1.5) * i
        oval = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                   ax + Inches(0.4), gy, Inches(1.4), Inches(1.4))
        oval.fill.solid(); oval.fill.fore_color.rgb = col
        oval.line.color.rgb = WHITE; oval.line.width = Pt(2)
        # Numero dans cercle
        nc = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                 ax + Inches(0.55), gy + Inches(0.1),
                                 Inches(0.35), Inches(0.35))
        nc.fill.solid(); nc.fill.fore_color.rgb = WHITE
        nc.line.color.rgb = col; nc.line.width = Pt(1.5)
        tf = nc.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; p.text = str(i + 1)
        for r in p.runs:
            r.font.bold = True; r.font.size = Pt(11); r.font.color.rgb = col

        add_text(s, ax + Inches(2.0), gy + Inches(0.3),
                 Inches(3.1), Inches(0.9),
                 txt, size=12, color=DARK, bold=True)

    # Boite verrou ACID en bas
    add_text(s, ax + Inches(0.3), ay + Inches(5), aw - Inches(0.6), Inches(0.4),
             "🔒 Atomicite ACID",
             size=14, bold=True, color=MINT)
    add_text(s, ax + Inches(0.3), ay + Inches(5.4), aw - Inches(0.6), Inches(2),
             "Tout ou rien. Si une etape echoue, toutes les operations precedentes sont annulees.",
             size=10, color=DARK, italic=True)

    add_footer(s, prs, page, TOTAL_SLIDES)


def slide_17_hybrid(prs, page):
    s = blank_slide(prs)
    slide_title(s, "Justification : La Persistance Hybride", size=32)

    # Trait vertical au milieu
    line = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                   prs.slide_width // 2, Inches(1.6),
                                   prs.slide_width // 2, Inches(5.8))
    line.line.color.rgb = GRAY; line.line.width = Pt(0.75)

    # Colonne gauche : MySQL
    add_text(s, Inches(0.7), Inches(1.6), Inches(0.5), Inches(0.5),
             "▦", size=36, color=MINT)
    add_text(s, Inches(1.5), Inches(1.7), Inches(4.5), Inches(0.5),
             "MySQL 8 - Le Roc", size=20, bold=True, color=DARK)
    add_text(s, Inches(1.5), Inches(2.1), Inches(4.5), Inches(0.4),
             "Relationnel", size=20, bold=True, color=DARK)

    items = [
        ("Paradigme", "SQL Relationnel\nStructure stricte", MINT),
        ("Propriete cle", "Transactions atomiques (ACID)", MINT),
        ("Cas d'usage E-STORE", "Stock, Panier, Commandes,\nFacturation, Utilisateurs", MINT),
    ]
    for i, (t, b, col) in enumerate(items):
        y = Inches(2.8) + i * Inches(0.95)
        card(s, Inches(0.7), y, Inches(5.5), Inches(0.85), t, b, accent=col)

    # Colonne droite : MongoDB
    add_text(s, Inches(7.0), Inches(1.6), Inches(0.5), Inches(0.5),
             "{ }", size=28, bold=True, color=LAVENDER, font="Consolas")
    add_text(s, Inches(7.7), Inches(1.7), Inches(5.5), Inches(0.5),
             "MongoDB 7 - La Flexibilite", size=20, bold=True, color=DARK)
    add_text(s, Inches(7.7), Inches(2.1), Inches(5.5), Inches(0.5),
             "Documentaire", size=20, bold=True, color=DARK)

    items_r = [
        ("Paradigme", "NoSQL (JSON Document)", LAVENDER),
        ("Propriete cle", "Schema flexible, optimise lectures\nfrequentes et donnees semi-structurees", LAVENDER),
        ("Cas d'usage E-STORE", "Les Avis (Reviews)\nfindByProductIdOrderByCreatedAtDesc()", LAVENDER),
    ]
    for i, (t, b, col) in enumerate(items_r):
        y = Inches(2.8) + i * Inches(0.95)
        card(s, Inches(7.0), y, Inches(5.7), Inches(0.85), t, b, accent=col)

    # Bandeau resilience en bas
    by = Inches(6.1); bx = Inches(0.5); bw = prs.slide_width - Inches(1); bh = Inches(0.8)
    bg = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx, by, bw, bh)
    bg.adjustments[0] = 0.2
    bg.fill.solid(); bg.fill.fore_color.rgb = RGBColor(0xFF, 0xF5, 0xDB)
    bg.line.color.rgb = GOLD; bg.line.width = Pt(1)
    add_text(s, bx + Inches(0.3), by + Inches(0.05), Inches(0.5), Inches(0.7),
             "🛡", size=24, color=GOLD)
    add_text(s, bx + Inches(0.9), by + Inches(0.1), bw - Inches(1.2), Inches(0.6),
             "Ingenierie de Resilience : si MongoDB est indisponible, l'application principale "
             "continue (avec un simple avertissement). Le seeder est defensif.",
             size=11, color=DARK)

    add_footer(s, prs, page, TOTAL_SLIDES)


def slide_18_review(prs, page):
    s = blank_slide(prs)
    slide_title(s, "Domaine review : MongoDB en pratique")

    # Bloc Document JSON a gauche
    code_x = Inches(0.5); code_y = Inches(1.6); code_w = Inches(6); code_h = Inches(3.5)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, code_x, code_y, code_w, code_h)
    bg.fill.solid(); bg.fill.fore_color.rgb = RGBColor(0x1E, 0x1E, 0x1E)
    bg.line.fill.background()
    add_text(s, code_x + Inches(0.2), code_y + Inches(0.1),
             code_w - Inches(0.4), Inches(0.3),
             "Document JSON dans la collection 'reviews'", size=11,
             color=RGBColor(0x80, 0xC0, 0xFF), font="Consolas", italic=True)
    code = """{
  "_id": ObjectId("671f3a2..."),
  "productId": 1,
  "userId": 2,
  "authorName": "Nouhaila Test",
  "rating": 5,
  "comment": "Excellent ordinateur,
              performances au top !",
  "createdAt": ISODate("2026-04-28T...")
}"""
    add_text(s, code_x + Inches(0.3), code_y + Inches(0.5),
             code_w - Inches(0.6), code_h - Inches(0.6),
             code, size=11, color=RGBColor(0xDC, 0xDC, 0xDC), font="Consolas")

    # Bloc Repository a droite
    code_x2 = Inches(6.8); code_w2 = Inches(6.2)
    bg2 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, code_x2, code_y, code_w2, code_h)
    bg2.fill.solid(); bg2.fill.fore_color.rgb = RGBColor(0x1E, 0x1E, 0x1E)
    bg2.line.fill.background()
    add_text(s, code_x2 + Inches(0.2), code_y + Inches(0.1),
             code_w2 - Inches(0.4), Inches(0.3),
             "ReviewRepository - methodes derivees", size=11,
             color=RGBColor(0x80, 0xC0, 0xFF), font="Consolas", italic=True)
    code2 = """public interface ReviewRepository
    extends MongoRepository<Review, String> {

    List<Review>
      findByProductIdOrderByCreatedAtDesc(
        Long productId);

    List<Review>
      findByUserIdOrderByCreatedAtDesc(
        Long userId);
}"""
    add_text(s, code_x2 + Inches(0.3), code_y + Inches(0.5),
             code_w2 - Inches(0.6), code_h - Inches(0.6),
             code2, size=11, color=RGBColor(0xDC, 0xDC, 0xDC), font="Consolas")

    # 3 cartes points cles
    items = [
        ("Schema flexible",
         "Les avis n'ont pas de relation stricte. Ajouter un champ ne casse pas l'existant.", LAVENDER),
        ("Methodes derivees",
         "Spring Data Mongo genere la requete a partir du nom de la methode.", MINT),
        ("Index sur productId",
         "@Indexed accelere les recherches par produit (le cas le plus frequent).", PEACH),
    ]
    cw = Inches(4.0); ch = Inches(1.4)
    for i, (t, b, col) in enumerate(items):
        x = Inches(0.5) + i * (cw + Inches(0.25))
        y = Inches(5.4)
        card(s, x, y, cw, ch, t, b, accent=col, icon="◈")

    add_footer(s, prs, page, TOTAL_SLIDES)


def slide_19_jwt_lifecycle(prs, page):
    s = blank_slide(prs)
    slide_title(s, "Mecanisme de Securite : Cycle de Vie JWT", size=30)

    # Deux colonnes
    add_text(s, Inches(1.5), Inches(1.4), Inches(4), Inches(0.5),
             "Navigateur Client", size=16, bold=True, color=DARK,
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(8), Inches(1.4), Inches(4), Inches(0.5),
             "Serveur Spring Boot 6", size=16, bold=True, color=DARK,
             align=PP_ALIGN.CENTER)

    # Boite client
    cb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             Inches(0.5), Inches(2.0), Inches(5.5), Inches(5))
    cb.adjustments[0] = 0.04
    cb.fill.solid(); cb.fill.fore_color.rgb = RGBColor(0xD8, 0xE6, 0xF0)
    cb.line.fill.background()

    # Boite serveur
    sb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             Inches(7.3), Inches(2.0), Inches(5.5), Inches(5))
    sb.adjustments[0] = 0.04
    sb.fill.solid(); sb.fill.fore_color.rgb = RGBColor(0xDA, 0xEB, 0xDD)
    sb.line.fill.background()

    # Etape 1 : POST login
    numbered_circle(s, Inches(0.95), Inches(2.4), 1, color=GRAY)
    add_arrow(s, Inches(1.3), Inches(2.55), Inches(7.7), Inches(2.55), color=GRAY, width=2)
    add_text(s, Inches(2.2), Inches(2.15), Inches(5), Inches(0.5),
             "1. POST /api/auth/login\n(email + password)",
             size=11, color=DARK, align=PP_ALIGN.CENTER, font="Consolas")

    # Etape 2 : validation BCrypt (cote serveur)
    numbered_circle(s, Inches(7.7), Inches(3.0), 2, color=GRAY)
    add_text(s, Inches(8.3), Inches(3.0), Inches(4.4), Inches(0.7),
             "2. Validation AuthService via\npasswordEncoder.matches()\n(BCrypt, cout 10)",
             size=10, color=DARK, font="Consolas")

    # Etape 3 : generation
    numbered_circle(s, Inches(7.7), Inches(4.0), 3, color=GRAY)
    add_text(s, Inches(8.3), Inches(4.0), Inches(4.4), Inches(0.9),
             "3. Generation JwtService :\nJwts.builder().signWith(key)\n(HS256, expir. 24h, claims:\nuid, role, name)",
             size=10, color=DARK, font="Consolas")

    # Token retour
    add_arrow(s, Inches(7.7), Inches(4.7), Inches(1.3), Inches(4.7),
              color=GRAY, width=2)
    add_text(s, Inches(2.0), Inches(4.4), Inches(5), Inches(0.4),
             "Token genere renvoye au client", size=10, italic=True,
             color=DARK, align=PP_ALIGN.CENTER)

    # Etape 4 : stockage
    numbered_circle(s, Inches(0.95), Inches(5.2), 4, color=GRAY)
    add_text(s, Inches(1.55), Inches(5.2), Inches(4), Inches(0.7),
             "4. Stockage du token dans\nlocalStorage\n('estore.token')",
             size=10, color=DARK, font="Consolas")

    # Etape 5 : nouvelle requete
    numbered_circle(s, Inches(0.95), Inches(6.2), 5, color=GRAY)
    add_arrow(s, Inches(1.3), Inches(6.4), Inches(7.7), Inches(6.4),
              color=GRAY, width=2)
    add_text(s, Inches(2.2), Inches(6.05), Inches(5), Inches(0.4),
             "5. Nouvelle requete avec AuthInterceptor :\nBearer <token>",
             size=10, color=DARK, align=PP_ALIGN.CENTER, font="Consolas")

    # Etape 6 : validation
    numbered_circle(s, Inches(7.7), Inches(6.5), 6, color=GRAY)
    add_text(s, Inches(8.3), Inches(6.4), Inches(4.4), Inches(0.5),
             "6. JwtAuthenticationFilter valide.\nAcces aux endpoints (ex:\n@PreAuthorize hasRole('ADMIN'))",
             size=10, color=DARK, font="Consolas")

    add_footer(s, prs, page, TOTAL_SLIDES)


def slide_20_endpoints(prs, page):
    s = blank_slide(prs)
    slide_title(s, "API REST : matrice d'endpoints")

    headers = ["Methode", "URL", "Acces", "Description"]
    rows = [
        ("POST",   "/api/auth/register",         "Public",     "Inscription"),
        ("POST",   "/api/auth/login",            "Public",     "Connexion (renvoie JWT)"),
        ("GET",    "/api/products",              "Public",     "Catalogue (q, categoryId, page, size)"),
        ("GET",    "/api/products/{id}",         "Public",     "Detail produit"),
        ("POST",   "/api/products",              "ADMIN",      "Creer un produit"),
        ("PUT",    "/api/products/{id}",         "ADMIN",      "Modifier un produit"),
        ("DELETE", "/api/products/{id}",         "ADMIN",      "Supprimer un produit"),
        ("GET",    "/api/categories",            "Public",     "Liste des categories"),
        ("POST",   "/api/categories",            "ADMIN",      "Creer une categorie"),
        ("GET",    "/api/cart",                  "Authentifie", "Panier de l'utilisateur courant"),
        ("POST",   "/api/cart/add",              "Authentifie", "Ajouter un produit au panier"),
        ("PUT",    "/api/cart/update",           "Authentifie", "Modifier la quantite"),
        ("DELETE", "/api/cart/remove/{itemId}",  "Authentifie", "Retirer un article"),
        ("DELETE", "/api/cart/clear",            "Authentifie", "Vider le panier"),
        ("POST",   "/api/orders",                "Authentifie", "Valider la commande (transactionnel)"),
        ("GET",    "/api/orders",                "Authentifie", "Historique des commandes"),
        ("POST",   "/api/reviews",               "Authentifie", "Deposer un avis (Mongo)"),
        ("GET",    "/api/reviews/product/{id}",  "Public",     "Avis d'un produit"),
    ]
    # Header
    hx = Inches(0.4); hy = Inches(1.3); hh = Inches(0.4)
    cols_w = [Inches(1.4), Inches(5.3), Inches(1.8), Inches(4.6)]
    cx = hx
    for i, (h, w) in enumerate(zip(headers, cols_w)):
        bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, cx, hy, w, hh)
        bg.fill.solid(); bg.fill.fore_color.rgb = DARK
        bg.line.fill.background()
        add_text(s, cx, hy, w, hh, h, size=11, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        cx += w

    # Lignes
    method_color = {
        "GET": MINT, "POST": TEAL, "PUT": LAVENDER, "DELETE": CORAL
    }
    access_color = {
        "Public": LIGHT_GRAY, "ADMIN": RGBColor(0xFF, 0xE5, 0xCC),
        "Authentifie": RGBColor(0xE5, 0xF1, 0xFB)
    }
    for i, row in enumerate(rows):
        y = hy + hh + Inches(0.32) * i
        cx = hx
        for j, (val, w) in enumerate(zip(row, cols_w)):
            bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, cx, y, w, Inches(0.32))
            if j == 0:
                bg.fill.solid(); bg.fill.fore_color.rgb = method_color.get(val, GRAY)
                bg.line.color.rgb = WHITE; bg.line.width = Pt(0.5)
                add_text(s, cx, y, w, Inches(0.32), val, size=9, bold=True,
                         color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                         font="Consolas")
            elif j == 2:
                bg.fill.solid(); bg.fill.fore_color.rgb = access_color.get(val, WHITE)
                bg.line.color.rgb = RGBColor(0xE0, 0xE5, 0xEA); bg.line.width = Pt(0.3)
                add_text(s, cx, y, w, Inches(0.32), val, size=9, color=DARK,
                         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            else:
                bg.fill.solid(); bg.fill.fore_color.rgb = WHITE if i % 2 == 0 else LIGHT_GRAY
                bg.line.color.rgb = RGBColor(0xE0, 0xE5, 0xEA); bg.line.width = Pt(0.3)
                font = "Consolas" if j == 1 else "Calibri"
                add_text(s, cx + Inches(0.1), y, w, Inches(0.32),
                         val, size=9, color=DARK, anchor=MSO_ANCHOR.MIDDLE, font=font)
            cx += w

    add_footer(s, prs, page, TOTAL_SLIDES)


# ── PART 3 — Frontend, Demo, Conclusion (slides 21-30) ─────────────
def slide_21_facade_angular(prs, page):
    s = blank_slide(prs)
    slide_title(s, "Facade Moderne : Ingenierie Angular 17", size=30)

    items = [
        ("⊞", "Architecture\nStandalone",
         "Suppression totale des NgModule.\nCode plus concis et\ncomposants isoles.\n\n"
         "Lazy loading active par route\ngenerant un bundle initial\nultra-leger de 13 ko, chargeant\nles features a la demande.", TEAL),
        ("∿", "Reactivite via\nSignals",
         "Gestion d'etat moderne\nremplacant les souscriptions\nRxJS complexes.\n\n"
         "Exemple :\ncart.itemCount() declenche la\nmise a jour instantanee et sans\nfriction du badge du panier dans\nle header global.", SKY),
        ("⇆", "Intercepteurs &\nCohesion",
         "Gestion centralisee du flux :\n'AuthInterceptor' attache le\nBearer token silencieusement.\n'ErrorInterceptor' capte les\n401/4xx en affichant des toasts.\n'AuthGuard' verrouille les\nroutes privees.\n\nUI cohesive via Bootstrap 5.", LAVENDER),
    ]
    cw = Inches(4); ch = Inches(5)
    spacing = Inches(0.3)
    total_w = cw * 3 + spacing * 2
    start_x = (prs.slide_width - total_w) // 2
    y = Inches(1.5)
    for i, (icon, title, body, col) in enumerate(items):
        x = start_x + (cw + spacing) * i
        bg = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, cw, ch)
        bg.adjustments[0] = 0.06
        bg.fill.solid(); bg.fill.fore_color.rgb = WHITE
        bg.line.color.rgb = col; bg.line.width = Pt(1)

        # Icone dans cercle
        cy_icon = y + Inches(0.3)
        ic = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                 x + (cw - Inches(0.7)) // 2, cy_icon,
                                 Inches(0.7), Inches(0.7))
        ic.fill.solid(); ic.fill.fore_color.rgb = LIGHT_GRAY
        ic.line.color.rgb = col; ic.line.width = Pt(1.5)
        tf = ic.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; p.text = icon
        for r in p.runs:
            r.font.size = Pt(22); r.font.color.rgb = col

        add_text(s, x + Inches(0.2), cy_icon + Inches(0.85),
                 cw - Inches(0.4), Inches(0.8),
                 title, size=15, bold=True, color=DARK, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.2), cy_icon + Inches(1.7),
                 cw - Inches(0.4), Inches(2.5),
                 body, size=10.5, color=DARK, align=PP_ALIGN.CENTER)

    # Bandeau experience utilisateur en bas
    by = Inches(6.6); bw = Inches(6); bh = Inches(0.6)
    bx = (prs.slide_width - bw) // 2
    bg = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx, by, bw, bh)
    bg.adjustments[0] = 0.5
    bg.fill.solid(); bg.fill.fore_color.rgb = WHITE
    bg.line.color.rgb = SKY; bg.line.width = Pt(1.5)
    add_text(s, bx, by, bw, bh, "Experience Utilisateur Fluide",
             size=14, bold=True, color=DARK, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)

    add_footer(s, prs, page, TOTAL_SLIDES)


def slide_22_frontend_arch(prs, page):
    s = blank_slide(prs)
    slide_title(s, "Architecture Frontend : core / shared / features")

    # 3 colonnes
    cols = [
        ("core/", TEAL,
         ["models/", "  user.model.ts", "  product.model.ts", "  cart.model.ts",
          "  order.model.ts", "  review.model.ts",
          "services/", "  auth.service.ts", "  cart.service.ts",
          "  product.service.ts", "  order.service.ts", "  review.service.ts",
          "  user.service.ts", "  toast.service.ts",
          "guards/", "  auth.guard.ts",
          "interceptors/", "  auth.interceptor.ts", "  error.interceptor.ts"]),
        ("shared/", LAVENDER,
         ["components/", "  header.component.ts", "  footer.component.ts",
          "  loader.component.ts", "  toast.component.ts"]),
        ("features/", PEACH,
         ["auth/", "  login.component.ts", "  register.component.ts",
          "catalog/", "  product-list.component.ts", "  product-detail.component.ts",
          "cart/", "  cart.component.ts",
          "orders/", "  orders.component.ts",
          "profile/", "  profile.component.ts",
          "reviews/", "  review-form.component.ts"]),
    ]
    cw = Inches(4); ch = Inches(5.5)
    spacing = Inches(0.2)
    start_x = Inches(0.4)
    y = Inches(1.5)
    for i, (title, color, files) in enumerate(cols):
        x = start_x + (cw + spacing) * i
        bg = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, cw, ch)
        bg.adjustments[0] = 0.04
        bg.fill.solid(); bg.fill.fore_color.rgb = WHITE
        bg.line.color.rgb = color; bg.line.width = Pt(1.5)
        # Header
        hd = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Inches(0.15), y, cw - Inches(0.3), Inches(0.05))
        hd.fill.solid(); hd.fill.fore_color.rgb = color
        hd.line.fill.background()
        # Titre
        add_text(s, x, y + Inches(0.2), cw, Inches(0.4),
                 title, size=18, bold=True, color=DARK,
                 align=PP_ALIGN.CENTER, font="Consolas")
        # Liste fichiers
        for j, f in enumerate(files):
            indent = Inches(0.3) if not f.startswith("  ") else Inches(0.6)
            add_text(s, x + indent, y + Inches(0.85) + Inches(0.21) * j,
                     cw - indent - Inches(0.2), Inches(0.21),
                     f.strip() if f.startswith("  ") else f,
                     size=9, color=DARK if f.endswith("/") else GRAY,
                     bold=f.endswith("/"), font="Consolas")

    add_text(s, Inches(0.4), Inches(7.0), Inches(12.5), Inches(0.3),
             "Convention : core (singleton, providedIn: 'root') / shared (composants reutilises) / features (pages metier).",
             size=10, italic=True, color=GRAY, align=PP_ALIGN.CENTER)

    add_footer(s, prs, page, TOTAL_SLIDES)


def slide_23_signals(prs, page):
    s = blank_slide(prs)
    slide_title(s, "Reactivite Signals : exemple CartService")

    # Code TypeScript
    code_x = Inches(0.5); code_y = Inches(1.5); code_w = Inches(7.5); code_h = Inches(5)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, code_x, code_y, code_w, code_h)
    bg.fill.solid(); bg.fill.fore_color.rgb = RGBColor(0x1E, 0x1E, 0x1E)
    bg.line.fill.background()
    add_text(s, code_x + Inches(0.2), code_y + Inches(0.1),
             code_w - Inches(0.4), Inches(0.3),
             "core/services/cart.service.ts", size=11,
             color=RGBColor(0x80, 0xC0, 0xFF), font="Consolas", italic=True)
    code = """@Injectable({ providedIn: 'root' })
export class CartService {
  private http = inject(HttpClient);
  private base = `${env.apiUrl}/cart`;

  readonly cart = signal<Cart | null>(null);
  readonly itemCount = signal<number>(0);

  add(productId: number, quantity: number) {
    return this.http
      .post<ApiResponse<Cart>>(
        `${this.base}/add`,
        { productId, quantity })
      .pipe(tap((r) => this.update(r.data)));
  }

  private update(cart: Cart | undefined): void {
    if (!cart) return;
    this.cart.set(cart);
    this.itemCount.set(cart.itemCount);
  }
}"""
    add_text(s, code_x + Inches(0.3), code_y + Inches(0.5),
             code_w - Inches(0.6), code_h - Inches(0.6),
             code, size=10, color=RGBColor(0xDC, 0xDC, 0xDC), font="Consolas")

    # Diagramme reactivite a droite
    add_text(s, Inches(8.3), Inches(1.5), Inches(4.7), Inches(0.5),
             "Mecanisme reactif", size=15, bold=True, color=DARK)

    items = [
        ("1. Composant Header", "lit cartSvc.itemCount()", MINT),
        ("2. Angular detecte", "la dependance au signal", TEAL),
        ("3. Si .set() change", "le DOM est invalide", LAVENDER),
        ("4. Re-render automatique", "du badge du panier", PEACH),
    ]
    for i, (t, b, col) in enumerate(items):
        y = Inches(2.2) + i * Inches(1)
        card(s, Inches(8.3), y, Inches(4.7), Inches(0.85), t, b, accent=col)

    add_footer(s, prs, page, TOTAL_SLIDES)


def slide_24_interceptors(prs, page):
    s = blank_slide(prs)
    slide_title(s, "Interceptors : Auth + Error")

    # AuthInterceptor (gauche)
    code_x = Inches(0.5); code_y = Inches(1.5); code_w = Inches(6.2); code_h = Inches(2.7)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, code_x, code_y, code_w, code_h)
    bg.fill.solid(); bg.fill.fore_color.rgb = RGBColor(0x1E, 0x1E, 0x1E)
    bg.line.fill.background()
    add_text(s, code_x + Inches(0.2), code_y + Inches(0.1),
             code_w - Inches(0.4), Inches(0.3),
             "auth.interceptor.ts", size=11,
             color=RGBColor(0x80, 0xC0, 0xFF), font="Consolas", italic=True)
    code1 = """export const authInterceptor:
    HttpInterceptorFn = (req, next) => {
  const auth = inject(AuthService);
  const token = auth.token;
  if (!token) return next(req);
  return next(req.clone({
    setHeaders: {
      Authorization: `Bearer ${token}`
    }
  }));
};"""
    add_text(s, code_x + Inches(0.3), code_y + Inches(0.5),
             code_w - Inches(0.6), code_h - Inches(0.6),
             code1, size=11, color=RGBColor(0xDC, 0xDC, 0xDC), font="Consolas")

    # ErrorInterceptor (droite)
    code_x2 = Inches(7); code_y = Inches(1.5); code_w2 = Inches(5.9)
    bg2 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, code_x2, code_y, code_w2, code_h)
    bg2.fill.solid(); bg2.fill.fore_color.rgb = RGBColor(0x1E, 0x1E, 0x1E)
    bg2.line.fill.background()
    add_text(s, code_x2 + Inches(0.2), code_y + Inches(0.1),
             code_w2 - Inches(0.4), Inches(0.3),
             "error.interceptor.ts", size=11,
             color=RGBColor(0x80, 0xC0, 0xFF), font="Consolas", italic=True)
    code2 = """return next(req).pipe(
  catchError((err: HttpErrorResponse) => {
    const msg = err.error?.message
              ?? `HTTP ${err.status}`;
    if (err.status === 401) {
      auth.logout();
      toast.error('Session expiree');
    } else if (err.status !== 0) {
      toast.error(msg);
    } else {
      toast.error('Serveur injoignable');
    }
    return throwError(() => err);
  })
);"""
    add_text(s, code_x2 + Inches(0.3), code_y + Inches(0.5),
             code_w2 - Inches(0.6), code_h - Inches(0.6),
             code2, size=11, color=RGBColor(0xDC, 0xDC, 0xDC), font="Consolas")

    # AuthGuard
    items = [
        ("AuthInterceptor", "Attache automatiquement le Bearer token JWT a chaque requete sortante.", TEAL),
        ("ErrorInterceptor", "Capte les erreurs HTTP : 401 -> logout, 4xx -> toast, 0 -> serveur down.", LAVENDER),
        ("AuthGuard", "Bloque l'acces a /cart, /orders, /profile pour les visiteurs non connectes.", MINT),
    ]
    cw = Inches(4.0); ch = Inches(2)
    for i, (t, b, col) in enumerate(items):
        x = Inches(0.4) + i * (cw + Inches(0.3))
        y = Inches(4.6)
        card(s, x, y, cw, ch, t, b, accent=col, icon="◆")

    add_footer(s, prs, page, TOTAL_SLIDES)


def slide_25_user_journey(prs, page):
    s = blank_slide(prs)
    slide_title(s, "Parcours Utilisateur : Tracabilite de la Demonstration",
                size=28)

    steps = [
        ("1.", "Exploration",      "Catalogue : 12 produits,\nrecherche, filtre.",       MINT),
        ("2.", "Inspection",        "Fiche produit : Prix, stock\nen temps reel,\navis Mongo.", TEAL),
        ("3.", "Friction",          "Tentative d'ajout au\npanier -> Redirection\n(Non-authentifie).", PEACH),
        ("4.", "Authentification",  "Login :\nuser@estore.ma",                            SKY),
        ("5.", "Conversion",        "Ajout de 2 produits\n(badge Signals\npasse a 2).",  MINT),
        ("6.", "Ajustement",        "Panier : Modification\ndes quantites, recalcul\nautomatique du total.", LAVENDER),
        ("7.", "Validation",        "Commande confirmee\n-> Routage vers\n/orders.",     MINT),
        ("8.", "Fidelisation",      "Depot d'un avis 5\netoiles persistant sur\nle produit.", GOLD),
    ]
    sw = Inches(1.55); sh = Inches(2.4)
    spacing = Inches(0.05)
    total_w = sw * 8 + spacing * 7
    start_x = (prs.slide_width - total_w) // 2
    y = Inches(3.5)
    for i, (n, label, body, color) in enumerate(steps):
        x = start_x + (sw + spacing) * i

        # Cercle au-dessus
        cy = Inches(2.6); cw_ = Inches(0.5)
        cx = x + (sw - cw_) // 2
        nc = s.shapes.add_shape(MSO_SHAPE.OVAL, cx, cy, cw_, cw_)
        nc.fill.solid(); nc.fill.fore_color.rgb = color
        nc.line.color.rgb = WHITE; nc.line.width = Pt(2)
        tf = nc.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; p.text = n.replace(".", "")
        for r in p.runs:
            r.font.bold = True; r.font.size = Pt(12); r.font.color.rgb = WHITE

        # Label sous le cercle
        add_text(s, x, cy + Inches(0.55), sw, Inches(0.4),
                 label, size=10, bold=True, color=DARK, align=PP_ALIGN.CENTER)

        # Carte
        bg = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, sw, sh)
        bg.adjustments[0] = 0.08
        bg.fill.solid(); bg.fill.fore_color.rgb = WHITE
        bg.line.color.rgb = color; bg.line.width = Pt(1)
        add_text(s, x + Inches(0.1), y + Inches(0.15),
                 sw - Inches(0.2), sh - Inches(0.3),
                 body, size=9, color=DARK, align=PP_ALIGN.CENTER)

    # Ligne reliant les cercles
    line = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                   start_x, Inches(2.85),
                                   start_x + total_w, Inches(2.85))
    line.line.color.rgb = GRAY; line.line.width = Pt(2)

    # Note "Intervention AuthGuard"
    add_text(s, Inches(3.0), Inches(2.0), Inches(2.5), Inches(0.4),
             "Intervention\nAuthGuard", size=9, italic=True, color=GRAY,
             align=PP_ALIGN.CENTER)

    add_footer(s, prs, page, TOTAL_SLIDES)


def slide_26_tests(prs, page):
    s = blank_slide(prs)
    slide_title(s, "Tests unitaires : 9 tests, BUILD SUCCESS")

    # Tableau
    tests = [
        ("ProductServiceTest",  "4", "findById ok / KO\nsearch avec / sans mot-cle", MINT),
        ("CartServiceTest",      "2", "addItem stock OK\naddItem stock insuffisant", TEAL),
        ("OrderServiceTest",     "3", "checkout panier vide\ncheckout valide\nmyOrders trie", LAVENDER),
    ]
    cw = Inches(4.0); ch = Inches(2.4)
    for i, (suite, n, cas, col) in enumerate(tests):
        x = Inches(0.4) + i * (cw + Inches(0.3))
        y = Inches(1.5)
        bg = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, cw, ch)
        bg.adjustments[0] = 0.06
        bg.fill.solid(); bg.fill.fore_color.rgb = WHITE
        bg.line.color.rgb = col; bg.line.width = Pt(1.5)
        # Big number
        add_text(s, x, y + Inches(0.2), cw, Inches(0.8),
                 n, size=48, bold=True, color=col, align=PP_ALIGN.CENTER, font="Calibri")
        add_text(s, x, y + Inches(1.05), cw, Inches(0.4),
                 "tests", size=12, color=GRAY, align=PP_ALIGN.CENTER)
        # Title
        add_text(s, x, y + Inches(1.45), cw, Inches(0.4),
                 suite, size=12, bold=True, color=DARK, align=PP_ALIGN.CENTER, font="Consolas")
        # Cas
        add_text(s, x + Inches(0.2), y + Inches(1.85), cw - Inches(0.4), Inches(0.6),
                 cas, size=9.5, color=DARK, align=PP_ALIGN.CENTER, italic=True)

    # Console output
    code_x = Inches(0.4); code_y = Inches(4.4); code_w = Inches(12.5); code_h = Inches(2.6)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, code_x, code_y, code_w, code_h)
    bg.fill.solid(); bg.fill.fore_color.rgb = RGBColor(0x1E, 0x1E, 0x1E)
    bg.line.fill.background()
    out = """$ mvn test
[INFO] -------------------------------------------------------
[INFO]  T E S T S
[INFO] -------------------------------------------------------
[INFO] Running com.estore.catalog.ProductServiceTest
[INFO] Tests run: 4, Failures: 0, Errors: 0, Skipped: 0
[INFO] Running com.estore.shopping.CartServiceTest
[INFO] Tests run: 2, Failures: 0, Errors: 0, Skipped: 0
[INFO] Running com.estore.billing.OrderServiceTest
[INFO] Tests run: 3, Failures: 0, Errors: 0, Skipped: 0
[INFO] Tests run: 9, Failures: 0, Errors: 0, Skipped: 0
[INFO] BUILD SUCCESS"""
    add_text(s, code_x + Inches(0.3), code_y + Inches(0.2),
             code_w - Inches(0.6), code_h - Inches(0.4),
             out, size=10, color=RGBColor(0xCC, 0xFF, 0xCC), font="Consolas")

    add_footer(s, prs, page, TOTAL_SLIDES)


def slide_27_docker(prs, page):
    s = blank_slide(prs)

    # Console d'erreur en haut
    code_x = Inches(2); code_y = Inches(0.5); code_w = Inches(9.3); code_h = Inches(1.3)
    bg = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, code_x, code_y, code_w, code_h)
    bg.adjustments[0] = 0.08
    bg.fill.solid(); bg.fill.fore_color.rgb = RGBColor(0x2D, 0x2D, 0x2D)
    bg.line.color.rgb = CORAL; bg.line.width = Pt(1.5)
    # Boutons macOS
    for i, color in enumerate([RGBColor(0xFF, 0x60, 0x57), RGBColor(0xFF, 0xBD, 0x2E),
                                RGBColor(0x28, 0xC9, 0x40)]):
        d = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                code_x + Inches(0.2) + Inches(0.25) * i,
                                code_y + Inches(0.15),
                                Inches(0.18), Inches(0.18))
        d.fill.solid(); d.fill.fore_color.rgb = color
        d.line.fill.background()
    add_text(s, code_x + Inches(0.3), code_y + Inches(0.4),
             code_w - Inches(0.6), code_h - Inches(0.5),
             "Obstacle d'Environnement :\nC:\\ProgramData\\DockerDesktop must be owned by an elevated\n"
             "account (Cause: Dossier residuel systeme)",
             size=11, color=RGBColor(0xDC, 0xDC, 0xDC), font="Consolas")

    # Titre central
    add_text(s, Inches(0.5), Inches(2), Inches(12.3), Inches(0.7),
             "Intelligence DevOps : Le Bifurcateur de Resolution",
             size=24, bold=True, color=DARK, align=PP_ALIGN.CENTER)

    # 2 solutions (gauche/droite)
    s1x = Inches(1); s1y = Inches(3.2); sw = Inches(5); sh = Inches(2)
    bg1 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, s1x, s1y, sw, sh)
    bg1.adjustments[0] = 0.06
    bg1.fill.solid(); bg1.fill.fore_color.rgb = WHITE
    bg1.line.color.rgb = GRAY; bg1.line.width = Pt(1.2)
    add_text(s, s1x, s1y + Inches(0.2), sw, Inches(0.5),
             "Solution 1 : Brute", size=14, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    add_text(s, s1x + Inches(0.3), s1y + Inches(0.8), sw - Inches(0.6), Inches(1.1),
             "Nettoyage via ligne de commande Administrateur (rmdir /s /q) "
             "suivi d'une reinstallation propre de l'environnement.",
             size=11, color=DARK, align=PP_ALIGN.CENTER)

    s2x = Inches(7.3); s2y = Inches(3.2)
    bg2 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, s2x, s2y, sw, sh)
    bg2.adjustments[0] = 0.06
    bg2.fill.solid(); bg2.fill.fore_color.rgb = RGBColor(0xE8, 0xF8, 0xEA)
    bg2.line.color.rgb = MINT; bg2.line.width = Pt(2)
    add_text(s, s2x, s2y + Inches(0.2), sw, Inches(0.5),
             "Solution 2 : Elegante", size=14, bold=True, color=MINT, align=PP_ALIGN.CENTER)
    add_text(s, s2x + Inches(0.3), s2y + Inches(0.8), sw - Inches(0.6), Inches(1.1),
             "Activation du profil 'dev' Spring Boot. Bascule automatique\n"
             "sur la base relationnelle 'H2 in-memory'. L'application demarre\n"
             "et tourne parfaitement sans Docker.",
             size=11, color=DARK, align=PP_ALIGN.CENTER)

    # Fleches partant du titre vers les 2 solutions
    add_arrow(s, Inches(4), Inches(2.8), Inches(3.5), Inches(3.2), color=GRAY)
    add_arrow(s, Inches(9.3), Inches(2.8), Inches(9.8), Inches(3.2), color=MINT)

    # Cercle qualite garantie en bas
    cx = (prs.slide_width - Inches(2.8)) // 2; cy = Inches(5.5)
    big = s.shapes.add_shape(MSO_SHAPE.OVAL, cx, cy, Inches(2.8), Inches(1.6))
    big.fill.solid(); big.fill.fore_color.rgb = RGBColor(0xE8, 0xF8, 0xEA)
    big.line.color.rgb = MINT; big.line.width = Pt(2)
    add_text(s, cx, cy + Inches(0.15), Inches(2.8), Inches(0.4),
             "Qualite Garantie :", size=12, bold=True, color=MINT, align=PP_ALIGN.CENTER)
    add_text(s, cx, cy + Inches(0.55), Inches(2.8), Inches(1.0),
             "9 tests unitaires valides\n(JUnit/Mockito en mode\n@MockitoSettings(LENIENT))\n0 echec\n(BUILD SUCCESS)",
             size=9, color=DARK, align=PP_ALIGN.CENTER)
    # Fleche solution 2 -> cercle
    add_arrow(s, s2x + sw // 2, s2y + sh, cx + Inches(1.4), cy, color=MINT, width=2)

    add_footer(s, prs, page, TOTAL_SLIDES)


def slide_28_bilan(prs, page):
    s = blank_slide(prs)
    slide_title(s, "Bilan : objectifs atteints")

    items = [
        ("✓ Architecture", "3 couches × 5 domaines rigoureusement compartimentee.", MINT),
        ("✓ Securite", "JWT + BCrypt + @PreAuthorize - etanche.", PEACH),
        ("✓ Persistance hybride", "SQL (MySQL/H2) + NoSQL (MongoDB) fonctionnelle.", LAVENDER),
        ("✓ Tests", "9 tests unitaires verts, BUILD SUCCESS reproductible.", TEAL),
        ("✓ DevOps", "Docker Compose + profil dev H2 (resilience).", SKY),
        ("✓ Documentation", "README + 7 PDF + Presentation NotebookLM-style.", GOLD),
    ]
    cw = Inches(6); ch = Inches(1.5)
    for i, (t, b, col) in enumerate(items):
        c = i % 2; r = i // 2
        x = Inches(0.4) + c * (cw + Inches(0.4))
        y = Inches(1.5) + r * (ch + Inches(0.2))
        card(s, x, y, cw, ch, t, b, accent=col)

    add_footer(s, prs, page, TOTAL_SLIDES)


def slide_29_roadmap(prs, page):
    s = blank_slide(prs)
    slide_title(s, "Feuille de Route : Projection")

    # Timeline en escalier
    steps = [
        ("Monetisation",   "Integration passerelles\nStripe / CMI", SKY),
        ("Engagement",      "Notifications emails &\nSMS automatisees", MINT),
        ("Intelligence",    "Recommandations ML\nbasees sur l'historique", PEACH),
        ("Scalabilite",     "Deploiement Cloud\nAWS / Azure", LAVENDER),
    ]
    sw = Inches(2.6); sh = Inches(1.5)
    base_y = Inches(5.5)
    base_x = Inches(1.5)
    for i, (title, body, col) in enumerate(steps):
        x = base_x + Inches(2.7) * i
        y = base_y - Inches(0.7) * i

        # Cube isometrique simulé par rectangles colores
        cube = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, sw, sh)
        cube.adjustments[0] = 0.08
        cube.fill.solid(); cube.fill.fore_color.rgb = col
        cube.line.fill.background()
        add_text(s, x, y + Inches(0.1), sw, Inches(0.4),
                 title, size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.2), y + Inches(0.55), sw - Inches(0.4), Inches(0.9),
                 body, size=10, color=WHITE, align=PP_ALIGN.CENTER, italic=True)

        if i < len(steps) - 1:
            x_next = base_x + Inches(2.7) * (i + 1)
            y_next = base_y - Inches(0.7) * (i + 1)
            add_arrow(s, x + sw, y + sh // 2, x_next, y_next + sh // 2,
                      color=col, width=2)

    # Fleche montante decorative
    add_text(s, Inches(0.5), Inches(1.5), Inches(12), Inches(0.5),
             "Direction d'evolution : du business immediat vers la scalabilite cloud",
             size=12, italic=True, color=GRAY, align=PP_ALIGN.CENTER)

    add_footer(s, prs, page, TOTAL_SLIDES)


def slide_30_thanks(prs, page):
    s = blank_slide(prs)
    add_grid_pattern(s, prs)

    # Texte "Merci pour votre attention"
    add_text(s, Inches(0.5), Inches(2.5), Inches(12.3), Inches(1.5),
             "Merci pour votre attention.", size=64, bold=True,
             color=DARK, align=PP_ALIGN.CENTER)

    # Sous-texte
    add_text(s, Inches(0.5), Inches(4.0), Inches(12.3), Inches(0.7),
             "Nous sommes prets a repondre a vos questions.",
             size=22, color=DARK, align=PP_ALIGN.CENTER, italic=True)

    # 3 pilules avec contacts
    pw = Inches(3.8); ph = Inches(1.0); spacing = Inches(0.3)
    total_w = pw * 3 + spacing * 2
    start_x = (prs.slide_width - total_w) // 2
    y = Inches(5.7)

    pill(s, start_x, y, pw, ph,
         "Auteurs :\nAkram Belmoussa\n& Nouhaila Ben Soumane",
         fill=WHITE, accent=TEAL, size=11, font="Consolas",
         align=PP_ALIGN.CENTER)

    pill(s, start_x + pw + spacing, y, pw, ph,
         "Module : Full-Stack\nPr. Omar Zahour",
         fill=WHITE, accent=MINT, size=12, font="Calibri",
         align=PP_ALIGN.CENTER)

    pill(s, start_x + 2 * (pw + spacing), y, pw, ph,
         "FSBM\nUniversite Hassan II\n2025-2026",
         fill=WHITE, accent=LAVENDER, size=11, font="Calibri",
         align=PP_ALIGN.CENTER)

    add_footer(s, prs, page, TOTAL_SLIDES)


# =====================================================================
# BUILDER PPTX
# =====================================================================
def build_pptx():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    builders = [
        slide_01_cover, slide_02_mandat, slide_03_plan, slide_04_specifications,
        slide_05_matrix, slide_06_layers, slide_07_domains, slide_08_ecosystem,
        slide_09_frontend_stack, slide_10_backend_stack,
        slide_11_data_model, slide_12_customer, slide_13_catalog, slide_14_inventory,
        slide_15_shopping, slide_16_billing, slide_17_hybrid, slide_18_review,
        slide_19_jwt_lifecycle, slide_20_endpoints,
        slide_21_facade_angular, slide_22_frontend_arch, slide_23_signals,
        slide_24_interceptors, slide_25_user_journey, slide_26_tests,
        slide_27_docker, slide_28_bilan, slide_29_roadmap, slide_30_thanks,
    ]
    for i, fn in enumerate(builders, 1):
        fn(prs, i)

    out = os.path.join(OUT_DIR, "presentation-avancee.pptx")
    prs.save(out)
    return out


# =====================================================================
# DOSSIER-SCRIPT POUR NOTEBOOKLM (3 PDF, 10 slides chacun)
# =====================================================================
PRIMARY_RL = HexColor("#1F2937")
TEAL_RL    = HexColor(TEAL_HEX)
MINT_RL    = HexColor(MINT_HEX)
LAV_RL     = HexColor(LAVENDER_HEX)
PEACH_RL   = HexColor(PEACH_HEX)
SKY_RL     = HexColor(SKY_HEX)
BG_RL      = HexColor(BG_HEX)


# Definitions detaillees pour NotebookLM (titre + body + visual cues)
NOTEBOOK_SCRIPTS = {
    "Part 1 - Fondations & Architecture": {
        "color": MINT_RL,
        "intro": (
            "Cette premiere partie pose les fondations du projet E-Store : contexte, "
            "objectifs pedagogiques, architecture en 3 couches techniques croisees avec "
            "5+1 domaines fonctionnels (DDD), et la stack technologique complete. "
            "Charger ce dossier dans NotebookLM pour generer les slides 1 a 10."
        ),
        "slides": [
            {
                "num": 1,
                "title": "Couverture - E-STORE",
                "subtitle": "Ingenierie d'une Plateforme E-Commerce Full-Stack : Conception, Persistance Hybride et Securite",
                "content": [
                    "Titre principal en tres grand : E-STORE.",
                    "Sous-titre sur deux lignes (rendu typographique large).",
                    "Trois pilules d'identification en bas :",
                    " - Auteurs : Akram Belmoussa & Nouhaila Ben Soumane",
                    " - Supervision academique : Pr. Omar Zahour",
                    " - Etablissement : FSBM - Universite Hassan II (Annee 2025-2026)",
                    "Style visuel : grille de fond legere (blueprint), typographie sans serif moderne, fond bleu glacier #E8EEF2.",
                ],
            },
            {
                "num": 2,
                "title": "Mandat d'Ingenierie",
                "subtitle": "Contexte & objectifs",
                "content": [
                    "Layout en deux colonnes :",
                    "Colonne gauche - Carte 'Contexte Projet' :",
                    " - Mini-projet pedagogique e-commerce.",
                    " - Perimetre d'application complet : Inscription, catalogue, panier, commandes, gestion des avis.",
                    " - Petit schema isometrique de boites empilees (cubes 3D) avec etiquettes.",
                    "Colonne droite - 5 cartes objectifs, chacune avec une bande couleur a gauche et un cercle-icone :",
                    " - Objectif 1 (vert mint) : Maitriser une architecture full-stack moderne.",
                    " - Objectif 2 (lavande) : Illustrer la puissance d'une persistance hybride SQL + NoSQL.",
                    " - Objectif 3 (gris) : Appliquer les bonnes pratiques (DTOs, @Transactional, tests unitaires).",
                    " - Objectif 4 (peche) : Verrouiller les acces via JWT et Spring Security 6.",
                    " - Objectif 5 (turquoise) : Produire une UI reactive avec Angular 17 en mode standalone.",
                ],
            },
            {
                "num": 3,
                "title": "Plan de la presentation",
                "subtitle": "30 slides en 3 parties",
                "content": [
                    "Trois colonnes egales (3 cartes), chacune surmontee d'un bandeau de couleur :",
                    "Colonne 1 - Fondations (vert mint) : couverture, mandat, plan, cahier des charges, matrice architecturale, couches techniques, domaines DDD, ecosysteme tech, stack frontend, stack backend.",
                    "Colonne 2 - Backend & Securite (lavande) : modele de donnees, customer, catalog, inventory, shopping, billing, persistance hybride, review (Mongo), cycle JWT, endpoints REST.",
                    "Colonne 3 - Frontend, Demo & Conclusion (peche) : facade Angular, architecture frontend, signals, interceptors, parcours utilisateur, tests, Docker, bilan, feuille de route, Q&R.",
                ],
            },
            {
                "num": 4,
                "title": "Cahier des charges fonctionnel",
                "subtitle": "6 fonctionnalites cles",
                "content": [
                    "Six cartes en grille 2x3, chacune avec une bande couleur et un titre + description :",
                    " - Inscription & connexion (vert) : email unique, BCrypt, USER/ADMIN.",
                    " - Catalogue produits (turquoise) : recherche par mot-cle, filtre, pagination, fiche detaillee.",
                    " - Panier (sky blue) : ajout/modif/suppression, prix fige, controle stock.",
                    " - Commandes (lavande) : validation transactionnelle, decrement stock, historique.",
                    " - Avis Reviews (peche) : note 1-5 + commentaire, MongoDB.",
                    " - Administration (corail) : CRUD reserve via @PreAuthorize.",
                ],
            },
            {
                "num": 5,
                "title": "Matrice Architecturale : 3 Couches x 6 Domaines",
                "subtitle": "Vue tabulaire avec connecteurs",
                "content": [
                    "Tableau matriciel : 6 lignes (domaines en italique a gauche : customer, catalog, inventory, shopping, billing, review) x 3 colonnes (PRESENTATION, LOGIQUE METIER, ACCES AUX DONNEES).",
                    "Chaque colonne a un bandeau de couleur en haut (turquoise / vert mint / lavande).",
                    "Cellules en pilules blanches avec texte en monospace (Composants Auth, User Profile JWT, JPA, etc).",
                    "Connecteurs gris (fleches) entre les cellules : Composants Auth -> User Profile JWT -> JPA, etc.",
                    "Derniere ligne (review) : la cellule 'Acces aux donnees' est coloree en lavande avec texte 'MongoDB' en blanc - mise en valeur de la persistance documentaire.",
                ],
            },
            {
                "num": 6,
                "title": "Justification : Les 3 Couches Techniques",
                "subtitle": "Vue verticale empilee",
                "content": [
                    "Trois grandes pilules verticales cote a cote :",
                    "Couche 1 (turquoise) PRESENTATION : Angular 17 standalone, Bootstrap 5 + signals, responsabilites = affichage UI + validation client.",
                    "Couche 2 (vert mint) LOGIQUE METIER : Spring Boot 3.3, Controllers + Services + DTOs, regles metier + orchestration + @Transactional + securite.",
                    "Couche 3 (lavande) ACCES AUX DONNEES : Spring Data JPA + Spring Data MongoDB, persistance + requetes + abstraction du SGBD.",
                    "Sous-titre en bas : 'Communication HTTP/JSON -> appels methodes -> requetes SQL/Mongo'.",
                ],
            },
            {
                "num": 7,
                "title": "Justification : Les Domaines (DDD)",
                "subtitle": "6 sous-packages autonomes",
                "content": [
                    "Sous-titre : 'Chaque domaine est un sous-package autonome (entity, dto, repository, service, controller).'",
                    "6 cartes en grille 2x3 (chaque carte a un cercle-icone + titre + corps) :",
                    " - customer (vert mint) : User, Profile, JWT, Auth.",
                    " - catalog (turquoise) : Categorie, Produit, recherche.",
                    " - inventory (sky blue) : Stock par produit.",
                    " - shopping (lavande) : Cart, CartItem, gestion panier.",
                    " - billing (peche) : Order, OrderItem, checkout transactionnel.",
                    " - review (corail) : Avis (NoSQL MongoDB).",
                ],
            },
            {
                "num": 8,
                "title": "Cartographie de l'Ecosysteme Technologique",
                "subtitle": "Mind map avec noyau central",
                "content": [
                    "Schema mind-map circulaire :",
                    "Au centre : un grand carre arrondi noir avec texte blanc 'Noyau Application E-STORE'.",
                    "Autour, 6 satellites (cartes blanches bordees de couleur) relies au noyau par des connecteurs colores :",
                    " - Frontend (turquoise) en haut : Angular 17 standalone, Bootstrap 5, RxJS, TypeScript 5.",
                    " - Backend (vert mint) a gauche : Spring Boot 3.3, Spring Security 6, JWT (jjwt 0.12).",
                    " - Persistance (lavande) en bas : Spring Data JPA + Hibernate, Spring Data MongoDB, MySQL 8 / H2 / MongoDB 7.",
                    " - Build (sky blue) en haut a droite : Maven 3.9, npm, Angular CLI.",
                    " - Tests (vert mint) a droite : JUnit 5 + Mockito (9 tests, BUILD SUCCESS).",
                    " - DevOps (lavande) a gauche : Docker Compose, MySQL + Mongo + phpMyAdmin + mongo-express.",
                ],
            },
            {
                "num": 9,
                "title": "Stack Frontend en detail",
                "subtitle": "4 piliers Angular",
                "content": [
                    "4 cartes en grille 2x2 (chaque carte avec un cercle-icone) :",
                    " - Angular 17 (turquoise) : Framework SPA, version standalone (pas de NgModule). TypeScript 5, RxJS 7, signals.",
                    " - Bootstrap 5 (lavande) : Framework CSS responsive. Grille 12 col, composants. Pas de JS Bootstrap.",
                    " - RxJS + Signals (vert mint) : RxJS pour les flux HTTP. Signals pour l'etat reactif local.",
                    " - Tooling (sky blue) : Angular CLI + npm. Lazy loading -> bundle initial 13 ko. AOT + tree-shaking.",
                ],
            },
            {
                "num": 10,
                "title": "Stack Backend & DevOps en detail",
                "subtitle": "4 piliers Spring",
                "content": [
                    "4 cartes en grille 2x2 :",
                    " - Spring Boot 3.3 (vert mint) : Web (REST), Validation, Data JPA + Data MongoDB, Security 6, Lombok.",
                    " - Persistance (lavande) : Hibernate, MySQL 8 prod / H2 in-memory dev, MongoDB 7 (avis).",
                    " - Securite (peche) : Spring Security 6 stateless, JWT HS256 (24h), BCrypt sur mots de passe.",
                    " - Build & Test (sky blue) : Maven 3.9, JUnit 5 + Mockito (9 tests), Docker Compose (4 services).",
                ],
            },
        ],
    },
    "Part 2 - Backend & Securite": {
        "color": LAV_RL,
        "intro": (
            "Cette deuxieme partie plonge dans le code backend, domaine par domaine, "
            "puis explore le mecanisme de securite JWT et la matrice complete des "
            "endpoints REST. Charger ce dossier dans NotebookLM pour generer les slides 11 a 20."
        ),
        "slides": [
            {
                "num": 11,
                "title": "Modele de donnees relationnel",
                "subtitle": "Diagramme ERD",
                "content": [
                    "Diagramme entite-relation avec 9 boites entites bordees de couleur :",
                    " - User (vert) : id, email UK, password, role.",
                    " - Profile (vert) : id, phone, address, city, country.",
                    " - Cart (turquoise) : id, createdAt, updatedAt.",
                    " - CartItem (turquoise) : id, quantity, unitPrice.",
                    " - Category (lavande) : id, name, description.",
                    " - Product (lavande) : id, name, description, price BigDecimal, imageUrl.",
                    " - Inventory (lavande) : id, quantity int.",
                    " - Order (peche) : id, orderDate, totalAmount, status OrderStatus.",
                    " - OrderItem (peche) : id, quantity, unitPrice.",
                    "Sous-titre : Relations 1-1 / 1-N entre toutes ces entites.",
                ],
            },
            {
                "num": 12,
                "title": "Domaine customer : User & Auth",
                "subtitle": "Code AuthService.register()",
                "content": [
                    "Layout 2 colonnes :",
                    "Colonne gauche - Bloc de code Java sur fond noir (~16 lignes) :",
                    " - Annotation @Transactional",
                    " - Verification existsByEmail -> BusinessException",
                    " - Construction User avec passwordEncoder.encode (BCrypt)",
                    " - Creation Profile en cascade",
                    " - Sauvegarde + generation token JWT",
                    "Colonne droite - 4 cartes points cles :",
                    " - BCrypt (vert) : hash one-way + sel aleatoire.",
                    " - @JsonIgnore (lavande) : password jamais serialise.",
                    " - Email unique (peche) : 409 Conflict si deja pris.",
                    " - Profile cascade (turquoise) : @OneToOne mappedBy.",
                ],
            },
            {
                "num": 13,
                "title": "Domaine catalog : recherche paginee",
                "subtitle": "JPQL avec parametres optionnels",
                "content": [
                    "Layout 2 colonnes :",
                    "Colonne gauche - Bloc de code JPQL (sur fond noir) avec ProductRepository.search() :",
                    " - @Query avec text block multi-lignes",
                    " - SELECT p FROM Product p WHERE (:categoryId IS NULL OR p.category.id = :categoryId)",
                    " - AND (:q IS NULL OR LOWER(p.name) LIKE LOWER(CONCAT('%', :q, '%')) OR LOWER(p.description) LIKE ...)",
                    " - Page<Product> search avec Pageable",
                    "Colonne droite - 4 cartes avantages :",
                    " - Requete unique (vert)",
                    " - Parametres optionnels (turquoise)",
                    " - Pagination (lavande)",
                    " - Insensible a la casse (peche)",
                ],
            },
            {
                "num": 14,
                "title": "Domaine inventory : gestion du stock",
                "subtitle": "Schema relation + 3 operations",
                "content": [
                    "Schema en haut : 2 boites Product (lavande) -- @OneToOne 1-1 --> Inventory (sky blue).",
                    "3 cartes operations en bas :",
                    " - checkAvailability (vert) : verifie sans modifier, leve BusinessException si insuffisant.",
                    " - decrement (peche) : verifie ET decrement, utilisee pendant le checkout.",
                    " - update (admin) (turquoise) : reapprovisionnement via PUT /api/inventory/{id}.",
                    "Note en bas : 'Cas d'usage : appele depuis CartService et OrderService.'",
                ],
            },
            {
                "num": 15,
                "title": "Domaine shopping : gestion du panier",
                "subtitle": "Diagramme relations + 3 operations",
                "content": [
                    "Diagramme horizontal : User (vert) -- 1-1 --> Cart (turquoise) -- 1-N --> CartItem (lavande) -- N-1 --> Product (peche).",
                    "Toutes les boites en pilules colorees + cardinalites au-dessus des fleches.",
                    "3 cartes operations en bas :",
                    " - addItem (vert) : verifie le stock global avant insertion. Fige le prix.",
                    " - updateItem (lavande) : met a jour la quantite avec checkAvailability.",
                    " - clear (peche) : vide le panier apres validation de commande.",
                ],
            },
            {
                "num": 16,
                "title": "Garantie d'Integrite : Le Compteur Atomique",
                "subtitle": "checkout() en transactionnel",
                "content": [
                    "Layout 2 colonnes :",
                    "Colonne gauche - Bloc de code Java sur fond noir (~18 lignes) :",
                    " - @Transactional public OrderDto checkout()",
                    " - Verification panier vide -> BusinessException",
                    " - Boucle 1 : checkAvailability pour TOUS les items",
                    " - Boucle 2 : creation OrderItem + inventoryService.decrement",
                    " - Sauvegarde + clearCart + return",
                    "Colonne droite - 3 cercles verticaux (engrenages) :",
                    " - 1. Verification simultanee des stocks (vert)",
                    " - 2. Creation commande + Decrementation (turquoise)",
                    " - 3. Vidage du panier (peche)",
                    "En bas : icone cadenas + texte 'Atomicite ACID : Tout ou rien. En cas d'erreur, toutes les operations precedentes sont annulees.'",
                ],
            },
            {
                "num": 17,
                "title": "Justification : La Persistance Hybride",
                "subtitle": "MySQL vs MongoDB",
                "content": [
                    "Layout 2 colonnes separees par un trait vertical :",
                    "Colonne gauche - MySQL 8 - Le Roc Relationnel :",
                    " - Icone tableau",
                    " - 3 cartes : Paradigme (SQL Relationnel, structure stricte) / Propriete cle (Transactions atomiques ACID) / Cas d'usage (Stock, Panier, Commandes, Facturation, Utilisateurs).",
                    "Colonne droite - MongoDB 7 - La Flexibilite Documentaire :",
                    " - Icone JSON",
                    " - 3 cartes : Paradigme (NoSQL JSON Document) / Propriete cle (Schema flexible, optimise lectures frequentes) / Cas d'usage (Avis Reviews, findByProductIdOrderByCreatedAtDesc).",
                    "Bandeau jaune en bas avec icone bouclier : 'Ingenierie de Resilience : Comportement defensif - l'application principale continue de fonctionner meme si MongoDB tombe.'",
                ],
            },
            {
                "num": 18,
                "title": "Domaine review : MongoDB en pratique",
                "subtitle": "Document JSON + Repository",
                "content": [
                    "Layout 2 colonnes - 2 blocs de code sur fond noir :",
                    "Colonne gauche : Document JSON exemple :",
                    " - _id ObjectId, productId, userId, authorName, rating, comment, createdAt ISODate",
                    "Colonne droite : ReviewRepository :",
                    " - extends MongoRepository<Review, String>",
                    " - findByProductIdOrderByCreatedAtDesc",
                    " - findByUserIdOrderByCreatedAtDesc",
                    "3 cartes points cles en bas :",
                    " - Schema flexible (lavande)",
                    " - Methodes derivees (vert) : Spring Data Mongo genere la requete a partir du nom.",
                    " - Index sur productId (peche) : @Indexed accelere les recherches.",
                ],
            },
            {
                "num": 19,
                "title": "Mecanisme de Securite : Cycle de Vie JWT",
                "subtitle": "Sequence diagram en 6 etapes",
                "content": [
                    "Schema sequence diagram en 2 colonnes :",
                    "Colonne gauche (bleue) - 'Navigateur Client'.",
                    "Colonne droite (verte) - 'Serveur Spring Boot 6'.",
                    "Etapes (cercles numerotes + fleches reliant les colonnes) :",
                    " 1. Client -> Serveur : POST /api/auth/login (email + password)",
                    " 2. Serveur : Validation AuthService via passwordEncoder.matches() (BCrypt, cout 10)",
                    " 3. Serveur : Generation JwtService - Jwts.builder().signWith(key) (HS256, expiration 24h, claims uid/role/name)",
                    " 4. Serveur -> Client : Token genere renvoye au client",
                    " 4. Client : Stockage du token dans localStorage ('estore.token')",
                    " 5. Client -> Serveur : Nouvelle requete avec AuthInterceptor : Bearer <token>",
                    " 6. Serveur : JwtAuthenticationFilter valide. Acces aux endpoints (ex: @PreAuthorize hasRole('ADMIN'))",
                ],
            },
            {
                "num": 20,
                "title": "API REST : matrice d'endpoints",
                "subtitle": "18 endpoints en tableau",
                "content": [
                    "Tableau de 18 lignes avec 4 colonnes : Methode, URL, Acces, Description.",
                    "Methodes colorees : GET (vert), POST (turquoise), PUT (lavande), DELETE (corail).",
                    "Acces colores : Public (gris clair), ADMIN (peche clair), Authentifie (bleu clair).",
                    "Endpoints couverts :",
                    " - /api/auth/register (POST), /api/auth/login (POST)",
                    " - /api/products GET (Public), POST/PUT/DELETE (ADMIN)",
                    " - /api/categories GET (Public), POST (ADMIN)",
                    " - /api/cart GET/POST/PUT/DELETE (Authentifie)",
                    " - /api/orders POST/GET (Authentifie)",
                    " - /api/reviews POST (Authentifie), GET product/{id} (Public)",
                ],
            },
        ],
    },
    "Part 3 - Frontend, Demo & Conclusion": {
        "color": PEACH_RL,
        "intro": (
            "Cette troisieme partie presente le frontend Angular 17, le parcours de "
            "demonstration en 8 etapes, les tests unitaires, la difficulte Docker "
            "rencontree et sa solution, le bilan et la feuille de route. "
            "Charger ce dossier dans NotebookLM pour generer les slides 21 a 30."
        ),
        "slides": [
            {
                "num": 21,
                "title": "Facade Moderne : Ingenierie Angular 17",
                "subtitle": "3 piliers Angular",
                "content": [
                    "3 cartes verticales avec icones :",
                    " - Architecture Standalone (turquoise, icone briques) : Suppression totale des NgModule. Code plus concis. Lazy loading par route -> bundle 13 ko.",
                    " - Reactivite via Signals (sky blue, icone onde) : Gestion d'etat moderne. Exemple cart.itemCount() declenche la mise a jour du badge dans le header.",
                    " - Intercepteurs & Cohesion (lavande, icone bouclier) : AuthInterceptor attache le Bearer token. ErrorInterceptor capte 401/4xx -> toasts. AuthGuard verrouille les routes privees. UI cohesive Bootstrap 5.",
                    "Bandeau pilule en bas : 'Experience Utilisateur Fluide'.",
                ],
            },
            {
                "num": 22,
                "title": "Architecture Frontend : core / shared / features",
                "subtitle": "Convention de structure",
                "content": [
                    "3 colonnes (cartes verticales) representant les 3 sous-dossiers :",
                    "Colonne 1 - core/ (turquoise) : models/ (User, Product, Cart, Order, Review), services/ (auth, cart, product, order, review, user, toast), guards/ (auth.guard.ts), interceptors/ (auth, error).",
                    "Colonne 2 - shared/ (lavande) : components/ (header, footer, loader, toast).",
                    "Colonne 3 - features/ (peche) : auth/ (login, register), catalog/ (product-list, product-detail), cart/, orders/, profile/, reviews/ (review-form).",
                    "Note en bas : 'Convention : core (singleton, providedIn: root) / shared (composants reutilises) / features (pages metier).'",
                ],
            },
            {
                "num": 23,
                "title": "Reactivite Signals : exemple CartService",
                "subtitle": "Code TS + diagramme reactif",
                "content": [
                    "Layout 2 colonnes :",
                    "Colonne gauche - Bloc de code TypeScript sur fond noir (~18 lignes) :",
                    " - @Injectable providedIn: 'root'",
                    " - readonly cart = signal<Cart | null>(null);",
                    " - readonly itemCount = signal<number>(0);",
                    " - methode add() qui pipe(tap()) -> update()",
                    " - private update() : this.cart.set(cart); this.itemCount.set(cart.itemCount);",
                    "Colonne droite - 4 cartes etapes du mecanisme reactif :",
                    " 1. Composant Header lit cartSvc.itemCount()",
                    " 2. Angular detecte la dependance au signal",
                    " 3. Si .set() change, le DOM est invalide",
                    " 4. Re-render automatique du badge",
                ],
            },
            {
                "num": 24,
                "title": "Interceptors : Auth + Error",
                "subtitle": "Cohesion HTTP centralisee",
                "content": [
                    "Layout 2 colonnes - 2 blocs de code TypeScript sur fond noir :",
                    "Colonne gauche - auth.interceptor.ts :",
                    " - HttpInterceptorFn fonctionnel",
                    " - inject(AuthService), recupere le token",
                    " - req.clone({ setHeaders: { Authorization: Bearer ${token} } })",
                    "Colonne droite - error.interceptor.ts :",
                    " - catchError sur next(req)",
                    " - 401 -> auth.logout() + toast.error('Session expiree')",
                    " - 4xx -> toast.error(msg)",
                    " - 0 -> toast.error('Serveur injoignable')",
                    "3 cartes resume en bas :",
                    " - AuthInterceptor (turquoise) : attache Bearer JWT.",
                    " - ErrorInterceptor (lavande) : capte 401/4xx/0.",
                    " - AuthGuard (vert) : bloque /cart, /orders, /profile pour anonymes.",
                ],
            },
            {
                "num": 25,
                "title": "Parcours Utilisateur : Tracabilite de la Demonstration",
                "subtitle": "Timeline en 8 etapes",
                "content": [
                    "Timeline horizontale en 8 jalons (cercles numerotes + cartes en dessous) :",
                    " 1. Exploration (vert) : Catalogue 12 produits, recherche, filtre.",
                    " 2. Inspection (turquoise) : Fiche produit - Prix, stock temps reel, avis Mongo.",
                    " 3. Friction (peche) : Tentative d'ajout au panier -> Redirection (Non-authentifie). Note 'Intervention AuthGuard'.",
                    " 4. Authentification (sky blue) : Login user@estore.ma / User@123.",
                    " 5. Conversion (vert) : Ajout de 2 produits (badge Signals passe a 2).",
                    " 6. Ajustement (lavande) : Modification quantites, recalcul total.",
                    " 7. Validation (vert) : Commande confirmee -> Routage vers /orders.",
                    " 8. Fidelisation (or) : Depot avis 5 etoiles persistant sur le produit.",
                    "Cercles relies par une ligne grise. Cartes blanches en dessous avec bordure de couleur.",
                ],
            },
            {
                "num": 26,
                "title": "Tests unitaires : 9 tests, BUILD SUCCESS",
                "subtitle": "Couverture des services critiques",
                "content": [
                    "3 grandes cartes en haut, chacune avec un GROS chiffre coloré :",
                    " - 4 (vert) : ProductServiceTest - findById ok/KO, search avec/sans mot-cle.",
                    " - 2 (turquoise) : CartServiceTest - addItem stock OK / stock insuffisant.",
                    " - 3 (lavande) : OrderServiceTest - checkout panier vide / valide / myOrders trie.",
                    "Bloc console (fond noir, texte vert clair) en bas reproduisant la sortie Maven :",
                    " $ mvn test",
                    " [INFO] T E S T S",
                    " [INFO] Tests run: 4, Failures: 0... ProductServiceTest",
                    " [INFO] Tests run: 2, Failures: 0... CartServiceTest",
                    " [INFO] Tests run: 3, Failures: 0... OrderServiceTest",
                    " [INFO] Tests run: 9, Failures: 0, Errors: 0, Skipped: 0",
                    " [INFO] BUILD SUCCESS",
                ],
            },
            {
                "num": 27,
                "title": "Intelligence DevOps : Le Bifurcateur de Resolution",
                "subtitle": "Difficulte Docker et solution",
                "content": [
                    "En haut : fenetre console style macOS (3 boutons rouge/jaune/vert) avec message d'erreur en blanc :",
                    " 'Obstacle d'Environnement : C:\\ProgramData\\DockerDesktop must be owned by an elevated account (Cause: Dossier residuel systeme)'",
                    "Titre central : 'Intelligence DevOps : Le Bifurcateur de Resolution'.",
                    "Deux fleches divergentes vers 2 cartes solutions :",
                    " - Solution 1 - Brute (gris) : Nettoyage via ligne de commande Administrateur (rmdir /s /q) suivi d'une reinstallation propre.",
                    " - Solution 2 - Elegante (vert) : Activation du profil 'dev' Spring Boot. Bascule automatique sur la base relationnelle 'H2 in-memory'. L'application demarre sans Docker.",
                    "Cercle 'Qualite Garantie' en bas (vert) avec : '9 tests unitaires valides, JUnit/Mockito mode @MockitoSettings(LENIENT), 0 echec, BUILD SUCCESS'.",
                ],
            },
            {
                "num": 28,
                "title": "Bilan : objectifs atteints",
                "subtitle": "6 axes valides",
                "content": [
                    "6 cartes en grille 2x3, chacune commencant par une coche verte (✓) :",
                    " - Architecture (vert) : 3 couches × 5 domaines rigoureusement compartimentee.",
                    " - Securite (peche) : JWT + BCrypt + @PreAuthorize - etanche.",
                    " - Persistance hybride (lavande) : SQL (MySQL/H2) + NoSQL (MongoDB) fonctionnelle.",
                    " - Tests (turquoise) : 9 tests unitaires verts, BUILD SUCCESS reproductible.",
                    " - DevOps (sky blue) : Docker Compose + profil dev H2 (resilience).",
                    " - Documentation (or) : README + 7 PDF + Presentation NotebookLM-style.",
                ],
            },
            {
                "num": 29,
                "title": "Feuille de Route : Projection",
                "subtitle": "4 paliers d'evolution",
                "content": [
                    "Timeline en escalier (4 cubes 3D montant en diagonale) :",
                    " 1. Monetisation (sky blue) : Integration passerelles Stripe / CMI.",
                    " 2. Engagement (vert) : Notifications emails & SMS automatisees.",
                    " 3. Intelligence (peche) : Recommandations ML basees sur l'historique.",
                    " 4. Scalabilite (lavande) : Deploiement Cloud AWS / Azure.",
                    "Fleche directrice montante en degrade.",
                    "Note : 'Direction d'evolution : du business immediat vers la scalabilite cloud'.",
                ],
            },
            {
                "num": 30,
                "title": "Merci pour votre attention",
                "subtitle": "Q & R",
                "content": [
                    "Slide finale style identique a la couverture (avec grille de fond) :",
                    "Texte central tres grand : 'Merci pour votre attention.'",
                    "Sous-texte italique : 'Nous sommes prets a repondre a vos questions.'",
                    "3 pilules colorees en bas :",
                    " - Auteurs (turquoise) : Akram Belmoussa & Nouhaila Ben Soumane.",
                    " - Module Full-Stack (vert) : Pr. Omar Zahour.",
                    " - Etablissement (lavande) : FSBM - Universite Hassan II - 2025-2026.",
                ],
            },
        ],
    },
}


def _get_part_styles():
    base = getSampleStyleSheet()
    return {
        "H1": ParagraphStyle("H1", parent=base["Heading1"], fontName="Helvetica-Bold",
                             fontSize=22, textColor=PRIMARY_RL, spaceAfter=8, leading=26),
        "H2": ParagraphStyle("H2", parent=base["Heading2"], fontName="Helvetica-Bold",
                             fontSize=14, textColor=PRIMARY_RL, spaceBefore=10, spaceAfter=4),
        "H3": ParagraphStyle("H3", parent=base["Heading3"], fontName="Helvetica-Bold",
                             fontSize=12, textColor=HexColor("#4B5563"), spaceBefore=6, spaceAfter=3),
        "body": ParagraphStyle("body", parent=base["BodyText"], fontName="Helvetica",
                               fontSize=11, textColor=PRIMARY_RL, leading=15, alignment=TA_JUSTIFY,
                               spaceAfter=4),
        "code": ParagraphStyle("code", parent=base["Code"], fontName="Courier",
                               fontSize=9.5, textColor=PRIMARY_RL,
                               backColor=HexColor("#F3F4F6"), borderPadding=6,
                               leftIndent=4, rightIndent=4, leading=12),
        "intro": ParagraphStyle("intro", parent=base["BodyText"], fontName="Helvetica-Oblique",
                                fontSize=11, textColor=HexColor("#4B5563"), leading=15,
                                alignment=TA_JUSTIFY, spaceAfter=10,
                                backColor=BG_RL, borderPadding=10,
                                borderColor=HexColor("#D1D5DB"), borderWidth=0.5),
    }


def build_notebooklm_part(part_name, part_data, idx, total_parts):
    """Genere un PDF pour une partie de la presentation NotebookLM."""
    safe_name = part_name.lower().replace(" ", "-").replace("&", "et")
    filename = f"script-notebooklm-part-{idx}.pdf"
    out = os.path.join(OUT_DIR, filename)

    doc = SimpleDocTemplate(out, pagesize=A4,
                            leftMargin=2 * cm, rightMargin=2 * cm,
                            topMargin=2.2 * cm, bottomMargin=1.8 * cm,
                            title=f"Script NotebookLM {part_name}",
                            author="A. Belmoussa & N. Ben Soumane")
    st = _get_part_styles()
    flow = []

    # Page de couverture
    flow.append(Spacer(1, 4 * cm))
    flow.append(Paragraph(
        f"<font color='{part_data['color'].hexval()}'>Script NotebookLM</font>",
        ParagraphStyle("c1", parent=st["body"], fontSize=14, textColor=part_data["color"],
                       alignment=TA_CENTER, fontName="Helvetica-Bold")))
    flow.append(Spacer(1, 0.3 * cm))
    flow.append(Paragraph(part_name, ParagraphStyle(
        "c2", parent=st["H1"], fontSize=32, alignment=TA_CENTER,
        textColor=PRIMARY_RL, leading=38)))
    flow.append(Spacer(1, 0.5 * cm))
    flow.append(Paragraph(
        f"Partie {idx} sur {total_parts}  ·  {len(part_data['slides'])} slides",
        ParagraphStyle("c3", parent=st["body"], fontSize=14, alignment=TA_CENTER,
                       textColor=HexColor("#6B7280"))))
    flow.append(Spacer(1, 1.5 * cm))

    info = Table([
        ["Projet",       "E-Store"],
        ["Auteurs",      "Akram Belmoussa & Nouhaila Ben Soumane"],
        ["Encadrant",    "Pr. Omar Zahour"],
        ["Etablissement", "FSBM - Universite Hassan II"],
        ["Annee",         "2025-2026"],
        ["Usage",         "A injecter dans NotebookLM (limite 10 slides)"],
    ], colWidths=[4 * cm, 11 * cm])
    info.setStyle(TableStyle([
        ("FONTNAME", (0, 0), (0, -1), "Helvetica-Bold"),
        ("FONTNAME", (1, 0), (1, -1), "Helvetica"),
        ("FONTSIZE", (0, 0), (-1, -1), 10.5),
        ("TEXTCOLOR", (0, 0), (0, -1), HexColor("#6B7280")),
        ("TEXTCOLOR", (1, 0), (1, -1), PRIMARY_RL),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 7),
        ("TOPPADDING", (0, 0), (-1, -1), 7),
        ("LINEBELOW", (0, 0), (-1, -1), 0.3, HexColor("#D1D5DB")),
        ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
    ]))
    flow.append(info)
    flow.append(PageBreak())

    # Mode d'emploi
    flow.append(Paragraph("Mode d'emploi", st["H1"]))
    flow.append(Paragraph(part_data["intro"], st["intro"]))
    flow.append(Paragraph("Procedure recommandee", st["H2"]))
    for step in [
        f"<b>1.</b> Ouvrir NotebookLM (notebooklm.google.com).",
        f"<b>2.</b> Creer un nouveau notebook intitule <i>'E-Store - {part_name}'</i>.",
        f"<b>3.</b> Importer ce PDF (script-notebooklm-part-{idx}.pdf) comme source.",
        f"<b>4.</b> Cliquer sur 'Studio' &gt; 'Video Overview' (ou 'Slide Deck').",
        f"<b>5.</b> Dans le prompt personnalise, indiquer : <i>'Genere une presentation visuelle de 10 slides en suivant strictement le contenu du document. Style : minimaliste, blueprint, pilules colorees, schemas type mind-map ou matrice, couleurs vert mint / lavande / sky blue / peche.'</i>",
        f"<b>6.</b> Lancer la generation et exporter en PDF.",
        f"<b>7.</b> Repeter avec les parties {idx % total_parts + 1}, etc.",
    ]:
        flow.append(Paragraph(step, st["body"]))
    flow.append(PageBreak())

    # Pour chaque slide
    for sd in part_data["slides"]:
        flow.append(Paragraph(f"Slide {sd['num']} — {sd['title']}", st["H1"]))
        if sd.get("subtitle"):
            flow.append(Paragraph(f"<i>{sd['subtitle']}</i>", st["H3"]))
        flow.append(Paragraph("Contenu detaille a faire generer par NotebookLM :", st["H2"]))
        for line in sd["content"]:
            if line.startswith(" -") or line.startswith("  -") or line.startswith(" ") and "." in line[:5]:
                flow.append(Paragraph(line, ParagraphStyle("li", parent=st["body"],
                                                            leftIndent=0.6 * cm)))
            else:
                flow.append(Paragraph(line, st["body"]))
        flow.append(PageBreak())

    # Header / footer
    def hf(canvas, doc):
        canvas.saveState()
        canvas.setStrokeColor(part_data["color"])
        canvas.setLineWidth(0.6)
        canvas.line(2 * cm, A4[1] - 1.6 * cm, A4[0] - 2 * cm, A4[1] - 1.6 * cm)
        canvas.setFillColor(part_data["color"])
        canvas.setFont("Helvetica-Bold", 9)
        canvas.drawString(2 * cm, A4[1] - 1.4 * cm, f"NOTEBOOKLM SCRIPT — {part_name.upper()}")
        canvas.setFillColor(HexColor("#6B7280"))
        canvas.drawRightString(A4[0] - 2 * cm, A4[1] - 1.4 * cm, f"Partie {idx}/{total_parts}")
        canvas.setFont("Helvetica", 8)
        canvas.drawString(2 * cm, 1.2 * cm, "E-Store · A. Belmoussa & N. Ben Soumane")
        canvas.drawRightString(A4[0] - 2 * cm, 1.2 * cm, f"Page {doc.page}")
        canvas.restoreState()

    doc.build(flow, onFirstPage=lambda c, d: None, onLaterPages=hf)
    return out


def build_all_notebooklm_scripts():
    paths = []
    parts = list(NOTEBOOK_SCRIPTS.items())
    for i, (name, data) in enumerate(parts, 1):
        path = build_notebooklm_part(name, data, i, len(parts))
        paths.append(path)
    return paths


# =====================================================================
# MAIN
# =====================================================================
if __name__ == "__main__":
    print("=" * 70)
    print("Generation E-Store : presentation avancee + scripts NotebookLM")
    print("=" * 70)
    pptx = build_pptx()
    print(f"  [OK] {os.path.basename(pptx):50s} ({os.path.getsize(pptx)//1024} ko)")
    for p in build_all_notebooklm_scripts():
        print(f"  [OK] {os.path.basename(p):50s} ({os.path.getsize(p)//1024} ko)")
    print("=" * 70)
    print(f"Fichiers generes dans : {OUT_DIR}")

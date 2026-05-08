# -*- coding: utf-8 -*-
"""
Genere la presentation de soutenance E-Store en :
  - presentation-soutenance.pptx (PowerPoint)
  - presentation-soutenance.pdf  (version PDF des slides)
  - guide-soutenance-discours.pdf (script de presentation pour Akram + Nouhaila)

Usage : python generate_presentation.py
"""
import os
import sys
import io

# Force UTF-8 sur la console Windows
if sys.stdout.encoding and sys.stdout.encoding.lower() != "utf-8":
    try:
        sys.stdout.reconfigure(encoding="utf-8", errors="replace")
    except Exception:
        sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from reportlab.lib.pagesizes import A4, landscape
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import cm
from reportlab.lib.colors import HexColor, black, white, grey
from reportlab.lib.enums import TA_LEFT, TA_CENTER, TA_JUSTIFY
from reportlab.platypus import (
    SimpleDocTemplate, Paragraph, Spacer, PageBreak,
    Table, TableStyle, Preformatted
)

OUT_DIR = os.path.dirname(os.path.abspath(__file__))

# Palette
PRIMARY_HEX = "#0d6efd"
DARK_HEX = "#212529"
LIGHT_HEX = "#f5f7fa"
PRIMARY = RGBColor(0x0D, 0x6E, 0xFD)
DARK = RGBColor(0x21, 0x25, 0x29)
LIGHT = RGBColor(0xF5, 0xF7, 0xFA)
ACCENT = RGBColor(0x19, 0x87, 0x54)  # vert
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREY = RGBColor(0x6c, 0x75, 0x7d)


# =====================================================================
# CONTENU DES SLIDES (commun PPTX et PDF)
# =====================================================================
SLIDES = [
    # ── Slide 1 ── Couverture
    {
        "type": "cover",
        "title": "E-Store",
        "subtitle": "Application e-commerce full-stack",
        "info": [
            ("Module", "Full-Stack"),
            ("Encadrant", "Pr. Omar Zahour"),
            ("Etablissement", "FSBM - Universite Hassan II"),
            ("Annee", "2025-2026"),
            ("Auteurs", "Akram Belmoussa & Nouhaila Ben Soumane"),
        ],
        "duration": "30 s",
    },

    # ── Slide 2 ── Plan
    {
        "type": "content",
        "title": "Plan de la presentation",
        "bullets": [
            "1.  Contexte et objectifs du projet",
            "2.  Architecture (3 couches x 5 domaines)",
            "3.  Stack technique",
            "4.  Backend Spring Boot - exemples de code",
            "5.  Securite (Spring Security + JWT)",
            "6.  MongoDB pour les avis (persistance hybride)",
            "7.  Frontend Angular 17",
            "8.  Demonstration en direct",
            "9.  Tests, DevOps et difficulte Docker",
            "10. Conclusion et perspectives",
        ],
        "duration": "30 s",
    },

    # ── Slide 3 ── Contexte
    {
        "type": "content",
        "title": "1. Contexte et objectifs",
        "subtitle": "Mini-projet pedagogique e-commerce",
        "bullets": [
            "Application complete : inscription, catalogue, panier, commandes, avis",
            "Objectif 1 : maitriser une architecture full-stack moderne",
            "Objectif 2 : illustrer la persistance hybride SQL + NoSQL",
            "Objectif 3 : appliquer les bonnes pratiques (DTOs, transactions, tests)",
            "Objectif 4 : securiser via JWT + Spring Security 6",
            "Objectif 5 : produire une UI moderne avec Angular 17 standalone",
        ],
        "duration": "1 min",
    },

    # ── Slide 4 ── Architecture
    {
        "type": "architecture",
        "title": "2. Architecture en 3 couches x 5 domaines",
        "duration": "1 min 30",
    },

    # ── Slide 5 ── Stack technique
    {
        "type": "table",
        "title": "3. Stack technique",
        "headers": ["Couche", "Technologies"],
        "rows": [
            ["Frontend", "Angular 17 standalone, Bootstrap 5, RxJS, TypeScript 5"],
            ["Backend", "Spring Boot 3.3, Spring Security 6, JWT (jjwt 0.12)"],
            ["ORM / Donnees", "Spring Data JPA + Hibernate, Spring Data MongoDB"],
            ["Bases", "MySQL 8 (prod) / H2 (dev) + MongoDB 7"],
            ["Build", "Maven 3.9, npm, Angular CLI"],
            ["Tests", "JUnit 5 + Mockito (9 tests, BUILD SUCCESS)"],
            ["DevOps", "Docker Compose : MySQL + Mongo + phpMyAdmin + mongo-express"],
        ],
        "duration": "30 s",
    },

    # ── Slide 6 ── Backend code
    {
        "type": "code",
        "title": "4. Backend - exemple : checkout transactionnel",
        "subtitle": "OrderService.checkout() - atomicite ACID",
        "code": """@Transactional
public OrderDto checkout() {
    Cart cart = cartService.getOrCreateCart(user);
    if (cart.getItems().isEmpty())
        throw new BusinessException("Votre panier est vide");

    // 1) Verifier le stock pour TOUS les items
    for (CartItem ci : cart.getItems())
        inventoryService.checkAvailability(ci.getProduct().getId(), ci.getQuantity());

    // 2) Creer la commande + decrementer stock
    Order order = Order.builder().user(user).status(PENDING).build();
    for (CartItem ci : cart.getItems()) {
        order.getItems().add(OrderItem.builder()...build());
        inventoryService.decrement(ci.getProduct().getId(), ci.getQuantity());
    }
    order.setStatus(CONFIRMED);
    Order saved = orderRepository.save(order);

    // 3) Vider le panier
    cartService.clearCart(cart);
    return OrderDto.from(saved);
}""",
        "duration": "1 min 30",
    },

    # ── Slide 7 ── Securite JWT
    {
        "type": "security",
        "title": "5. Securite - JWT + Spring Security 6",
        "duration": "1 min",
    },

    # ── Slide 8 ── MongoDB
    {
        "type": "content",
        "title": "6. MongoDB pour les avis (persistance hybride)",
        "subtitle": "Pourquoi un second SGBD ?",
        "bullets": [
            "Avis = donnees semi-structurees, volumineuses, en lecture frequente",
            "Pas de relations strictes - schema flexible",
            "Spring Data MongoDB : meme abstraction que JPA",
            "Repository : findByProductIdOrderByCreatedAtDesc(productId)",
            "Document type :",
            "    { productId, userId, authorName, rating(1-5), comment, createdAt }",
            "Comportement defensif : si Mongo absent, l'app continue (warning)",
        ],
        "duration": "30 s",
    },

    # ── Slide 9 ── Frontend Angular
    {
        "type": "content",
        "title": "7. Frontend Angular 17",
        "subtitle": "Composants standalone + signals",
        "bullets": [
            "Standalone components (pas de NgModule) - moderne et concis",
            "Lazy loading par route : bundle initial 13 ko + features a la demande",
            "Signals pour l'etat reactif : cart.itemCount() dans le header",
            "AuthService : BehaviorSubject<User> + localStorage persistance",
            "AuthInterceptor : ajoute Bearer <token> a chaque requete",
            "ErrorInterceptor : gestion centralisee 401/4xx/0 (toast)",
            "AuthGuard sur /cart, /orders, /profile",
            "Bootstrap 5 + thematique cohesive, responsive",
        ],
        "duration": "1 min",
    },

    # ── Slide 10 ── Demonstration
    {
        "type": "demo",
        "title": "8. Demonstration en direct",
        "duration": "1 min",
    },

    # ── Slide 11 ── Tests + Docker
    {
        "type": "two_columns",
        "title": "9. Tests, DevOps et la difficulte Docker",
        "left_title": "Tests unitaires",
        "left_bullets": [
            "ProductServiceTest (4 tests)",
            "CartServiceTest (2 tests)",
            "OrderServiceTest (3 tests)",
            "TOTAL : 9 tests, 0 failure",
            "Mockito @MockitoSettings(LENIENT)",
        ],
        "right_title": "Difficulte Docker rencontree",
        "right_bullets": [
            "Erreur : 'C:\\\\ProgramData\\\\DockerDesktop must be owned",
            "    by an elevated account'",
            "Cause : dossier residuel d'une install precedente",
            "Solution 1 : rmdir /s /q en admin + reinstall",
            "Solution 2 : profil dev (H2 in-memory) - pas d'install",
            "=> Architecture pensee pour 2 modes des le depart",
        ],
        "duration": "30 s",
    },

    # ── Slide 12 ── Conclusion
    {
        "type": "content",
        "title": "10. Conclusion et perspectives",
        "subtitle": "Bilan",
        "bullets": [
            "[+] Architecture 3 couches x 5 domaines respectee",
            "[+] Securite robuste (JWT + BCrypt + @PreAuthorize)",
            "[+] Persistance hybride SQL + NoSQL operationnelle",
            "[+] 9 tests unitaires verts, build reproductible",
            "",
            "Perspectives :",
            "  - Paiement en ligne (Stripe / CMI)",
            "  - Notifications (emails de confirmation, SMS)",
            "  - Recommandations produits (ML)",
            "  - Deploiement cloud (Heroku, AWS, Azure)",
            "",
            "Merci pour votre attention - Questions ?",
        ],
        "duration": "30 s",
    },
]


# =====================================================================
# GENERATEUR PPTX
# =====================================================================
def add_footer(slide, prs, page_num, total):
    """Ajoute un footer discret sur la slide."""
    left = Inches(0.3)
    top = prs.slide_height - Inches(0.4)
    width = prs.slide_width - Inches(0.6)
    tx = slide.shapes.add_textbox(left, top, width, Inches(0.3))
    tf = tx.text_frame
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.text = f"E-Store  |  A. Belmoussa & N. Ben Soumane  |  Pr. O. Zahour  |  {page_num}/{total}"
    p.alignment = PP_ALIGN.CENTER
    for run in p.runs:
        run.font.size = Pt(9)
        run.font.color.rgb = GREY


def add_title_bar(slide, prs, title, subtitle=None):
    """Bandeau bleu en haut avec titre."""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0,
        prs.slide_width, Inches(0.9)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()

    tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), prs.slide_width - Inches(1), Inches(0.6))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.runs[0].font.size = Pt(28)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = WHITE

    if subtitle:
        tx2 = slide.shapes.add_textbox(Inches(0.5), Inches(0.95), prs.slide_width - Inches(1), Inches(0.4))
        tf2 = tx2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.runs[0].font.size = Pt(14)
        p2.runs[0].font.italic = True
        p2.runs[0].font.color.rgb = DARK


def add_bullets(slide, prs, bullets, top=Inches(1.6)):
    """Liste a puces."""
    tx = slide.shapes.add_textbox(
        Inches(0.7), top,
        prs.slide_width - Inches(1.4), prs.slide_height - top - Inches(0.7)
    )
    tf = tx.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "  " + b if b else ""
        p.space_after = Pt(6)
        for run in p.runs:
            run.font.size = Pt(16)
            run.font.color.rgb = DARK


def add_text_box(slide, left, top, width, height, text, size=14, color=DARK, bold=False, align=None):
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    if align:
        p.alignment = align
    for run in p.runs:
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tx


def build_cover_slide(prs, slide_data, page, total):
    """Slide de couverture."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Bandeau gauche bleu (decoration)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(1.5), prs.slide_height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()

    # Titre principal
    add_text_box(slide, Inches(2), Inches(1.5), Inches(11), Inches(1.5),
                 slide_data["title"], size=72, color=PRIMARY, bold=True)

    # Sous-titre
    add_text_box(slide, Inches(2), Inches(3), Inches(11), Inches(0.7),
                 slide_data["subtitle"], size=22, color=DARK)

    # Trait de separation
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2), Inches(3.8),
                                  Inches(2), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = PRIMARY
    line.line.fill.background()

    # Bloc info (label + valeur)
    top = Inches(4.3)
    for label, value in slide_data["info"]:
        add_text_box(slide, Inches(2), top, Inches(2.5), Inches(0.4),
                     label, size=12, color=GREY, bold=True)
        add_text_box(slide, Inches(4.5), top, Inches(8.5), Inches(0.4),
                     value, size=14, color=DARK)
        top += Inches(0.45)

    add_footer(slide, prs, page, total)
    return slide


def build_content_slide(prs, slide_data, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, prs, slide_data["title"], slide_data.get("subtitle"))
    add_bullets(slide, prs, slide_data["bullets"])
    add_footer(slide, prs, page, total)
    return slide


def build_table_slide(prs, slide_data, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, prs, slide_data["title"], slide_data.get("subtitle"))

    rows = len(slide_data["rows"]) + 1
    cols = len(slide_data["headers"])
    left = Inches(0.7)
    top = Inches(1.8)
    width = prs.slide_width - Inches(1.4)
    height = Inches(0.4) * rows

    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    tbl = table_shape.table

    # Header
    for i, h in enumerate(slide_data["headers"]):
        cell = tbl.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.bold = True
                r.font.size = Pt(14)
                r.font.color.rgb = WHITE

    # Body
    for r_idx, row in enumerate(slide_data["rows"], 1):
        for c_idx, val in enumerate(row):
            cell = tbl.cell(r_idx, c_idx)
            cell.text = val
            if r_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(12)
                    run.font.color.rgb = DARK
                    if c_idx == 0:
                        run.font.bold = True

    add_footer(slide, prs, page, total)
    return slide


def build_code_slide(prs, slide_data, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, prs, slide_data["title"], slide_data.get("subtitle"))

    # Cadre code (fond gris fonce)
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(0.5), Inches(1.7),
                                 prs.slide_width - Inches(1), prs.slide_height - Inches(2.4))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0x1e, 0x1e, 0x1e)
    box.line.fill.background()

    # Texte du code
    tx = slide.shapes.add_textbox(Inches(0.7), Inches(1.85),
                                  prs.slide_width - Inches(1.4),
                                  prs.slide_height - Inches(2.7))
    tf = tx.text_frame
    tf.word_wrap = True
    lines = slide_data["code"].split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        for run in p.runs:
            run.font.name = "Consolas"
            run.font.size = Pt(11)
            run.font.color.rgb = RGBColor(0xdc, 0xdc, 0xdc)

    add_footer(slide, prs, page, total)
    return slide


def build_architecture_slide(prs, slide_data, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, prs, slide_data["title"])

    # 3 boites verticales (couches)
    layers = [
        ("PRESENTATION (Angular)",        "Composants, services HTTP, formulaires"),
        ("LOGIQUE METIER (Spring Boot)",  "Controllers REST, Services, DTOs"),
        ("ACCES AUX DONNEES",             "JpaRepository + MongoRepository"),
    ]
    top = Inches(1.7)
    for title, desc in layers:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(0.7), top,
                                     Inches(6.5), Inches(1.0))
        box.fill.solid()
        box.fill.fore_color.rgb = PRIMARY
        box.line.fill.background()

        tx = box.text_frame
        tx.margin_left = Pt(10)
        p = tx.paragraphs[0]
        p.text = title
        for r in p.runs:
            r.font.bold = True
            r.font.size = Pt(16)
            r.font.color.rgb = WHITE
        p2 = tx.add_paragraph()
        p2.text = desc
        for r in p2.runs:
            r.font.size = Pt(11)
            r.font.italic = True
            r.font.color.rgb = LIGHT
        top += Inches(1.15)

    # 5 boites a droite (domaines)
    add_text_box(slide, Inches(7.7), Inches(1.65), Inches(5), Inches(0.4),
                 "5 DOMAINES FONCTIONNELS", size=14, color=DARK, bold=True)

    domains = [
        ("customer",  "User, Profile, JWT"),
        ("catalog",   "Categorie, Produit, recherche"),
        ("inventory", "Stock par produit"),
        ("shopping",  "Cart, CartItem"),
        ("billing",   "Order, OrderItem - transactionnel"),
        ("review",    "Avis (MongoDB)"),
    ]
    top = Inches(2.1)
    for name, desc in domains:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(7.7), top, Inches(5), Inches(0.55))
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT
        box.line.color.rgb = PRIMARY

        tx = box.text_frame
        tx.margin_left = Pt(8)
        p = tx.paragraphs[0]
        p.text = f"{name}  -  {desc}"
        for r in p.runs:
            r.font.size = Pt(12)
            r.font.color.rgb = DARK
        top += Inches(0.62)

    add_footer(slide, prs, page, total)
    return slide


def build_security_slide(prs, slide_data, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, prs, slide_data["title"])

    # Schema flux JWT
    steps = [
        ("1.", "Client envoie email + password",     "POST /api/auth/login"),
        ("2.", "AuthService valide via BCrypt",       "passwordEncoder.matches(...)"),
        ("3.", "JwtService genere token signe HS256", "Jwts.builder().signWith(key)"),
        ("4.", "Client stocke le token (localStorage)", "estore.token"),
        ("5.", "Toute requete suivante : Bearer ...",   "AuthInterceptor"),
        ("6.", "JwtFilter valide a chaque appel",      "doFilterInternal()"),
    ]
    top = Inches(1.7)
    for n, label, code in steps:
        # Numero (cercle)
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), top,
                                        Inches(0.5), Inches(0.5))
        circle.fill.solid()
        circle.fill.fore_color.rgb = PRIMARY
        circle.line.fill.background()
        tf = circle.text_frame
        tf.margin_top = Pt(0); tf.margin_bottom = Pt(0)
        p = tf.paragraphs[0]
        p.text = n; p.alignment = PP_ALIGN.CENTER
        for r in p.runs:
            r.font.bold = True; r.font.size = Pt(14); r.font.color.rgb = WHITE

        # Label
        add_text_box(slide, Inches(1.4), top + Inches(0.05), Inches(7.5), Inches(0.5),
                     label, size=14, color=DARK)
        # Code
        add_text_box(slide, Inches(8.5), top + Inches(0.08), Inches(4.7), Inches(0.5),
                     code, size=11, color=PRIMARY)
        top += Inches(0.6)

    # Bloc bas : points cles
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(5.6),
                                 prs.slide_width - Inches(1.4), Inches(1.4))
    box.fill.solid(); box.fill.fore_color.rgb = LIGHT; box.line.fill.background()

    tx = box.text_frame
    tx.margin_left = Pt(15); tx.margin_top = Pt(8)
    p = tx.paragraphs[0]
    p.text = "Points cles"
    for r in p.runs: r.font.bold = True; r.font.size = Pt(14); r.font.color.rgb = PRIMARY
    for txt in [
        "- Mots de passe hashes BCrypt",
        "- Token JWT HS256, expiration 24h",
        "- Endpoints publics : /api/auth/**, GET /api/products /categories /reviews",
        "- @PreAuthorize(\"hasRole('ADMIN')\") sur les endpoints d'administration",
    ]:
        p2 = tx.add_paragraph()
        p2.text = txt
        for r in p2.runs: r.font.size = Pt(11); r.font.color.rgb = DARK

    add_footer(slide, prs, page, total)
    return slide


def build_two_columns_slide(prs, slide_data, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, prs, slide_data["title"])

    half = (prs.slide_width - Inches(2)) / 2
    # Colonne gauche
    add_text_box(slide, Inches(0.7), Inches(1.7), half, Inches(0.5),
                 slide_data["left_title"], size=18, color=PRIMARY, bold=True)
    tx = slide.shapes.add_textbox(Inches(0.7), Inches(2.3), half, Inches(4.5))
    tf = tx.text_frame; tf.word_wrap = True
    for i, b in enumerate(slide_data["left_bullets"]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "- " + b
        p.space_after = Pt(6)
        for r in p.runs: r.font.size = Pt(13); r.font.color.rgb = DARK

    # Colonne droite
    add_text_box(slide, Inches(0.9) + half, Inches(1.7), half, Inches(0.5),
                 slide_data["right_title"], size=18, color=ACCENT, bold=True)
    tx = slide.shapes.add_textbox(Inches(0.9) + half, Inches(2.3), half, Inches(4.5))
    tf = tx.text_frame; tf.word_wrap = True
    for i, b in enumerate(slide_data["right_bullets"]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "- " + b
        p.space_after = Pt(6)
        for r in p.runs: r.font.size = Pt(13); r.font.color.rgb = DARK

    add_footer(slide, prs, page, total)
    return slide


def build_demo_slide(prs, slide_data, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, prs, slide_data["title"])

    add_text_box(slide, Inches(1), Inches(2), Inches(11), Inches(0.7),
                 "Scenario de demonstration en 8 etapes (passage au navigateur)",
                 size=18, color=PRIMARY, bold=True, align=PP_ALIGN.CENTER)

    steps = [
        "1. Catalogue : 12 produits, recherche, filtre par categorie",
        "2. Fiche produit : prix, stock, avis",
        "3. Tentative panier sans connexion -> redirection /login",
        "4. Connexion user@estore.ma / User@123",
        "5. Ajout de 2 produits au panier (badge passe a 2)",
        "6. Panier : modification quantite, total recalcule",
        "7. Validation commande -> /orders, commande #1 confirmee",
        "8. Depot d'un avis 5 etoiles sur un produit",
    ]
    top = Inches(3)
    for st in steps:
        add_text_box(slide, Inches(1.5), top, Inches(11), Inches(0.4),
                     st, size=14, color=DARK)
        top += Inches(0.4)

    add_footer(slide, prs, page, total)
    return slide


def build_pptx():
    """Construit la presentation .pptx."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    total = len(SLIDES)
    for i, sd in enumerate(SLIDES, 1):
        t = sd["type"]
        if t == "cover":
            build_cover_slide(prs, sd, i, total)
        elif t == "content":
            build_content_slide(prs, sd, i, total)
        elif t == "table":
            build_table_slide(prs, sd, i, total)
        elif t == "code":
            build_code_slide(prs, sd, i, total)
        elif t == "architecture":
            build_architecture_slide(prs, sd, i, total)
        elif t == "security":
            build_security_slide(prs, sd, i, total)
        elif t == "two_columns":
            build_two_columns_slide(prs, sd, i, total)
        elif t == "demo":
            build_demo_slide(prs, sd, i, total)

    out = os.path.join(OUT_DIR, "presentation-soutenance.pptx")
    prs.save(out)
    return out


# =====================================================================
# GENERATEUR PDF (slides en paysage)
# =====================================================================
PRIMARY_RL = HexColor(PRIMARY_HEX)
DARK_RL = HexColor(DARK_HEX)
LIGHT_RL = HexColor(LIGHT_HEX)
ACCENT_RL = HexColor("#198754")

def slide_pdf_styles():
    base = getSampleStyleSheet()
    return {
        "title": ParagraphStyle("title", parent=base["Heading1"], fontName="Helvetica-Bold",
                                fontSize=26, textColor=white, leading=30, spaceAfter=4),
        "subtitle": ParagraphStyle("subtitle", parent=base["Heading2"], fontName="Helvetica-Oblique",
                                   fontSize=14, textColor=DARK_RL, leading=18, spaceAfter=10),
        "bullet": ParagraphStyle("bullet", parent=base["BodyText"], fontName="Helvetica",
                                 fontSize=15, textColor=DARK_RL, leading=22, leftIndent=20,
                                 bulletIndent=8, spaceAfter=4),
        "body": ParagraphStyle("body", parent=base["BodyText"], fontName="Helvetica",
                               fontSize=14, textColor=DARK_RL, leading=20, alignment=TA_LEFT),
        "code": ParagraphStyle("code", parent=base["Code"], fontName="Courier",
                               fontSize=10, textColor=HexColor("#dcdcdc"),
                               backColor=HexColor("#1e1e1e"), borderPadding=10,
                               leftIndent=4, rightIndent=4, leading=12),
        "cover_title": ParagraphStyle("cover_title", parent=base["Heading1"],
                                      fontName="Helvetica-Bold", fontSize=64,
                                      textColor=PRIMARY_RL, alignment=TA_LEFT, leading=68),
        "cover_subtitle": ParagraphStyle("cover_sub", parent=base["BodyText"],
                                         fontName="Helvetica", fontSize=20,
                                         textColor=DARK_RL, alignment=TA_LEFT, spaceAfter=20),
    }


def slide_header_footer(canvas, doc):
    """Bandeau bleu en haut + footer."""
    canvas.saveState()
    # Bandeau bleu titre (35 mm de haut)
    if hasattr(doc, "slide_title") and doc.slide_title:
        canvas.setFillColor(PRIMARY_RL)
        canvas.rect(0, A4[0] - 2.5 * cm, A4[1], 2.5 * cm, fill=1, stroke=0)
        canvas.setFillColor(white)
        canvas.setFont("Helvetica-Bold", 22)
        canvas.drawString(1.5 * cm, A4[0] - 1.5 * cm, doc.slide_title)
        if hasattr(doc, "slide_subtitle") and doc.slide_subtitle:
            canvas.setFont("Helvetica-Oblique", 12)
            canvas.drawString(1.5 * cm, A4[0] - 2.1 * cm, doc.slide_subtitle)
    # Footer
    canvas.setFillColor(grey)
    canvas.setFont("Helvetica", 9)
    canvas.drawCentredString(A4[1] / 2, 0.8 * cm,
                             f"E-Store  |  A. Belmoussa & N. Ben Soumane  |  Pr. O. Zahour  |  {doc.page}/{len(SLIDES)}")
    canvas.restoreState()


class SlideDocTemplate(SimpleDocTemplate):
    """DocTemplate qui transmet le titre courant au header_footer."""
    def __init__(self, *args, **kwargs):
        super().__init__(*args, **kwargs)
        self.slide_title = ""
        self.slide_subtitle = ""


def build_pdf():
    """Genere le PDF des slides en format paysage."""
    out = os.path.join(OUT_DIR, "presentation-soutenance.pdf")
    page = landscape(A4)
    doc = SlideDocTemplate(
        out, pagesize=page,
        leftMargin=1.5 * cm, rightMargin=1.5 * cm,
        topMargin=3.5 * cm, bottomMargin=1.5 * cm,
        title="E-Store - Presentation de soutenance",
        author="A. Belmoussa & N. Ben Soumane"
    )
    st = slide_pdf_styles()

    flow = []
    total = len(SLIDES)

    def render_cover(sd, page_num):
        # Pour la couverture on fait sans bandeau — on ajoute manuellement
        flow.append(Spacer(1, 1 * cm))
        flow.append(Paragraph("E-STORE", st["cover_title"]))
        flow.append(Paragraph(sd["subtitle"], st["cover_subtitle"]))
        flow.append(Spacer(1, 0.5 * cm))
        rows = [[lab, val] for lab, val in sd["info"]]
        t = Table(rows, colWidths=[5 * cm, 14 * cm])
        t.setStyle(TableStyle([
            ("FONTNAME", (0, 0), (0, -1), "Helvetica-Bold"),
            ("FONTNAME", (1, 0), (1, -1), "Helvetica"),
            ("FONTSIZE", (0, 0), (-1, -1), 13),
            ("TEXTCOLOR", (0, 0), (0, -1), grey),
            ("TEXTCOLOR", (1, 0), (1, -1), DARK_RL),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 8),
            ("TOPPADDING", (0, 0), (-1, -1), 8),
            ("LINEBELOW", (0, 0), (-1, -1), 0.3, grey),
            ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
        ]))
        flow.append(t)

    def render_content(sd):
        if sd.get("subtitle"):
            flow.append(Paragraph(sd["subtitle"], st["subtitle"]))
        for b in sd["bullets"]:
            if not b:
                flow.append(Spacer(1, 0.2 * cm))
                continue
            flow.append(Paragraph(b, st["bullet"]))

    def render_table(sd):
        rows = [sd["headers"]] + sd["rows"]
        t = Table(rows, colWidths=[5 * cm, 18 * cm])
        t.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), PRIMARY_RL),
            ("TEXTCOLOR", (0, 0), (-1, 0), white),
            ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
            ("FONTNAME", (0, 1), (0, -1), "Helvetica-Bold"),
            ("FONTSIZE", (0, 0), (-1, -1), 11),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 7),
            ("TOPPADDING", (0, 0), (-1, -1), 7),
            ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
            ("GRID", (0, 0), (-1, -1), 0.4, grey),
            ("ROWBACKGROUNDS", (0, 1), (-1, -1), [white, LIGHT_RL]),
        ]))
        flow.append(t)

    def render_code(sd):
        if sd.get("subtitle"):
            flow.append(Paragraph(sd["subtitle"], st["subtitle"]))
        flow.append(Preformatted(sd["code"], st["code"]))

    def render_arch(sd):
        # 3 couches gauche
        flow.append(Paragraph("3 couches techniques", st["subtitle"]))
        layers = Table([
            ["PRESENTATION (Angular)", "Composants standalone, services HTTP, formulaires"],
            ["LOGIQUE METIER (Spring Boot)", "Controllers REST, Services @Transactional, DTOs"],
            ["ACCES AUX DONNEES", "Spring Data JPA + Spring Data MongoDB"],
        ], colWidths=[7 * cm, 16 * cm])
        layers.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (0, -1), PRIMARY_RL),
            ("TEXTCOLOR", (0, 0), (0, -1), white),
            ("FONTNAME", (0, 0), (0, -1), "Helvetica-Bold"),
            ("FONTSIZE", (0, 0), (-1, -1), 12),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 10),
            ("TOPPADDING", (0, 0), (-1, -1), 10),
            ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
            ("GRID", (0, 0), (-1, -1), 0.4, grey),
        ]))
        flow.append(layers)

        flow.append(Spacer(1, 0.4 * cm))
        flow.append(Paragraph("6 sous-packages domaine", st["subtitle"]))
        domains = Table([
            ["customer",  "User, Profile, JWT, Auth"],
            ["catalog",   "Categorie, Produit, recherche"],
            ["inventory", "Stock par produit"],
            ["shopping",  "Cart, CartItem"],
            ["billing",   "Order, OrderItem (transactionnel)"],
            ["review",    "Avis (MongoDB)"],
        ], colWidths=[5 * cm, 18 * cm])
        domains.setStyle(TableStyle([
            ("FONTNAME", (0, 0), (0, -1), "Helvetica-Bold"),
            ("BACKGROUND", (0, 0), (0, -1), LIGHT_RL),
            ("FONTSIZE", (0, 0), (-1, -1), 11),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 6),
            ("GRID", (0, 0), (-1, -1), 0.4, grey),
        ]))
        flow.append(domains)

    def render_security(sd):
        steps = [
            ["1", "Client envoie email + password",      "POST /api/auth/login"],
            ["2", "AuthService valide via BCrypt",       "passwordEncoder.matches(...)"],
            ["3", "JwtService genere token signe HS256", "Jwts.builder().signWith(key)"],
            ["4", "Client stocke le token",              "localStorage 'estore.token'"],
            ["5", "Toute requete : Bearer <token>",      "AuthInterceptor"],
            ["6", "JwtFilter valide a chaque appel",     "JwtAuthenticationFilter"],
        ]
        t = Table(steps, colWidths=[1.2 * cm, 11 * cm, 11 * cm])
        t.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (0, -1), PRIMARY_RL),
            ("TEXTCOLOR", (0, 0), (0, -1), white),
            ("FONTNAME", (0, 0), (0, -1), "Helvetica-Bold"),
            ("FONTNAME", (2, 0), (2, -1), "Courier-Bold"),
            ("TEXTCOLOR", (2, 0), (2, -1), PRIMARY_RL),
            ("ALIGN", (0, 0), (0, -1), "CENTER"),
            ("FONTSIZE", (0, 0), (-1, -1), 11),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 8),
            ("TOPPADDING", (0, 0), (-1, -1), 8),
            ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
            ("LINEBELOW", (0, 0), (-1, -1), 0.3, grey),
        ]))
        flow.append(t)
        flow.append(Spacer(1, 0.4 * cm))
        for txt in [
            "<b>BCrypt</b> pour les mots de passe (cout 10).",
            "<b>JWT HS256</b>, expiration 24h, claims : uid, role, name.",
            "<b>Endpoints publics</b> : /api/auth/**, GET /products /categories /reviews.",
            "<b>@PreAuthorize(\"hasRole('ADMIN')\")</b> sur les operations d'administration.",
        ]:
            flow.append(Paragraph("- " + txt, st["bullet"]))

    def render_two_columns(sd):
        rows = [
            [
                Paragraph(f"<b><font color='{PRIMARY_HEX}'>{sd['left_title']}</font></b>", st["body"]),
                Paragraph(f"<b><font color='#198754'>{sd['right_title']}</font></b>", st["body"]),
            ]
        ]
        max_len = max(len(sd["left_bullets"]), len(sd["right_bullets"]))
        for i in range(max_len):
            l = sd["left_bullets"][i] if i < len(sd["left_bullets"]) else ""
            r = sd["right_bullets"][i] if i < len(sd["right_bullets"]) else ""
            rows.append([
                Paragraph("- " + l if l else "", st["body"]),
                Paragraph("- " + r if r else "", st["body"]),
            ])
        t = Table(rows, colWidths=[12 * cm, 12 * cm])
        t.setStyle(TableStyle([
            ("VALIGN", (0, 0), (-1, -1), "TOP"),
            ("FONTSIZE", (0, 0), (-1, -1), 11),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
        ]))
        flow.append(t)

    def render_demo(sd):
        flow.append(Paragraph("Scenario de demonstration en 8 etapes", st["subtitle"]))
        steps = [
            "1. Catalogue : 12 produits, recherche, filtre par categorie",
            "2. Fiche produit : prix, stock, avis",
            "3. Tentative panier sans connexion -> redirection /login",
            "4. Connexion user@estore.ma / User@123",
            "5. Ajout de 2 produits au panier (badge passe a 2)",
            "6. Panier : modification quantite, total recalcule",
            "7. Validation commande -> /orders, commande #1 confirmee",
            "8. Depot d'un avis 5 etoiles sur un produit",
        ]
        for s_ in steps:
            flow.append(Paragraph(s_, st["bullet"]))

    for i, sd in enumerate(SLIDES, 1):
        # Pour chaque slide on assigne le titre courant qui sera dessine par le canvas
        # On utilise un renderer specifique par type
        # Comme reportlab nous oblige a un seul Document, on contourne en stockant
        # le titre courant dans une liste, et on lit le bon titre dans header_footer
        # via doc.page.
        # Approche plus simple : pre-calculer les titres en list
        pass

    # Approche correcte : on pousse une PageBreak entre les slides, et on calcule
    # le titre via un dictionnaire indexe par numero de page.
    titles_by_page = {i: (sd.get("title", ""), sd.get("subtitle") or "")
                      for i, sd in enumerate(SLIDES, 1)}

    def header_footer_dyn(canvas, doc):
        canvas.saveState()
        title, subtitle = titles_by_page.get(doc.page, ("", ""))
        # Pour la couverture on n'affiche pas de bandeau (page 1)
        if doc.page != 1 and title:
            canvas.setFillColor(PRIMARY_RL)
            canvas.rect(0, A4[0] - 2.5 * cm, A4[1], 2.5 * cm, fill=1, stroke=0)
            canvas.setFillColor(white)
            canvas.setFont("Helvetica-Bold", 22)
            canvas.drawString(1.5 * cm, A4[0] - 1.6 * cm, title)
        # Footer
        canvas.setFillColor(grey)
        canvas.setFont("Helvetica", 9)
        canvas.drawCentredString(A4[1] / 2, 0.8 * cm,
                                 f"E-Store  |  A. Belmoussa & N. Ben Soumane  |  Pr. O. Zahour  |  {doc.page}/{total}")
        canvas.restoreState()

    flow = []
    for i, sd in enumerate(SLIDES, 1):
        t = sd["type"]
        if t == "cover":
            render_cover(sd, i)
        elif t == "content":
            render_content(sd)
        elif t == "table":
            render_table(sd)
        elif t == "code":
            render_code(sd)
        elif t == "architecture":
            render_arch(sd)
        elif t == "security":
            render_security(sd)
        elif t == "two_columns":
            render_two_columns(sd)
        elif t == "demo":
            render_demo(sd)
        if i < total:
            flow.append(PageBreak())

    doc.build(flow, onFirstPage=header_footer_dyn, onLaterPages=header_footer_dyn)
    return out


# =====================================================================
# GUIDE DE DISCOURS POUR LES PRESENTATEURS (PDF)
# =====================================================================
def build_speaker_guide_pdf():
    """Guide complet de soutenance avec script discours pour Akram + Nouhaila."""
    out = os.path.join(OUT_DIR, "guide-soutenance-discours.pdf")

    base = getSampleStyleSheet()
    H1 = ParagraphStyle("H1", parent=base["Heading1"], fontName="Helvetica-Bold",
                        fontSize=20, textColor=PRIMARY_RL, spaceAfter=12, leading=24)
    H2 = ParagraphStyle("H2", parent=base["Heading2"], fontName="Helvetica-Bold",
                        fontSize=14, textColor=DARK_RL, spaceBefore=12, spaceAfter=6, leading=18)
    Akram = ParagraphStyle("Akram", parent=base["BodyText"], fontName="Helvetica",
                           fontSize=10.5, textColor=DARK_RL, leading=15,
                           alignment=TA_JUSTIFY, leftIndent=0.5 * cm, spaceAfter=4)
    Nouhaila = ParagraphStyle("Nouhaila", parent=base["BodyText"], fontName="Helvetica",
                              fontSize=10.5, textColor=DARK_RL, leading=15,
                              alignment=TA_JUSTIFY, leftIndent=0.5 * cm, spaceAfter=4)
    Body = ParagraphStyle("Body", parent=base["BodyText"], fontName="Helvetica",
                          fontSize=10.5, textColor=DARK_RL, leading=15,
                          alignment=TA_JUSTIFY, spaceAfter=6)
    Quote = ParagraphStyle("Quote", parent=Body, fontName="Helvetica-Oblique",
                           leftIndent=1 * cm, rightIndent=1 * cm,
                           textColor=DARK_RL, spaceAfter=6,
                           backColor=LIGHT_RL, borderPadding=8)
    Bullet = ParagraphStyle("Bullet", parent=Body, leftIndent=0.6 * cm, bulletIndent=0.2 * cm)
    Tip = ParagraphStyle("Tip", parent=Body, fontName="Helvetica-Oblique",
                         textColor=ACCENT_RL, leftIndent=0.5 * cm, spaceAfter=6)

    doc = SimpleDocTemplate(out, pagesize=A4,
                            leftMargin=2 * cm, rightMargin=2 * cm,
                            topMargin=2 * cm, bottomMargin=1.8 * cm,
                            title="Guide de soutenance E-Store",
                            author="A. Belmoussa & N. Ben Soumane")

    flow = []

    # Couverture
    flow.append(Spacer(1, 5 * cm))
    flow.append(Paragraph("Guide de soutenance", H1))
    flow.append(Spacer(1, 0.4 * cm))
    flow.append(Paragraph("Script de presentation pour Akram & Nouhaila",
                          ParagraphStyle("c", parent=Body, fontSize=16,
                                         textColor=DARK_RL, leading=20)))
    flow.append(Spacer(1, 1 * cm))
    flow.append(Paragraph("Duree totale : 10 minutes  |  12 slides",
                          ParagraphStyle("c", parent=Body, fontSize=12,
                                         textColor=grey, leading=18)))
    flow.append(PageBreak())

    # ─── Section 1 : Comment lire ce guide ─────────────────────────────
    flow.append(Paragraph("Comment utiliser ce guide", H1))
    flow.append(Paragraph(
        "Ce document contient le <b>script complet</b> que vous devez prononcer "
        "pendant la soutenance. Il est decoupe slide par slide, avec :", Body))
    flow.append(Paragraph("- <b>Qui parle</b> (Akram ou Nouhaila) — alternance equilibree.", Bullet))
    flow.append(Paragraph("- <b>Le discours mot pour mot</b> a dire (en italique).", Bullet))
    flow.append(Paragraph("- <b>La duree</b> de chaque slide.", Bullet))
    flow.append(Paragraph("- <b>Conseils visuels</b> (ce que vous montrez a l'ecran).", Bullet))
    flow.append(Paragraph("- <b>Transitions</b> entre presentateurs.", Bullet))

    flow.append(Paragraph("Repartition globale", H2))
    repart = Table([
        ["Akram", "Slides 1, 3, 5, 7, 9, 12 (couverture, contexte, stack, securite, frontend, conclusion)", "5 min"],
        ["Nouhaila", "Slides 2, 4, 6, 8, 10, 11 (plan, architecture, code, mongo, demo, tests/docker)", "5 min"],
    ], colWidths=[2.5 * cm, 12 * cm, 1.5 * cm])
    repart.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (0, -1), PRIMARY_RL),
        ("TEXTCOLOR", (0, 0), (0, -1), white),
        ("FONTNAME", (0, 0), (0, -1), "Helvetica-Bold"),
        ("FONTSIZE", (0, 0), (-1, -1), 10),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 8),
        ("TOPPADDING", (0, 0), (-1, -1), 8),
        ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
        ("GRID", (0, 0), (-1, -1), 0.4, grey),
        ("ALIGN", (-1, 0), (-1, -1), "CENTER"),
    ]))
    flow.append(repart)

    flow.append(Paragraph("Conseils generaux avant de commencer", H2))
    for tip in [
        "Arriver 10 minutes en avance, tester la connexion HDMI/projecteur.",
        "Ouvrir le navigateur sur localhost:4200 et le terminal avec Spring Boot deja lance.",
        "Avoir le PDF des slides ouvert en plein ecran (touche F5).",
        "Respirer profondement avant de commencer — ralentir le debit, articuler.",
        "Regarder le jury (Pr. Zahour) et non pas l'ecran.",
        "Ne JAMAIS dire 'euh' ou 'voila' a la fin d'une phrase — silence vaut mieux.",
        "Si on bug : ne pas paniquer, dire 'Excusez-moi, je reprends' et continuer.",
    ]:
        flow.append(Paragraph("- " + tip, Bullet))

    flow.append(PageBreak())

    # ─── Section 2 : Discours slide par slide ──────────────────────────
    discours = [
        # Slide 1
        {
            "num": 1, "speaker": "Akram", "duration": "30 s",
            "title": "Couverture",
            "visuel": "Slide de couverture affichee.",
            "text": [
                ("Akram", "Bonjour Monsieur Zahour, bonjour a tout le monde. Je suis Akram Belmoussa, "
                          "et avec ma binome Nouhaila Ben Soumane, nous avons le plaisir de vous "
                          "presenter aujourd'hui notre projet de fin de module Full-Stack."),
                ("Akram", "Notre projet s'intitule <b>E-Store</b>. C'est une application e-commerce complete "
                          "developpee dans le cadre du module Full-Stack, sous votre direction, "
                          "Pr. Zahour. La presentation va durer environ 10 minutes."),
            ],
            "tip": "Ne PAS lire ces lignes — apprendre par coeur l'idee, pas les mots exacts.",
        },
        # Slide 2
        {
            "num": 2, "speaker": "Nouhaila", "duration": "30 s",
            "title": "Plan",
            "visuel": "Plan de la presentation (10 points).",
            "text": [
                ("Nouhaila", "Voici le plan de notre presentation. Nous allons d'abord poser le contexte "
                             "et les objectifs, puis presenter notre architecture en 3 couches et "
                             "5 domaines fonctionnels. Nous detaillerons ensuite la stack technique, "
                             "puis nous plongerons dans le code du backend, la securite via JWT, "
                             "et notre choix de MongoDB pour les avis."),
                ("Nouhaila", "Nous enchainerons avec le frontend Angular, une demonstration en direct "
                             "de l'application, puis nous parlerons des tests, du DevOps et d'une "
                             "difficulte rencontree avec Docker. Nous conclurons sur les perspectives."),
            ],
            "tip": "Pointer chaque numero a l'ecran pendant que vous parlez.",
            "transition": "Akram : 'Akram, a toi pour le contexte.'",
        },
        # Slide 3
        {
            "num": 3, "speaker": "Akram", "duration": "1 min",
            "title": "Contexte et objectifs",
            "visuel": "6 bullets des objectifs.",
            "text": [
                ("Akram", "Merci Nouhaila. E-Store est une application e-commerce complete : un utilisateur "
                          "peut s'inscrire, se connecter, parcourir un catalogue de produits, ajouter "
                          "des articles a son panier, valider une commande et consulter son historique."),
                ("Akram", "Au-dela de la fonctionnalite, nous nous sommes fixe cinq objectifs pedagogiques. "
                          "Premierement : <b>maitriser une architecture full-stack moderne</b> avec Spring Boot "
                          "cote serveur et Angular cote client. Deuxiemement : illustrer la "
                          "<b>persistance hybride</b> — c'est-a-dire utiliser une base de donnees "
                          "relationnelle MySQL pour les donnees structurees, ET une base documentaire "
                          "MongoDB pour les avis."),
                ("Akram", "Troisiemement : appliquer les <b>bonnes pratiques de developpement</b> — DTOs, "
                          "transactions, tests unitaires. Quatriemement : <b>securiser l'application</b> "
                          "par JWT et Spring Security 6, qui sont les standards actuels du marche. "
                          "Et enfin cinquiemement : produire une <b>interface utilisateur moderne</b> "
                          "avec Angular 17 en mode standalone."),
            ],
            "tip": "Insister sur le mot 'hybride' — c'est ce qui differencie le projet.",
        },
        # Slide 4
        {
            "num": 4, "speaker": "Nouhaila", "duration": "1 min 30",
            "title": "Architecture 3 couches x 5 domaines",
            "visuel": "Schema avec 3 couches a gauche, 6 boites domaines a droite.",
            "text": [
                ("Nouhaila", "Pour atteindre ces objectifs, nous avons retenu une architecture en deux dimensions."),
                ("Nouhaila", "<b>Premiere dimension : trois couches techniques.</b> La couche <b>presentation</b>, "
                             "implementee par Angular, est responsable de l'interface utilisateur, des "
                             "interactions et de la validation cote client. La couche <b>logique metier</b>, "
                             "en Spring Boot, contient les controllers REST, les services qui appliquent "
                             "les regles metier, et les DTOs qui evitent d'exposer directement les entites. "
                             "Enfin, la couche <b>acces aux donnees</b> utilise Spring Data JPA pour le "
                             "relationnel et Spring Data MongoDB pour le documentaire."),
                ("Nouhaila", "<b>Deuxieme dimension : six domaines fonctionnels.</b> Chaque domaine est un "
                             "sous-package autonome du backend, avec ses entites, repositories, services "
                             "et controllers. Nous avons : <b>customer</b> pour les utilisateurs et "
                             "l'authentification ; <b>catalog</b> pour les categories et produits ; "
                             "<b>inventory</b> pour le stock ; <b>shopping</b> pour le panier ; "
                             "<b>billing</b> pour les commandes ; et <b>review</b> pour les avis stockes en MongoDB."),
                ("Nouhaila", "Cette double organisation suit les principes du Domain-Driven Design : "
                             "elle permet de scaler facilement, d'isoler les responsabilites, et de "
                             "garantir une lecture rapide du code."),
            ],
            "tip": "Pointer du doigt chaque couche puis chaque domaine pendant que vous les nommez.",
        },
        # Slide 5
        {
            "num": 5, "speaker": "Akram", "duration": "30 s",
            "title": "Stack technique",
            "visuel": "Tableau Couche -> Technologies.",
            "text": [
                ("Akram", "Voici concretement la stack technique. Nous utilisons Angular 17 avec Bootstrap 5 "
                          "pour le frontend ; Spring Boot 3.3 et Spring Security 6 pour le backend ; "
                          "MySQL 8 en production avec H2 in-memory pour le profil developpement ; "
                          "et MongoDB 7 pour les avis."),
                ("Akram", "Cote outillage : Maven 3.9 pour le build Java, npm pour le frontend, "
                          "JUnit 5 et Mockito pour les tests, et Docker Compose qui orchestre les bases "
                          "de donnees ainsi que phpMyAdmin et mongo-express pour leur inspection."),
            ],
            "tip": "Tu peux survoler rapidement — le tableau parle de lui-meme.",
        },
        # Slide 6
        {
            "num": 6, "speaker": "Nouhaila", "duration": "1 min 30",
            "title": "Backend — checkout transactionnel",
            "visuel": "Code Java de OrderService.checkout().",
            "text": [
                ("Nouhaila", "Pour illustrer la richesse du backend, je vais detailler notre operation la plus "
                             "critique : la validation d'une commande, dans la methode "
                             "<b>OrderService.checkout()</b>."),
                ("Nouhaila", "Vous voyez en haut l'annotation <b>@Transactional</b>. Elle est essentielle : "
                             "elle garantit que toutes les operations qui suivent forment une unite "
                             "atomique — soit elles reussissent toutes, soit aucune n'est appliquee."),
                ("Nouhaila", "L'algorithme se deroule en trois etapes. <b>Etape 1</b> : nous parcourons le "
                             "panier et nous verifions, pour chaque article, que le stock est suffisant. "
                             "Si l'un d'eux est en rupture, on leve une <b>BusinessException</b> qui "
                             "se traduira par un HTTP 409 Conflict cote API."),
                ("Nouhaila", "<b>Etape 2</b> : nous creons l'objet <b>Order</b> avec ses <b>OrderItem</b>, "
                             "et nous decrementons le stock pour chaque produit. <b>Etape 3</b> : nous "
                             "passons la commande en statut <b>CONFIRMED</b>, nous la sauvegardons, et "
                             "nous vidons le panier."),
                ("Nouhaila", "Le point cle : si une exception survient pendant les etapes 2 ou 3 — par exemple "
                             "une coupure base de donnees — toutes les modifications precedentes sont "
                             "annulees. C'est l'<b>atomicite ACID</b> en action, garantie par Spring."),
            ],
            "tip": "Pointer @Transactional puis les 3 commentaires // 1) // 2) // 3) au fur et a mesure.",
        },
        # Slide 7
        {
            "num": 7, "speaker": "Akram", "duration": "1 min",
            "title": "Securite — JWT + Spring Security 6",
            "visuel": "6 etapes numerotees du flux JWT.",
            "text": [
                ("Akram", "Concernant la securite, nous avons implemente un systeme stateless base sur "
                          "JSON Web Tokens, avec Spring Security 6."),
                ("Akram", "Le flux que vous voyez ici se deroule en six etapes. Quand un client se connecte, "
                          "il envoie email et mot de passe a l'endpoint /api/auth/login. Le mot de passe "
                          "est verifie via <b>BCrypt</b> — jamais stocke en clair. Si la verification "
                          "reussit, notre <b>JwtService</b> genere un token signe en HMAC-SHA256, "
                          "contenant l'email, l'identifiant et le role de l'utilisateur. Ce token "
                          "expire au bout de 24 heures."),
                ("Akram", "Le client stocke ce token cote navigateur dans le localStorage. A chaque requete "
                          "suivante, notre <b>AuthInterceptor</b> Angular l'ajoute automatiquement dans "
                          "l'entete Authorization Bearer. Cote serveur, notre <b>JwtAuthenticationFilter</b> "
                          "intercepte chaque appel, valide la signature et l'expiration, puis place "
                          "l'utilisateur dans le SecurityContext de Spring."),
                ("Akram", "Resultat : nos endpoints d'administration — comme la creation de produits — "
                          "sont proteges par <code>@PreAuthorize hasRole ADMIN</code>, et les endpoints "
                          "de lecture publique restent accessibles sans authentification."),
            ],
            "tip": "Insister sur 'stateless' — gros avantage par rapport aux sessions HTTP classiques.",
        },
        # Slide 8
        {
            "num": 8, "speaker": "Nouhaila", "duration": "30 s",
            "title": "MongoDB pour les avis",
            "visuel": "Bullets + structure du document JSON.",
            "text": [
                ("Nouhaila", "Pourquoi MongoDB pour les avis et pas une simple table SQL ? Parce que les "
                             "avis sont des donnees <b>semi-structurees</b> : peu de relations, "
                             "lecture frequente, ecriture en flux. C'est exactement le cas d'usage "
                             "ou MongoDB excelle, et cela permet de demontrer la coexistence des deux "
                             "paradigmes dans une meme application."),
                ("Nouhaila", "Concretement, nous utilisons Spring Data MongoDB qui offre la meme abstraction "
                             "que JPA : un repository <b>ReviewRepository</b> avec des methodes derivees "
                             "comme <b>findByProductIdOrderByCreatedAtDesc</b>. Aucune SQL ni JPQL a ecrire."),
            ],
            "tip": "Si Mongo n'est pas lance pendant la demo, expliquer que l'app continue grace au seeder defensif.",
        },
        # Slide 9
        {
            "num": 9, "speaker": "Akram", "duration": "1 min",
            "title": "Frontend Angular 17",
            "visuel": "Bullets sur les choix Angular.",
            "text": [
                ("Akram", "Cote frontend, nous avons fait le choix d'Angular 17 en mode <b>standalone</b>. "
                          "C'est la version la plus moderne du framework : pas de NgModule a maintenir, "
                          "moins de boilerplate, lazy loading natif."),
                ("Akram", "Notre application charge un bundle initial leger de 13 kilo-octets de "
                          "JavaScript metier, plus 245 kilo-octets de framework partage. Chaque "
                          "fonctionnalite — login, catalogue, panier, commandes — est un bundle "
                          "independant charge a la demande."),
                ("Akram", "Pour la reactivite de l'UI, nous utilisons les <b>signals</b>, une nouveaute "
                          "d'Angular 17. Par exemple, le badge du panier dans le header est lie au "
                          "signal <code>cart.itemCount()</code>. Quand on ajoute un produit, le badge "
                          "se met a jour automatiquement, sans aucune ligne de code dans le composant Header."),
                ("Akram", "Notre <b>AuthService</b> persiste l'utilisateur en localStorage et expose un "
                          "BehaviorSubject. Notre <b>AuthInterceptor</b> ajoute le JWT a chaque requete. "
                          "Notre <b>ErrorInterceptor</b> centralise la gestion des erreurs HTTP — "
                          "401 declenche la deconnexion automatique, 4xx affiche un toast d'erreur. "
                          "Aucun composant n'a a gerer les erreurs manuellement."),
            ],
            "tip": "Si tu connais bien le code, tu peux mentionner que les composants sont <50 lignes.",
        },
        # Slide 10
        {
            "num": 10, "speaker": "Nouhaila", "duration": "1 min",
            "title": "Demonstration en direct",
            "visuel": "Basculer sur le navigateur localhost:4200.",
            "text": [
                ("Nouhaila", "Place a la demonstration. Je passe au navigateur."),
                ("Nouhaila", "[Action 1] Voici notre catalogue avec les 12 produits seeded. "
                             "Je peux filtrer par categorie — choisissons Sport — et je peux faire une "
                             "recherche textuelle."),
                ("Nouhaila", "[Action 2] Je clique sur un produit, on arrive sur sa fiche detaillee, "
                             "avec image, description, prix, stock disponible et la section avis."),
                ("Nouhaila", "[Action 3] Si je tente d'ajouter au panier sans etre connectee, je suis "
                             "redirigee vers /login — c'est notre AuthGuard."),
                ("Nouhaila", "[Action 4] Je me connecte avec user@estore.ma. Le header affiche maintenant "
                             "mon prenom et le badge panier."),
                ("Nouhaila", "[Action 5] J'ajoute deux produits au panier — vous voyez le badge passer a 2."),
                ("Nouhaila", "[Action 6] Je vais sur la page panier, je modifie une quantite, le total "
                             "se recalcule automatiquement."),
                ("Nouhaila", "[Action 7] Je valide la commande. Je suis redirigee vers /orders, et je "
                             "vois la commande #1 confirmee avec son detail expandable."),
                ("Nouhaila", "[Action 8] Je retourne sur le produit et je depose un avis 5 etoiles. "
                             "L'avis apparait immediatement — il a ete enregistre dans MongoDB."),
            ],
            "tip": "ENTRAINEZ-VOUS — la demo est le moment le plus risque. Avoir un plan B si quelque chose plante.",
        },
        # Slide 11
        {
            "num": 11, "speaker": "Nouhaila", "duration": "30 s",
            "title": "Tests, DevOps et Docker",
            "visuel": "Deux colonnes : Tests / Docker.",
            "text": [
                ("Nouhaila", "Sur le plan DevOps : nous avons ecrit 9 tests unitaires JUnit 5 + Mockito, "
                             "qui passent tous en BUILD SUCCESS. Trois suites — ProductService, CartService, "
                             "OrderService — couvrent les services metier critiques."),
                ("Nouhaila", "Concernant Docker : nous avons rencontre une difficulte d'installation "
                             "sur Windows 11 avec le message <i>'C:/ProgramData/DockerDesktop must be "
                             "owned by an elevated account'</i>. La cause etait un dossier residuel "
                             "d'une installation precedente. Nous avons applique deux solutions : "
                             "d'abord la suppression du dossier en mode admin et la reinstallation, "
                             "puis nous avons mis en place un <b>profil dev avec H2 in-memory</b> qui "
                             "permet de demarrer l'application sans aucune installation prealable."),
                ("Nouhaila", "Cette dualite Docker / H2 a ete pensee des le depart — c'est une bonne "
                             "pratique professionnelle qui nous a permis de poursuivre la demonstration "
                             "malgre l'incident."),
            ],
            "tip": "Le jury peut poser une question sur Docker — vous avez la reponse complete.",
        },
        # Slide 12
        {
            "num": 12, "speaker": "Akram", "duration": "30 s",
            "title": "Conclusion",
            "visuel": "Bullets bilan + perspectives + 'Merci'.",
            "text": [
                ("Akram", "Pour conclure : ce projet nous a permis de mettre en pratique les enseignements "
                          "du module dans une application complete. L'architecture en 3 couches et "
                          "5 domaines a ete respectee, la securite JWT est operationnelle, la "
                          "persistance hybride SQL plus NoSQL est fonctionnelle, et les 9 tests "
                          "unitaires garantissent la stabilite du code."),
                ("Akram", "Comme perspectives, nous envisageons d'integrer un module de paiement en ligne "
                          "via Stripe, des notifications par email lors de la confirmation des "
                          "commandes, un systeme de recommandation produit base sur l'historique, "
                          "et un deploiement cloud sur Heroku ou AWS."),
                ("Akram", "Nous vous remercions pour votre attention. Nous sommes prets a repondre a "
                          "vos questions."),
            ],
            "tip": "Sourire, regarder le jury, tendre legerement la main pour inviter les questions.",
        },
    ]

    for d in discours:
        speaker_color = PRIMARY_HEX if d["speaker"] == "Akram" else "#198754"
        flow.append(Paragraph(f"Slide {d['num']} — {d['title']}", H1))
        flow.append(Paragraph(
            f"<b>Locuteur :</b> <font color='{speaker_color}'>{d['speaker']}</font>  "
            f"<b>|  Duree :</b> {d['duration']}", H2))
        flow.append(Paragraph(f"<i>Visuel a l'ecran :</i> {d['visuel']}", Body))

        for who, text in d["text"]:
            who_color = PRIMARY_HEX if who == "Akram" else "#198754"
            flow.append(Paragraph(
                f"<b><font color='{who_color}'>{who} :</font></b>", Body))
            flow.append(Paragraph(f"<i>« {text} »</i>", Quote))

        if d.get("tip"):
            flow.append(Paragraph("Conseil : " + d["tip"], Tip))
        if d.get("transition"):
            flow.append(Paragraph("Transition : " + d["transition"], Tip))
        flow.append(PageBreak())

    # ─── Section 3 : Q & R prevues ────────────────────────────────────
    flow.append(Paragraph("Annexe : questions probables du jury", H1))
    flow.append(Paragraph(
        "Voici une liste de questions courantes avec des reponses preparees. Lire et "
        "memoriser au moins les 3-4 premieres — la reponse doit etre fluide.", Body))

    qa = [
        ("Pourquoi MongoDB pour les avis et pas du SQL ?",
         "Akram",
         "Les avis sont des donnees semi-structurees, peu contraintes relationnellement, "
         "avec beaucoup de lectures et d'ecritures concurrentes. MongoDB est plus performant "
         "pour ce cas et illustre la coexistence de deux paradigmes de persistance."),

        ("Comment garantissez-vous la coherence du panier en cas de stock insuffisant ?",
         "Nouhaila",
         "Le checkout est annote @Transactional. Toutes les verifications de stock sont "
         "faites AVANT toute modification. Si une exception est levee, le rollback automatique "
         "annule toutes les modifications precedentes. C'est l'atomicite ACID."),

        ("Pourquoi H2 en plus de MySQL ?",
         "Akram",
         "H2 est utilisee pour le profil dev — demarrage instantane sans installation, "
         "ideale pour les demos et le developpement local. MySQL est utilisee pour la "
         "production. Le code est strictement identique grace a JPA."),

        ("Comment le JWT est-il securise ?",
         "Akram",
         "Signe en HMAC-SHA256 avec une cle secrete configuree dans application.properties. "
         "Expiration 24h. Verifie a chaque requete par le JwtAuthenticationFilter avant "
         "execution du controller."),

        ("Pourquoi Angular standalone et pas un NgModule classique ?",
         "Akram",
         "C'est le standard depuis Angular 17. Plus concis, plus performant grace au "
         "lazy loading natif, pas de boilerplate. Les composants declarent leurs "
         "dependances directement."),

        ("Que se passe-t-il si l'utilisateur tente de commander un panier vide ?",
         "Nouhaila",
         "OrderService leve une BusinessException avec le message 'Votre panier est vide', "
         "qui est traduite par le GlobalExceptionHandler en HTTP 409 Conflict."),

        ("Le mot de passe est-il chiffre ?",
         "Akram",
         "Il est hashe avec BCrypt — un algorithme one-way avec sel aleatoire et facteur "
         "de cout. Impossible de le retrouver meme en cas de fuite de la base."),

        ("Comment avez-vous gere l'erreur Docker ?",
         "Nouhaila",
         "Le dossier C:/ProgramData/DockerDesktop avait de mauvaises permissions, residu "
         "d'une installation precedente. Nous l'avons supprime en mode administrateur, "
         "et nous avons mis en place un profil dev H2 qui evite la dependance a Docker "
         "pour la soutenance."),

        ("Combien de tests avez-vous ecrit et que couvrent-ils ?",
         "Nouhaila",
         "9 tests unitaires repartis sur 3 suites : ProductServiceTest (4), CartServiceTest (2), "
         "OrderServiceTest (3). Ils couvrent la recherche, la gestion du stock, et le checkout "
         "transactionnel — les 3 services metier critiques."),

        ("Combien de temps avez-vous mis a developper le projet ?",
         "Akram",
         "Environ X semaines, en partageant le travail entre nous deux : Akram principalement "
         "sur le backend et la securite, Nouhaila sur le frontend et la modelisation des donnees, "
         "et nous avons fait l'integration et les tests ensemble."),
    ]

    for q, who, a in qa:
        who_color = PRIMARY_HEX if who == "Akram" else "#198754"
        flow.append(Paragraph(f"<b>Q : {q}</b>", H2))
        flow.append(Paragraph(f"<b><font color='{who_color}'>Reponse ({who}) :</font></b>", Body))
        flow.append(Paragraph(a, Quote))

    flow.append(PageBreak())

    # ─── Section 4 : Checklist du jour J ──────────────────────────────
    flow.append(Paragraph("Checklist du jour J", H1))
    flow.append(Paragraph("Avant la soutenance (1h avant)", H2))
    for item in [
        "[ ] Backend Spring Boot lance et accessible sur localhost:8080",
        "[ ] Frontend Angular lance et accessible sur localhost:4200",
        "[ ] Tester un parcours complet : login -> ajout panier -> commande",
        "[ ] PDF des slides ouvert en plein ecran (mode presentation)",
        "[ ] Guide de discours imprime ou sur le telephone",
        "[ ] Cable HDMI / adaptateur si demo sur projecteur",
        "[ ] Souris (plus pratique que le trackpad pour la demo)",
        "[ ] Bouteille d'eau pour les deux presentateurs",
        "[ ] Telephone en mode silencieux",
    ]:
        flow.append(Paragraph(item, Bullet))

    flow.append(Paragraph("Pendant la presentation", H2))
    for item in [
        "[ ] Saluer le jury avant de commencer",
        "[ ] Annoncer son nom + nom du binome + nom du projet",
        "[ ] Garder un debit calme — ne pas accelerer si stress",
        "[ ] Regarder le jury, pas l'ecran",
        "[ ] Pointer ce dont on parle a l'ecran",
        "[ ] Faire des transitions visibles entre presentateurs",
        "[ ] Sourire, etre positif",
        "[ ] En cas de bug demo : ne pas paniquer, expliquer ce qui devait se passer",
    ]:
        flow.append(Paragraph(item, Bullet))

    flow.append(Paragraph("Apres la presentation", H2))
    for item in [
        "[ ] Remercier le jury ('Merci pour votre attention')",
        "[ ] Rester calme et a l'ecoute pendant les questions",
        "[ ] Si on ne sait pas : 'C'est une bonne question, nous n'avons pas explore ce point'",
        "[ ] Ne pas couper la parole au jury",
        "[ ] Remercier a la fin meme si la note semble decevante",
    ]:
        flow.append(Paragraph(item, Bullet))

    # Page footer for all pages
    def hf(canvas, doc):
        canvas.saveState()
        canvas.setStrokeColor(PRIMARY_RL)
        canvas.setLineWidth(0.6)
        canvas.line(2 * cm, A4[1] - 1.6 * cm, A4[0] - 2 * cm, A4[1] - 1.6 * cm)
        canvas.setFillColor(PRIMARY_RL)
        canvas.setFont("Helvetica-Bold", 9)
        canvas.drawString(2 * cm, A4[1] - 1.4 * cm, "GUIDE DE SOUTENANCE")
        canvas.drawRightString(A4[0] - 2 * cm, A4[1] - 1.4 * cm, "E-STORE")
        canvas.setFont("Helvetica", 8)
        canvas.setFillColor(grey)
        canvas.drawString(2 * cm, 1.2 * cm, "A. Belmoussa & N. Ben Soumane — Pr. O. Zahour")
        canvas.drawRightString(A4[0] - 2 * cm, 1.2 * cm, f"Page {doc.page}")
        canvas.restoreState()

    doc.build(flow, onFirstPage=lambda c, d: None, onLaterPages=hf)
    return out


# =====================================================================
# Main
# =====================================================================
if __name__ == "__main__":
    print("=" * 60)
    print("Generation de la presentation E-Store")
    print("=" * 60)
    pptx = build_pptx()
    print(f"  [OK] {os.path.basename(pptx):50s} ({os.path.getsize(pptx)//1024} ko)")
    pdf = build_pdf()
    print(f"  [OK] {os.path.basename(pdf):50s} ({os.path.getsize(pdf)//1024} ko)")
    guide = build_speaker_guide_pdf()
    print(f"  [OK] {os.path.basename(guide):50s} ({os.path.getsize(guide)//1024} ko)")
    print("=" * 60)
    print(f"3 fichiers generes dans : {OUT_DIR}")

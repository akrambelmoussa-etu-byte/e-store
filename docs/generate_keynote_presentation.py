# -*- coding: utf-8 -*-
"""
E-STORE | KEYNOTE EDITION
=========================
Genere une presentation .pptx premium style Apple Keynote / Tesla AI Day
avec transitions Morph fonctionnelles et 38 slides chainees.

USAGE :
    python generate_keynote_presentation.py
    -> docs/E-Store-Keynote.pptx
"""
import os
import sys
import io
import math

if sys.stdout.encoding and sys.stdout.encoding.lower() != "utf-8":
    try:
        sys.stdout.reconfigure(encoding="utf-8", errors="replace")
    except Exception:
        sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn, nsmap
from lxml import etree

OUT_DIR = os.path.dirname(os.path.abspath(__file__))

# =====================================================================
#  PALETTE PREMIUM — STYLE TESLA / OPENAI / NVIDIA GTC
# =====================================================================
DEEP        = RGBColor(0x05, 0x08, 0x1E)   # Deep space - fond principal
NIGHT       = RGBColor(0x0B, 0x0F, 0x2C)
INK         = RGBColor(0x14, 0x1A, 0x42)
CARD        = RGBColor(0x1A, 0x21, 0x4D)   # Glass card
SUBTLE      = RGBColor(0x29, 0x32, 0x5E)
LINE        = RGBColor(0x3A, 0x46, 0x78)

# Accents neon
CYAN        = RGBColor(0x00, 0xE5, 0xFF)
VIOLET      = RGBColor(0xB5, 0x7A, 0xFF)
PINK        = RGBColor(0xFF, 0x4F, 0x8B)
LIME        = RGBColor(0x00, 0xFF, 0xB2)
AMBER       = RGBColor(0xFF, 0xC1, 0x07)
CORAL       = RGBColor(0xFF, 0x6F, 0x61)

# Texte
WHITE       = RGBColor(0xF5, 0xF7, 0xFF)
SOFT        = RGBColor(0xC8, 0xCD, 0xE5)
DIM         = RGBColor(0x8B, 0x92, 0xAB)
FAINT       = RGBColor(0x52, 0x5C, 0x80)


# =====================================================================
#  UTILITAIRE MORPH — INJECTION XML
# =====================================================================
def add_morph_transition(slide, option="byObject", speed="med"):
    """
    Ajoute une transition Morph (PowerPoint 2016+) sur la slide.
    `option` : 'byObject' | 'byWord' | 'byChar'
    `speed`  : 'slow' | 'med' | 'fast'
    """
    sld = slide._element
    # Supprimer toute transition existante
    for t in sld.findall(qn('p:transition')):
        sld.remove(t)
    # Construire l'element de transition Morph
    P14 = "http://schemas.microsoft.com/office/powerpoint/2010/main"
    P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
    nsmap_full = {None: P_NS, "p14": P14, "mc": "http://schemas.openxmlformats.org/markup-compatibility/2006"}

    transition = etree.SubElement(sld, f"{{{P_NS}}}transition", attrib={
        "spd": speed,
        "advClick": "1",
    }, nsmap={"p14": P14, "mc": "http://schemas.openxmlformats.org/markup-compatibility/2006"})
    transition.set("{http://schemas.openxmlformats.org/markup-compatibility/2006}Ignorable", "p14")
    morph = etree.SubElement(transition, f"{{{P14}}}morph")
    morph.set("option", option)
    return transition


def add_fade_transition(slide, speed="med"):
    """Transition fondu standard."""
    sld = slide._element
    for t in sld.findall(qn('p:transition')):
        sld.remove(t)
    P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
    transition = etree.SubElement(sld, f"{{{P_NS}}}transition", attrib={
        "spd": speed, "advClick": "1"
    })
    etree.SubElement(transition, f"{{{P_NS}}}fade")


def set_speaker_notes(slide, text):
    """Ajoute des notes orateur."""
    notes = slide.notes_slide
    tf = notes.notes_text_frame
    tf.text = text
    for p in tf.paragraphs:
        for r in p.runs:
            r.font.size = Pt(11)


# =====================================================================
#  HELPERS DESIGN
# =====================================================================
def set_bg(slide, color=DEEP):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def named_shape(slide, shape_type, x, y, w, h, name):
    """Cree une forme nommee (essentiel pour Morph)."""
    s = slide.shapes.add_shape(shape_type, x, y, w, h)
    s.name = name
    return s


def add_text(slide, x, y, w, h, text, *, name=None, size=14, bold=False,
             italic=False, color=WHITE, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, font="Calibri", letter_spacing=None):
    tx = slide.shapes.add_textbox(x, y, w, h)
    if name:
        tx.name = name
    tf = tx.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(0); tf.margin_right = Pt(0)
    tf.margin_top = Pt(0); tf.margin_bottom = Pt(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.text = text
    for r in p.runs:
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
        if letter_spacing:
            # Espacement de caracteres via XML
            rPr = r._r.get_or_add_rPr()
            rPr.set("spc", str(letter_spacing))
    return tx


def gradient_band(slide, x, y, w, h, color, name=None):
    """Bande horizontale colore (effet halo)."""
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if name: s.name = name
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def glass_card(slide, x, y, w, h, *, name=None, accent=None, accent_pos="top"):
    """Carte glassmorphism : fond translucide + bordure subtile."""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    if name: card.name = name
    card.adjustments[0] = 0.05
    card.fill.solid(); card.fill.fore_color.rgb = CARD
    card.line.color.rgb = LINE; card.line.width = Pt(0.5)

    if accent:
        if accent_pos == "top":
            bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                          x + Inches(0.15), y,
                                          w - Inches(0.3), Inches(0.04))
        else:  # left
            bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                          x, y + Inches(0.15),
                                          Inches(0.05), h - Inches(0.3))
        bar.fill.solid(); bar.fill.fore_color.rgb = accent
        bar.line.fill.background()
    return card


def neon_circle(slide, cx, cy, r, color, name=None, glow=True):
    """Cercle néon avec effet de halo (un cercle plus large transparent)."""
    if glow:
        halo = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                       cx - r * 1.6, cy - r * 1.6,
                                       r * 3.2, r * 3.2)
        halo.fill.solid(); halo.fill.fore_color.rgb = color
        halo.line.fill.background()
        # transparence
        sf = halo.fill._xPr
        # Approche : ajouter opacity via XML
        spPr = halo.fill._xPr
        # Plus simple : on laisse opaque mais en couleur sombre
        halo.fill.fore_color.rgb = NIGHT  # fallback safe
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                cx - r, cy - r, r * 2, r * 2)
    if name: c.name = name
    c.fill.solid(); c.fill.fore_color.rgb = color
    c.line.color.rgb = WHITE; c.line.width = Pt(0.5)
    return c


def line_connector(slide, x1, y1, x2, y2, color=LINE, width=1.0, name=None):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    if name: conn.name = name
    conn.line.color.rgb = color
    conn.line.width = Pt(width)
    return conn


def arrow_connector(slide, x1, y1, x2, y2, color=CYAN, width=1.5, name=None):
    conn = line_connector(slide, x1, y1, x2, y2, color, width, name)
    line = conn.line._get_or_add_ln()
    head = etree.SubElement(line, qn('a:tailEnd'))
    head.set('type', 'triangle'); head.set('w', 'med'); head.set('len', 'med')
    return conn


def chip(slide, x, y, w, h, text, *, name=None, color=CYAN,
         bg_color=None, size=10, font="Consolas"):
    """Petite pilule (chip) néon."""
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    if name: s.name = name
    s.adjustments[0] = 0.5
    s.fill.solid(); s.fill.fore_color.rgb = bg_color if bg_color else INK
    s.line.color.rgb = color; s.line.width = Pt(1)
    tf = s.text_frame
    tf.margin_left = Pt(8); tf.margin_right = Pt(8)
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    p.text = text
    for r in p.runs:
        r.font.name = font; r.font.size = Pt(size)
        r.font.color.rgb = color; r.font.bold = True
    return s


def dot(slide, cx, cy, r, color, name=None):
    return neon_circle(slide, cx, cy, r, color, name, glow=False)


def big_number(slide, x, y, w, h, num, color=CYAN, name=None):
    """Nombre geant style hero."""
    return add_text(slide, x, y, w, h, num, name=name,
                    size=160, bold=True, color=color,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                    font="Calibri")


def progress_bar(slide, x, y, w, h, percent, color=CYAN, name=None):
    """Barre de progression neon."""
    bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    if name: bg.name = name + "_bg"
    bg.adjustments[0] = 0.5
    bg.fill.solid(); bg.fill.fore_color.rgb = INK
    bg.line.color.rgb = SUBTLE; bg.line.width = Pt(0.5)
    fill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   x, y, int(w * percent), h)
    if name: fill.name = name + "_fill"
    fill.adjustments[0] = 0.5
    fill.fill.solid(); fill.fill.fore_color.rgb = color
    fill.line.fill.background()


def page_number(slide, prs, n, total):
    add_text(slide, prs.slide_width - Inches(1.4), prs.slide_height - Inches(0.4),
             Inches(1), Inches(0.3),
             f"{n:02d} / {total:02d}",
             size=9, color=DIM, align=PP_ALIGN.RIGHT, font="Consolas",
             letter_spacing=200)


def footer_brand(slide, prs):
    add_text(slide, Inches(0.5), prs.slide_height - Inches(0.4),
             Inches(4), Inches(0.3),
             "E-STORE  ·  KEYNOTE",
             size=9, color=DIM, font="Consolas", letter_spacing=400)


def section_marker(slide, prs, label, n):
    """Etiquette de section discrete en haut."""
    add_text(slide, Inches(0.5), Inches(0.4), Inches(4), Inches(0.3),
             f"// {label}",
             size=10, color=CYAN, font="Consolas", letter_spacing=300)
    add_text(slide, Inches(0.5), Inches(0.65), Inches(4), Inches(0.3),
             f"CHAPITRE {n:02d}",
             size=8, color=DIM, font="Consolas", letter_spacing=500)


def blank(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, DEEP)
    return s


# =====================================================================
#  STORYBOARD - 38 SLIDES en 11 ACTES
# =====================================================================
TOTAL = 38


# ─── ACTE 1 — OUVERTURE CINEMATIQUE (slides 1-5) ─────────────────────
def s01_black(prs):
    """Slide noire - anticipation."""
    s = blank(prs)
    # Un seul point lumineux au centre
    cx = prs.slide_width // 2
    cy = prs.slide_height // 2
    dot(s, cx, cy, Inches(0.05), CYAN, name="logo_dot")
    add_fade_transition(s)
    set_speaker_notes(s, "Pause cinématique : laisser quelques secondes de silence avant de parler. "
                         "L'attention se concentre sur le point lumineux. Quand le silence devient dense, "
                         "appuyer sur la touche pour déclencher la transition Morph.")
    return s


def s02_logo_zoom(prs):
    """Le point grandit en cercle - zoom Morph."""
    s = blank(prs)
    cx = prs.slide_width // 2
    cy = prs.slide_height // 2
    dot(s, cx, cy, Inches(1.2), CYAN, name="logo_dot")
    # Halo violet
    halo = s.shapes.add_shape(MSO_SHAPE.OVAL,
                              cx - Inches(2.4), cy - Inches(2.4),
                              Inches(4.8), Inches(4.8))
    halo.name = "halo"
    halo.fill.solid(); halo.fill.fore_color.rgb = VIOLET
    halo.line.fill.background()
    add_morph_transition(s)
    set_speaker_notes(s, "Le point grossit (Morph) — moment d'impact. Annoncer : 'Bonjour à tous.' "
                         "puis attendre que la transition suivante révèle le titre.")
    return s


def s03_title_reveal(prs):
    """Le cercle devient le 'O' de E-STORE - typographie geante."""
    s = blank(prs)
    cx = prs.slide_width // 2
    cy = prs.slide_height // 2
    # Le dot reste mais se decale et passe en arriere-plan via couleur
    dot(s, Inches(7.5), cy - Inches(0.05), Inches(0.6), CYAN, name="logo_dot")
    # Texte E—STORE
    add_text(s, Inches(0.5), Inches(2.8), Inches(12.3), Inches(2),
             "E—STORE", name="hero_title",
             size=180, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, font="Calibri", letter_spacing=-100)
    add_morph_transition(s)
    set_speaker_notes(s, "REVELATION DU NOM. Annoncer : 'Nous vous présentons E-STORE.' "
                         "Marquer une pause après le nom pour laisser l'impact se faire.")
    return s


def s04_tagline(prs):
    """Tagline — la promesse du produit."""
    s = blank(prs)
    cx = prs.slide_width // 2
    cy = prs.slide_height // 2
    # Le titre se compresse en haut
    dot(s, Inches(7.5), Inches(1.55), Inches(0.4), CYAN, name="logo_dot")
    add_text(s, Inches(0.5), Inches(0.9), Inches(12.3), Inches(1.2),
             "E—STORE", name="hero_title",
             size=72, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, font="Calibri", letter_spacing=-50)

    # Tagline
    add_text(s, Inches(0.5), Inches(3), Inches(12.3), Inches(0.8),
             "Une plateforme e-commerce full-stack",
             name="tagline_1",
             size=36, color=SOFT, align=PP_ALIGN.CENTER, font="Calibri")
    add_text(s, Inches(0.5), Inches(3.8), Inches(12.3), Inches(0.8),
             "conçue pour la performance, sécurisée par construction.",
             name="tagline_2",
             size=36, italic=True, color=CYAN, align=PP_ALIGN.CENTER, font="Calibri")

    # Trois mots-cles dans des chips
    keywords = [("MODERNE", CYAN), ("HYBRIDE", VIOLET), ("SECURISE", PINK)]
    cw = Inches(2.5); ch = Inches(0.5); spacing = Inches(0.3)
    total = cw * 3 + spacing * 2
    sx = (prs.slide_width - total) // 2
    for i, (kw, col) in enumerate(keywords):
        chip(s, sx + (cw + spacing) * i, Inches(5.2), cw, ch, kw,
             name=f"kw_{i}", color=col, size=11)

    add_morph_transition(s)
    set_speaker_notes(s, "Annoncer la promesse. 'E-Store, c'est une plateforme e-commerce full-stack "
                         "conçue pour la performance, sécurisée par construction.' "
                         "Insister sur les trois mots-clés : moderne, hybride, sécurisée.")
    return s


def s05_authors(prs):
    """Carte d'identite du projet."""
    s = blank(prs)
    # Titre encore plus petit en haut
    dot(s, Inches(7.5), Inches(0.95), Inches(0.25), CYAN, name="logo_dot")
    add_text(s, Inches(0.5), Inches(0.6), Inches(12.3), Inches(0.6),
             "E—STORE", name="hero_title",
             size=32, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, font="Calibri")

    # Trois cartes (auteurs / encadrant / établissement)
    cards = [
        ("AUTEURS", "Akram Belmoussa\nNouhaila Ben Soumane", CYAN, "name_card_1"),
        ("ENCADRANT", "Pr. Omar Zahour", VIOLET, "name_card_2"),
        ("ÉTABLISSEMENT", "FSBM\nUniversité Hassan II", PINK, "name_card_3"),
    ]
    cw = Inches(3.6); ch = Inches(2.5); spacing = Inches(0.3)
    total = cw * 3 + spacing * 2
    sx = (prs.slide_width - total) // 2
    y = Inches(2.4)
    for i, (lbl, txt, col, nm) in enumerate(cards):
        x = sx + (cw + spacing) * i
        glass_card(s, x, y, cw, ch, name=nm, accent=col, accent_pos="top")
        add_text(s, x, y + Inches(0.4), cw, Inches(0.4),
                 lbl, size=10, color=col, align=PP_ALIGN.CENTER,
                 font="Consolas", letter_spacing=400)
        add_text(s, x + Inches(0.3), y + Inches(0.9),
                 cw - Inches(0.6), Inches(1.5),
                 txt, size=18, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, font="Calibri")

    add_text(s, Inches(0.5), Inches(5.5), Inches(12.3), Inches(0.4),
             "MODULE FULL-STACK · ANNÉE 2025-2026 · MAROC",
             size=10, color=DIM, align=PP_ALIGN.CENTER,
             font="Consolas", letter_spacing=600)

    page_number(s, prs, 5, TOTAL)
    add_morph_transition(s)
    set_speaker_notes(s, "Présenter brièvement : 'Je suis Akram, voici ma binôme Nouhaila. "
                         "Ce projet a été réalisé sous la direction du Pr. Omar Zahour, "
                         "à la Faculté des Sciences Ben M'Sick.' Tendre la main vers Nouhaila à ce moment.")
    return s


# ─── ACTE 2 — VISION & PROBLEMATIQUE (slides 6-9) ────────────────────
def s06_vision(prs):
    s = blank(prs)
    section_marker(s, prs, "VISION", 1)

    add_text(s, Inches(0.5), Inches(2.8), Inches(12.3), Inches(1.5),
             "Et si construire", name="big_q1",
             size=64, color=SOFT, align=PP_ALIGN.CENTER, font="Calibri")
    add_text(s, Inches(0.5), Inches(3.8), Inches(12.3), Inches(1.5),
             "une vraie plateforme e-commerce", name="big_q2",
             size=64, color=WHITE, align=PP_ALIGN.CENTER, font="Calibri", bold=True)
    add_text(s, Inches(0.5), Inches(4.8), Inches(12.3), Inches(1.5),
             "tenait dans un seul projet ?", name="big_q3",
             size=64, italic=True, color=CYAN, align=PP_ALIGN.CENTER, font="Calibri")

    page_number(s, prs, 6, TOTAL)
    add_morph_transition(s)
    set_speaker_notes(s, "Question rhétorique. Marquer une pause après chaque ligne. "
                         "C'est le moment où le jury comprend l'ampleur de ce qu'on a construit.")
    return s


def s07_challenge(prs):
    s = blank(prs)
    section_marker(s, prs, "DEFI", 1)

    add_text(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(1),
             "Un défi à 6 dimensions",
             size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Calibri")

    challenges = [
        ("●", "Authentification\nrobuste",          CYAN),
        ("◆", "Catalogue\nperformant",               VIOLET),
        ("▲", "Stock & panier\nthread-safe",         PINK),
        ("◇", "Commandes\ntransactionnelles",        LIME),
        ("◈", "Persistance\nhybride",                AMBER),
        ("◆", "Déploiement\nmulti-environnement",    CORAL),
    ]
    cw = Inches(2.0); ch = Inches(2.5); spacing = Inches(0.15)
    total = cw * 6 + spacing * 5
    sx = (prs.slide_width - total) // 2
    y = Inches(3)
    for i, (icon, txt, col) in enumerate(challenges):
        x = sx + (cw + spacing) * i
        glass_card(s, x, y, cw, ch, name=f"chal_{i}", accent=col, accent_pos="top")
        add_text(s, x, y + Inches(0.5), cw, Inches(0.6),
                 icon, size=40, color=col, align=PP_ALIGN.CENTER, bold=True)
        add_text(s, x + Inches(0.1), y + Inches(1.3), cw - Inches(0.2), Inches(1.2),
                 txt, size=12, color=WHITE, align=PP_ALIGN.CENTER, font="Calibri")

    page_number(s, prs, 7, TOTAL)
    add_morph_transition(s)
    set_speaker_notes(s, "Énumérer rapidement les 6 défis. Insister : 'Chacun de ces points est "
                         "déjà un projet en soi. Notre engagement : tout livrer dans une architecture cohérente.'")
    return s


def s08_metrics_preview(prs):
    """Aperçu des chiffres clés du projet."""
    s = blank(prs)
    section_marker(s, prs, "PAR LES CHIFFRES", 1)

    metrics = [
        ("38", "slides keynote",      CYAN,    "metric_38"),
        ("30+", "fichiers source",     VIOLET,  "metric_30"),
        ("18", "endpoints REST",      PINK,    "metric_18"),
        ("9", "tests unitaires",      LIME,    "metric_9"),
        ("0", "failure",              AMBER,   "metric_0"),
    ]
    cw = Inches(2.4); ch = Inches(2.8); spacing = Inches(0.25)
    total = cw * 5 + spacing * 4
    sx = (prs.slide_width - total) // 2
    y = Inches(2.5)
    for i, (num, label, col, nm) in enumerate(metrics):
        x = sx + (cw + spacing) * i
        glass_card(s, x, y, cw, ch, name=nm, accent=col, accent_pos="top")
        add_text(s, x, y + Inches(0.5), cw, Inches(1.4),
                 num, size=72, bold=True, color=col,
                 align=PP_ALIGN.CENTER, font="Calibri", letter_spacing=-50)
        add_text(s, x + Inches(0.1), y + Inches(1.95), cw - Inches(0.2), Inches(0.8),
                 label, size=11, color=SOFT, align=PP_ALIGN.CENTER,
                 font="Consolas", letter_spacing=200)

    page_number(s, prs, 8, TOTAL)
    add_morph_transition(s)
    set_speaker_notes(s, "Énoncer les chiffres : '30+ fichiers, 18 endpoints REST, 9 tests, 0 échec.' "
                         "Ces métriques sont la signature de la rigueur appliquée.")
    return s


def s09_promise(prs):
    s = blank(prs)
    section_marker(s, prs, "PROMESSE", 1)

    add_text(s, Inches(0.5), Inches(1.8), Inches(12.3), Inches(1),
             "Notre promesse",
             size=42, color=DIM, align=PP_ALIGN.CENTER, font="Calibri", italic=True)

    add_text(s, Inches(0.5), Inches(3), Inches(12.3), Inches(1.2),
             "Architecture propre.",
             name="promise_1",
             size=64, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Calibri")
    add_text(s, Inches(0.5), Inches(4), Inches(12.3), Inches(1.2),
             "Code testé.",
             name="promise_2",
             size=64, bold=True, color=CYAN, align=PP_ALIGN.CENTER, font="Calibri")
    add_text(s, Inches(0.5), Inches(5), Inches(12.3), Inches(1.2),
             "Démo qui fonctionne.",
             name="promise_3",
             size=64, bold=True, color=VIOLET, align=PP_ALIGN.CENTER, font="Calibri")

    page_number(s, prs, 9, TOTAL)
    add_morph_transition(s)
    set_speaker_notes(s, "Trois engagements simples. Les énoncer un à un avec conviction. "
                         "Faire une pause après chaque phrase pour que le jury intègre.")
    return s


# ─── ACTE 3 — ARCHITECTURE LAYER PAR LAYER (slides 10-15) ────────────
def s10_arch_intro(prs):
    s = blank(prs)
    section_marker(s, prs, "ARCHITECTURE", 2)

    add_text(s, Inches(0.5), Inches(2.5), Inches(12.3), Inches(1.5),
             "Architecture en couches",
             size=58, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Calibri")
    add_text(s, Inches(0.5), Inches(3.6), Inches(12.3), Inches(1),
             "3 strates techniques  ·  6 domaines fonctionnels",
             size=24, color=CYAN, align=PP_ALIGN.CENTER, font="Calibri", italic=True)

    # Cube vide central qui va se remplir
    cx = prs.slide_width // 2 - Inches(1.5)
    cy = Inches(5)
    cube = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, cy, Inches(3), Inches(1.5))
    cube.name = "arch_cube"
    cube.adjustments[0] = 0.06
    cube.fill.solid(); cube.fill.fore_color.rgb = INK
    cube.line.color.rgb = LINE; cube.line.width = Pt(1)
    add_text(s, cx, cy + Inches(0.5), Inches(3), Inches(0.5),
             "● ● ●", size=24, color=DIM, align=PP_ALIGN.CENTER)

    page_number(s, prs, 10, TOTAL)
    add_morph_transition(s)
    set_speaker_notes(s, "Annoncer le chapitre architecture. Le cube va prendre vie progressivement "
                         "à chaque slide suivante via Morph.")
    return s


def _arch_layer_slide(prs, n_layers_visible, title=None, sub=None):
    """Construit une slide architecture avec n_layers (1 a 3) visibles."""
    s = blank(prs)
    section_marker(s, prs, "ARCHITECTURE", 2)

    # Titre
    add_text(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.7),
             title or "Architecture en couches",
             size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Calibri")
    if sub:
        add_text(s, Inches(0.5), Inches(1.9), Inches(12.3), Inches(0.5),
                 sub, size=14, color=CYAN, align=PP_ALIGN.CENTER,
                 italic=True, font="Calibri")

    layers = [
        ("PRESENTATION",   "Angular 17 standalone\nsignals · Bootstrap 5",    CYAN,   "layer_1"),
        ("LOGIQUE METIER", "Spring Boot 3.3\n@Transactional · DTOs",          VIOLET, "layer_2"),
        ("ACCES DONNEES",  "Spring Data JPA + Hibernate\nSpring Data MongoDB", PINK,   "layer_3"),
    ]
    cw = Inches(3.8); ch = Inches(1.5); spacing = Inches(0.15)
    sx = (prs.slide_width - (cw * 3 + spacing * 2)) // 2
    y = Inches(3)

    for i, (name, body, col, nm) in enumerate(layers):
        x = sx + (cw + spacing) * i
        if i < n_layers_visible:
            # Visible
            glass_card(s, x, y, cw, ch, name=nm, accent=col, accent_pos="left")
            add_text(s, x + Inches(0.3), y + Inches(0.2),
                     cw - Inches(0.5), Inches(0.5),
                     name, size=15, bold=True, color=col,
                     font="Consolas", letter_spacing=200)
            add_text(s, x + Inches(0.3), y + Inches(0.7),
                     cw - Inches(0.5), Inches(0.8),
                     body, size=11, color=SOFT, font="Calibri")
        else:
            # Placeholder gris (preserve shape pour Morph)
            placeholder = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                              x, y, cw, ch)
            placeholder.name = nm
            placeholder.adjustments[0] = 0.05
            placeholder.fill.solid(); placeholder.fill.fore_color.rgb = NIGHT
            placeholder.line.color.rgb = SUBTLE; placeholder.line.width = Pt(0.5)
            add_text(s, x, y + Inches(0.5), cw, Inches(0.5),
                     "...", size=24, color=FAINT, align=PP_ALIGN.CENTER)

    return s


def s11_arch_layer1(prs):
    s = _arch_layer_slide(prs, 1, "Couche 1 : Présentation",
                           "L'utilisateur interagit ici - Angular 17 standalone")
    page_number(s, prs, 11, TOTAL)
    add_morph_transition(s)
    set_speaker_notes(s, "Apparition (Morph) de la couche 1 : Présentation. "
                         "'Tout commence dans le navigateur, avec Angular 17 et ses signals.'")
    return s


def s12_arch_layer2(prs):
    s = _arch_layer_slide(prs, 2, "Couche 2 : Logique métier",
                           "Le cerveau - Spring Boot orchestre les règles")
    page_number(s, prs, 12, TOTAL)
    add_morph_transition(s)
    set_speaker_notes(s, "Apparition (Morph) de la couche 2 : Logique métier. "
                         "'Spring Boot encaisse les requêtes, valide, oriente vers les services.'")
    return s


def s13_arch_layer3(prs):
    s = _arch_layer_slide(prs, 3, "Couche 3 : Accès aux données",
                           "Persistance - JPA + MongoDB côte à côte")
    page_number(s, prs, 13, TOTAL)
    add_morph_transition(s)
    set_speaker_notes(s, "Apparition (Morph) de la couche 3 : Accès aux données. "
                         "'Et au plus profond, deux bases de données complémentaires : MySQL et MongoDB.'")
    return s


def s14_domains(prs):
    """Decomposition en 6 domaines."""
    s = blank(prs)
    section_marker(s, prs, "ARCHITECTURE", 2)
    add_text(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.7),
             "6 domaines fonctionnels",
             size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Calibri")
    add_text(s, Inches(0.5), Inches(1.9), Inches(12.3), Inches(0.5),
             "Chaque domaine est un sous-package autonome (DDD)",
             size=14, color=CYAN, align=PP_ALIGN.CENTER, italic=True)

    domains = [
        ("customer",  "User · Profile\nJWT · Auth",          CYAN),
        ("catalog",   "Catégorie · Produit\nRecherche",       VIOLET),
        ("inventory", "Stock · Décrément\natomique",           PINK),
        ("shopping",  "Cart · CartItem\nPrix figé",            LIME),
        ("billing",   "Order · OrderItem\nCheckout TX",        AMBER),
        ("review",    "Avis MongoDB\nLecture frequente",       CORAL),
    ]
    cw = Inches(3.8); ch = Inches(1.5); spacing = Inches(0.15)
    total_w = cw * 3 + spacing * 2
    sx = (prs.slide_width - total_w) // 2
    for i, (name, body, col) in enumerate(domains):
        c = i % 3; r = i // 3
        x = sx + (cw + spacing) * c
        y = Inches(2.9) + (ch + Inches(0.25)) * r
        glass_card(s, x, y, cw, ch, name=f"dom_{i}", accent=col, accent_pos="left")
        add_text(s, x + Inches(0.3), y + Inches(0.2),
                 cw - Inches(0.5), Inches(0.5),
                 name, size=16, bold=True, color=col,
                 font="Consolas", letter_spacing=300)
        add_text(s, x + Inches(0.3), y + Inches(0.7),
                 cw - Inches(0.5), Inches(0.8),
                 body, size=11, color=SOFT, font="Calibri")

    page_number(s, prs, 14, TOTAL)
    add_morph_transition(s)
    set_speaker_notes(s, "Six domaines apparaissent ensemble (Morph). "
                         "Souligner : 'Chaque domaine a son entité, son DTO, son service, son controller.' "
                         "C'est l'application directe du Domain-Driven Design.")
    return s


def s15_domain_zoom(prs):
    """Zoom sur le domaine 'billing' - exemple concret."""
    s = blank(prs)
    section_marker(s, prs, "ARCHITECTURE", 2)
    add_text(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.7),
             "Anatomie d'un domaine : billing",
             size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Calibri")

    # Pyramide inversee : controller -> service -> repository -> entity
    layers = [
        ("OrderController",  "@RestController · routes /api/orders",  CYAN),
        ("OrderService",     "@Service @Transactional · checkout()",  VIOLET),
        ("OrderRepository",  "extends JpaRepository<Order, Long>",    PINK),
        ("Order + OrderItem","@Entity · OneToMany cascade",           LIME),
    ]
    y = Inches(2.5)
    for i, (name, body, col) in enumerate(layers):
        # Largeur decroissante pour effet pyramide
        w = Inches(11 - i * 1.5)
        h = Inches(0.85)
        x = (prs.slide_width - w) // 2
        glass_card(s, x, y, w, h, name=f"billing_{i}", accent=col, accent_pos="left")
        add_text(s, x + Inches(0.3), y + Inches(0.05), Inches(4), Inches(0.4),
                 name, size=14, bold=True, color=col, font="Consolas")
        add_text(s, x + Inches(4.5), y + Inches(0.15), w - Inches(5), Inches(0.6),
                 body, size=12, color=SOFT, font="Consolas")
        y += Inches(1.05)

    page_number(s, prs, 15, TOTAL)
    add_morph_transition(s)
    set_speaker_notes(s, "Zoom sur un domaine. La pyramide se construit (Morph). "
                         "'Voici le schéma typique : controller, service, repository, entity. "
                         "Cette structure se retrouve dans tous les domaines.'")
    return s


# ─── ACTE 4 — STACK TECHNIQUE (slides 16-20) ─────────────────────────
def s16_stack_galaxy(prs):
    """Galaxie de technologies."""
    s = blank(prs)
    section_marker(s, prs, "STACK", 3)
    add_text(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.7),
             "Une galaxie de technologies",
             size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Calibri")

    # Noyau central
    cx = prs.slide_width // 2
    cy = Inches(4.3)
    nucleus = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                  cx - Inches(0.9), cy - Inches(0.9),
                                  Inches(1.8), Inches(1.8))
    nucleus.name = "nucleus"
    nucleus.fill.solid(); nucleus.fill.fore_color.rgb = INK
    nucleus.line.color.rgb = CYAN; nucleus.line.width = Pt(2)
    add_text(s, cx - Inches(1.5), cy - Inches(0.3), Inches(3), Inches(0.6),
             "E-STORE", size=20, bold=True, color=CYAN,
             align=PP_ALIGN.CENTER, font="Calibri")

    # Satellites en cercle
    sats = [
        ("Angular 17",      CYAN,    -90),
        ("Spring Boot",     VIOLET,  -30),
        ("MySQL · H2",      PINK,     30),
        ("MongoDB",         LIME,     90),
        ("Spring Security", AMBER,   150),
        ("Maven · npm",     CORAL,  -150),
    ]
    R = Inches(2.8)
    for label, col, angle in sats:
        rad = math.radians(angle)
        sx = cx + int(R * math.cos(rad)) - Inches(1)
        sy = cy + int(R * math.sin(rad)) - Inches(0.4)
        # Connecteur depuis le noyau
        line_connector(s, cx, cy,
                       cx + int(R * math.cos(rad)),
                       cy + int(R * math.sin(rad)),
                       color=SUBTLE, width=0.75)
        # Pilule satellite
        chip(s, sx, sy, Inches(2), Inches(0.6), label,
             color=col, size=11, font="Calibri")

    page_number(s, prs, 16, TOTAL)
    add_morph_transition(s)
    set_speaker_notes(s, "Le système solaire de E-Store. Au centre l'application, "
                         "autour ses satellites technologiques. Ils tous rayonnent vers le noyau.")
    return s


def s17_frontend_stack(prs):
    s = blank(prs)
    section_marker(s, prs, "FRONTEND", 3)
    add_text(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.7),
             "Frontend Angular 17",
             size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Calibri")

    items = [
        ("Standalone components", "Pas de NgModule. Imports directs.\nLazy-loading natif par route.", CYAN),
        ("Signals reactifs",       "Etat reactif sans RxJS.\ncart.itemCount() -> badge auto-update.",  VIOLET),
        ("Interceptors fonctionnels","AuthInterceptor + ErrorInterceptor.\nGestion centralisee.",      PINK),
        ("Bootstrap 5",             "CSS-only. Grid responsive.\nThematique cohesive.",               LIME),
    ]
    cw = Inches(6); ch = Inches(1.7); spacing = Inches(0.25)
    sx = Inches(0.4)
    for i, (t, b, col) in enumerate(items):
        c = i % 2; r = i // 2
        x = sx + (cw + spacing) * c
        y = Inches(2.5) + (ch + Inches(0.2)) * r
        glass_card(s, x, y, cw, ch, name=f"fe_{i}", accent=col, accent_pos="left")
        add_text(s, x + Inches(0.3), y + Inches(0.25), cw - Inches(0.5), Inches(0.5),
                 t, size=16, bold=True, color=col, font="Calibri")
        add_text(s, x + Inches(0.3), y + Inches(0.85), cw - Inches(0.5), Inches(0.8),
                 b, size=12, color=SOFT, font="Calibri")

    # Stat de bundle
    add_text(s, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.4),
             "BUNDLE INITIAL : 13 KO  ·  TIME TO INTERACTIVE < 1 S",
             size=11, color=DIM, align=PP_ALIGN.CENTER,
             font="Consolas", letter_spacing=400)

    page_number(s, prs, 17, TOTAL)
    add_morph_transition(s)
    set_speaker_notes(s, "Quatre piliers frontend. Insister sur les signals : 'C'est la nouveauté Angular 17 "
                         "que nous avons exploitée pour la réactivité du panier.' "
                         "Et sur les 13 ko de bundle initial : exceptionnel.")
    return s


def s18_backend_stack(prs):
    s = blank(prs)
    section_marker(s, prs, "BACKEND", 3)
    add_text(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.7),
             "Backend Spring Boot 3.3",
             size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Calibri")

    items = [
        ("Spring Boot 3.3", "Java 17 · Auto-configuration\nProductivite × 10",     CYAN),
        ("Spring Security 6", "Filter chain · @PreAuthorize\nStateless JWT",        VIOLET),
        ("Spring Data JPA",  "JpaRepository derive\nRequetes JPQL declarees",       PINK),
        ("Spring Data Mongo","MongoRepository derive\nIndexed @Indexed",            LIME),
    ]
    cw = Inches(6); ch = Inches(1.7); spacing = Inches(0.25)
    sx = Inches(0.4)
    for i, (t, b, col) in enumerate(items):
        c = i % 2; r = i // 2
        x = sx + (cw + spacing) * c
        y = Inches(2.5) + (ch + Inches(0.2)) * r
        glass_card(s, x, y, cw, ch, name=f"be_{i}", accent=col, accent_pos="left")
        add_text(s, x + Inches(0.3), y + Inches(0.25), cw - Inches(0.5), Inches(0.5),
                 t, size=16, bold=True, color=col, font="Calibri")
        add_text(s, x + Inches(0.3), y + Inches(0.85), cw - Inches(0.5), Inches(0.8),
                 b, size=12, color=SOFT, font="Calibri")

    add_text(s, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.4),
             "DEMARRAGE A FROID < 8 S  ·  AUTO-RELOAD EN DEV",
             size=11, color=DIM, align=PP_ALIGN.CENTER,
             font="Consolas", letter_spacing=400)

    page_number(s, prs, 18, TOTAL)
    add_morph_transition(s)
    set_speaker_notes(s, "Le quatuor Spring : Boot, Security, Data JPA, Data MongoDB. "
                         "'Spring permet de faire en une annotation ce qui prendrait 100 lignes ailleurs.' "
                         "Démarrage à froid en moins de 8 secondes.")
    return s


def s19_data_hybrid(prs):
    s = blank(prs)
    section_marker(s, prs, "DATA", 3)
    add_text(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.7),
             "Persistance hybride",
             size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Calibri")
    add_text(s, Inches(0.5), Inches(1.9), Inches(12.3), Inches(0.5),
             "Le bon outil pour la bonne donnée",
             size=14, color=CYAN, align=PP_ALIGN.CENTER, italic=True)

    # Deux grandes cartes cote a cote
    # Gauche : SQL
    sql_x = Inches(0.6); sql_y = Inches(2.7); sql_w = Inches(6); sql_h = Inches(4)
    glass_card(s, sql_x, sql_y, sql_w, sql_h, name="sql_card",
               accent=CYAN, accent_pos="top")
    add_text(s, sql_x + Inches(0.4), sql_y + Inches(0.3), sql_w - Inches(0.8), Inches(0.5),
             "MySQL 8 · H2", size=22, bold=True, color=CYAN, font="Calibri")
    add_text(s, sql_x + Inches(0.4), sql_y + Inches(0.9), sql_w - Inches(0.8), Inches(0.4),
             "RELATIONNEL · ACID", size=10, color=DIM,
             font="Consolas", letter_spacing=400)
    sql_items = [
        "● Utilisateurs, profils, sessions",
        "● Catalogue produits + catégories",
        "● Stock & inventaire en temps réel",
        "● Panier et items du panier",
        "● Commandes transactionnelles",
    ]
    for i, it in enumerate(sql_items):
        add_text(s, sql_x + Inches(0.5), sql_y + Inches(1.5) + Inches(0.4) * i,
                 sql_w - Inches(0.9), Inches(0.4),
                 it, size=12, color=SOFT, font="Calibri")

    # Droite : NoSQL
    no_x = Inches(6.8); no_y = Inches(2.7); no_w = Inches(6); no_h = Inches(4)
    glass_card(s, no_x, no_y, no_w, no_h, name="nosql_card",
               accent=VIOLET, accent_pos="top")
    add_text(s, no_x + Inches(0.4), no_y + Inches(0.3), no_w - Inches(0.8), Inches(0.5),
             "MongoDB 7", size=22, bold=True, color=VIOLET, font="Calibri")
    add_text(s, no_x + Inches(0.4), no_y + Inches(0.9), no_w - Inches(0.8), Inches(0.4),
             "DOCUMENTAIRE · FLEXIBLE", size=10, color=DIM,
             font="Consolas", letter_spacing=400)
    no_items = [
        "◆ Avis et notes des produits",
        "◆ Schema flexible (JSON)",
        "◆ Lecture haute frequence",
        "◆ findByProductIdOrderByDate",
        "◆ Index automatique @Indexed",
    ]
    for i, it in enumerate(no_items):
        add_text(s, no_x + Inches(0.5), no_y + Inches(1.5) + Inches(0.4) * i,
                 no_w - Inches(0.9), Inches(0.4),
                 it, size=12, color=SOFT, font="Calibri")

    page_number(s, prs, 19, TOTAL)
    add_morph_transition(s)
    set_speaker_notes(s, "Persistance hybride : 'Pas de religion technique. MySQL pour les transactions, "
                         "MongoDB pour les avis. Le bon outil pour la bonne donnée.'")
    return s


def s20_devops(prs):
    s = blank(prs)
    section_marker(s, prs, "DEVOPS", 3)
    add_text(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.7),
             "DevOps & environnements",
             size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Calibri")

    # Deux profils
    cw = Inches(5.8); ch = Inches(4); spacing = Inches(0.3)
    sx = (prs.slide_width - (cw * 2 + spacing)) // 2
    y = Inches(2.5)

    # Profil DEV
    glass_card(s, sx, y, cw, ch, name="profile_dev", accent=CYAN, accent_pos="top")
    add_text(s, sx + Inches(0.4), y + Inches(0.3), cw - Inches(0.8), Inches(0.4),
             "PROFIL DEV", size=10, color=CYAN,
             font="Consolas", letter_spacing=500)
    add_text(s, sx + Inches(0.4), y + Inches(0.7), cw - Inches(0.8), Inches(0.5),
             "H2 in-memory",
             size=24, bold=True, color=WHITE, font="Calibri")
    items = [
        "▸ Aucune installation requise",
        "▸ Demarrage en 8 secondes",
        "▸ Console H2 sur /h2-console",
        "▸ Idéal soutenance & démo",
    ]
    for i, it in enumerate(items):
        add_text(s, sx + Inches(0.5), y + Inches(1.5) + Inches(0.4) * i,
                 cw - Inches(0.9), Inches(0.4),
                 it, size=12, color=SOFT, font="Calibri")

    # Profil PROD
    glass_card(s, sx + cw + spacing, y, cw, ch,
               name="profile_prod", accent=VIOLET, accent_pos="top")
    add_text(s, sx + cw + spacing + Inches(0.4), y + Inches(0.3),
             cw - Inches(0.8), Inches(0.4),
             "PROFIL PROD", size=10, color=VIOLET,
             font="Consolas", letter_spacing=500)
    add_text(s, sx + cw + spacing + Inches(0.4), y + Inches(0.7),
             cw - Inches(0.8), Inches(0.5),
             "Docker Compose",
             size=24, bold=True, color=WHITE, font="Calibri")
    items_prod = [
        "▸ MySQL 8 + MongoDB 7",
        "▸ phpMyAdmin (port 8081)",
        "▸ mongo-express (port 8082)",
        "▸ Volumes persistents",
    ]
    for i, it in enumerate(items_prod):
        add_text(s, sx + cw + spacing + Inches(0.5), y + Inches(1.5) + Inches(0.4) * i,
                 cw - Inches(0.9), Inches(0.4),
                 it, size=12, color=SOFT, font="Calibri")

    page_number(s, prs, 20, TOTAL)
    add_morph_transition(s)
    set_speaker_notes(s, "Deux environnements pensés dès le départ. La même application, deux modes. "
                         "C'est ce qui nous a sauvés quand Docker a refusé de s'installer.")
    return s


# ─── ACTE 5 — JWT FLOW MORPH (slides 21-25) ──────────────────────────
def _jwt_slide(prs, visible_steps, title_emphasis=None):
    """Slide JWT avec etapes progressives."""
    s = blank(prs)
    section_marker(s, prs, "SECURITE", 4)
    add_text(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.7),
             "Cycle de vie d'un JWT",
             size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Calibri")
    if title_emphasis:
        add_text(s, Inches(0.5), Inches(1.9), Inches(12.3), Inches(0.5),
                 title_emphasis, size=14, color=CYAN, align=PP_ALIGN.CENTER, italic=True)

    # Deux colonnes Client / Serveur
    add_text(s, Inches(1.5), Inches(2.6), Inches(4), Inches(0.4),
             "▣ NAVIGATEUR CLIENT", size=11, color=CYAN,
             font="Consolas", letter_spacing=300, align=PP_ALIGN.CENTER)
    add_text(s, Inches(8), Inches(2.6), Inches(4), Inches(0.4),
             "▣ SERVEUR SPRING BOOT", size=11, color=VIOLET,
             font="Consolas", letter_spacing=300, align=PP_ALIGN.CENTER)

    # Colonnes (zones glass)
    col_y = Inches(3.1); col_h = Inches(3.7)
    glass_card(s, Inches(0.6), col_y, Inches(5.5), col_h,
               name="col_client", accent=CYAN, accent_pos="left")
    glass_card(s, Inches(7.2), col_y, Inches(5.5), col_h,
               name="col_server", accent=VIOLET, accent_pos="left")

    # 6 etapes
    steps = [
        (1, "client", "POST /api/auth/login\n(email + password)", CYAN),
        (2, "server", "Validation BCrypt\npasswordEncoder.matches()", VIOLET),
        (3, "server", "Generation JWT HS256\nsubject + role + uid (24h)", VIOLET),
        (4, "client", "Stockage localStorage\n'estore.token'", CYAN),
        (5, "client", "Requete suivante :\nAuthInterceptor → Bearer", CYAN),
        (6, "server", "JwtAuthFilter valide.\n@PreAuthorize verifie role", VIOLET),
    ]
    step_h = Inches(0.75); step_y_start = col_y + Inches(0.25); spacing = Inches(0.1)
    for i, (n, side, txt, col) in enumerate(steps):
        if n > visible_steps:
            continue
        y = step_y_start + (step_h + spacing) * i
        x_col = Inches(0.6) if side == "client" else Inches(7.2)
        # Pilule etape
        x = x_col + Inches(0.2)
        # Numero
        nc = s.shapes.add_shape(MSO_SHAPE.OVAL, x, y + Inches(0.18),
                                 Inches(0.4), Inches(0.4))
        nc.name = f"step_circle_{n}"
        nc.fill.solid(); nc.fill.fore_color.rgb = col
        nc.line.color.rgb = WHITE; nc.line.width = Pt(1)
        tf = nc.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        p.text = str(n)
        for r in p.runs:
            r.font.bold = True; r.font.size = Pt(13); r.font.color.rgb = DEEP

        # Texte
        add_text(s, x + Inches(0.55), y + Inches(0.1), Inches(4.5), step_h,
                 txt, name=f"step_text_{n}",
                 size=11, color=WHITE, font="Consolas", anchor=MSO_ANCHOR.MIDDLE)

    return s


def s21_jwt_step1(prs):
    s = _jwt_slide(prs, 1, "Étape 1 : le client demande l'accès")
    page_number(s, prs, 21, TOTAL)
    add_morph_transition(s)
    set_speaker_notes(s, "Première étape. 'L'utilisateur tape son email et mot de passe. "
                         "POST vers /api/auth/login.'")
    return s


def s22_jwt_step3(prs):
    s = _jwt_slide(prs, 3, "Étapes 2-3 : le serveur valide et signe")
    page_number(s, prs, 22, TOTAL)
    add_morph_transition(s)
    set_speaker_notes(s, "Étapes 2 et 3 (Morph). 'Le serveur vérifie BCrypt — impossible de retrouver "
                         "le mot de passe — puis génère un JWT signé HS256 valable 24 heures.'")
    return s


def s23_jwt_step5(prs):
    s = _jwt_slide(prs, 5, "Étapes 4-5 : le token voyage")
    page_number(s, prs, 23, TOTAL)
    add_morph_transition(s)
    set_speaker_notes(s, "Étapes 4 et 5 (Morph). 'Le token est stocké côté navigateur. "
                         "L'AuthInterceptor l'attache à chaque requête sortante.'")
    return s


def s24_jwt_complete(prs):
    s = _jwt_slide(prs, 6, "Boucle complète : 6 étapes verrouillées")
    page_number(s, prs, 24, TOTAL)
    add_morph_transition(s)
    set_speaker_notes(s, "La boucle est complète. 'À chaque requête, le filtre JWT vérifie. "
                         "Si le rôle ne correspond pas à @PreAuthorize, c'est 403 Forbidden.'")
    return s


def s25_security_summary(prs):
    s = blank(prs)
    section_marker(s, prs, "SECURITE", 4)
    add_text(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.7),
             "4 verrous de sécurité",
             size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Calibri")

    locks = [
        ("BCrypt", "Hash one-way avec sel.\nCoût 10 (>1500 ms par tentative).", CYAN),
        ("JWT HS256", "Signature HMAC-SHA-256.\nClé secrète 256 bits.",          VIOLET),
        ("Filter chain", "JwtAuthenticationFilter sur\nchaque requête HTTP.",     PINK),
        ("@PreAuthorize", "Verrou méthode.\nADMIN sur les CRUD sensibles.",       LIME),
    ]
    cw = Inches(2.95); ch = Inches(3.3); spacing = Inches(0.15)
    sx = (prs.slide_width - (cw * 4 + spacing * 3)) // 2
    y = Inches(2.5)
    for i, (t, b, col) in enumerate(locks):
        x = sx + (cw + spacing) * i
        glass_card(s, x, y, cw, ch, name=f"lock_{i}", accent=col, accent_pos="top")
        add_text(s, x, y + Inches(0.5), cw, Inches(0.7),
                 "🔒", size=42, color=col, align=PP_ALIGN.CENTER)
        add_text(s, x, y + Inches(1.4), cw, Inches(0.5),
                 t, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Calibri")
        add_text(s, x + Inches(0.2), y + Inches(2), cw - Inches(0.4), Inches(1.2),
                 b, size=11, color=SOFT, align=PP_ALIGN.CENTER, font="Calibri")

    page_number(s, prs, 25, TOTAL)
    add_morph_transition(s)
    set_speaker_notes(s, "Quatre verrous indépendants. Pour qu'un attaquant passe, il devrait casser les quatre. "
                         "C'est la défense en profondeur.")
    return s


# ─── ACTE 6 — TRANSACTIONS ATOMIQUES (slides 26-28) ──────────────────
def s26_atomic_intro(prs):
    s = blank(prs)
    section_marker(s, prs, "ATOMICITE", 5)
    add_text(s, Inches(0.5), Inches(2.5), Inches(12.3), Inches(1.5),
             "Tout ou rien.",
             size=84, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
             font="Calibri", letter_spacing=-50)
    add_text(s, Inches(0.5), Inches(4.3), Inches(12.3), Inches(0.8),
             "L'atomicité ACID, en pratique.",
             size=24, color=CYAN, align=PP_ALIGN.CENTER, italic=True)

    # Trois engrenages
    gx = (prs.slide_width - Inches(7)) // 2
    gy = Inches(5.5)
    for i in range(3):
        c = s.shapes.add_shape(MSO_SHAPE.OVAL,
                               gx + Inches(2.4) * i, gy,
                               Inches(0.8), Inches(0.8))
        c.name = f"gear_{i}"
        c.fill.solid(); c.fill.fore_color.rgb = [CYAN, VIOLET, PINK][i]
        c.line.fill.background()
        add_text(s, gx + Inches(2.4) * i, gy + Inches(0.2),
                 Inches(0.8), Inches(0.5),
                 ["1", "2", "3"][i], size=18, bold=True, color=DEEP,
                 align=PP_ALIGN.CENTER)

    page_number(s, prs, 26, TOTAL)
    add_morph_transition(s)
    set_speaker_notes(s, "Annonce dramatique. 'Quand un client clique sur Valider la commande, "
                         "trois choses doivent se passer ensemble. Ou aucune.'")
    return s


def s27_checkout_code(prs):
    s = blank(prs)
    section_marker(s, prs, "ATOMICITE", 5)
    add_text(s, Inches(0.5), Inches(0.95), Inches(12.3), Inches(0.5),
             "OrderService.checkout()",
             size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Consolas")
    add_text(s, Inches(0.5), Inches(1.45), Inches(12.3), Inches(0.4),
             "Le coeur transactionnel de l'application",
             size=12, color=CYAN, align=PP_ALIGN.CENTER, italic=True)

    # Cartouche code
    cx = Inches(0.5); cy = Inches(2.0); cw = Inches(7.5); ch = Inches(4.5)
    code_bg = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, cy, cw, ch)
    code_bg.name = "code_block"
    code_bg.adjustments[0] = 0.04
    code_bg.fill.solid(); code_bg.fill.fore_color.rgb = NIGHT
    code_bg.line.color.rgb = LINE; code_bg.line.width = Pt(0.5)

    # Petit dot row macOS
    for i, col in enumerate([RGBColor(0xFF, 0x60, 0x57),
                             RGBColor(0xFF, 0xBD, 0x2E),
                             RGBColor(0x28, 0xC9, 0x40)]):
        d = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                cx + Inches(0.2) + Inches(0.25) * i,
                                cy + Inches(0.18),
                                Inches(0.18), Inches(0.18))
        d.fill.solid(); d.fill.fore_color.rgb = col
        d.line.fill.background()

    code = """@Transactional
public OrderDto checkout() {
  Cart cart = cartService.getOrCreateCart(user);
  if (cart.getItems().isEmpty())
    throw new BusinessException("Panier vide");

  // 1) Verifier le stock pour TOUS les items
  for (CartItem ci : cart.getItems())
    inventoryService.checkAvailability(...);

  // 2) Creer Order + decrement stock
  Order order = Order.builder()...build();
  for (CartItem ci : cart.getItems()) {
    order.getItems().add(OrderItem...);
    inventoryService.decrement(...);
  }

  // 3) Confirmer + vider panier
  order.setStatus(CONFIRMED);
  cartService.clearCart(cart);
  return OrderDto.from(saved);
}"""
    add_text(s, cx + Inches(0.4), cy + Inches(0.5), cw - Inches(0.6), ch - Inches(0.6),
             code, size=11, color=WHITE, font="Consolas")

    # 3 highlights a droite
    high = [
        ("@Transactional", "Tout dans une transaction. Si une etape echoue, rollback automatique.", CYAN),
        ("checkAvailability", "Verifie AVANT toute modification. Pre-requis fiable.", VIOLET),
        ("decrement", "Decrement reel du stock. Coordonne avec la creation OrderItem.", PINK),
    ]
    hx = Inches(8.3); hy = Inches(2.0); hw = Inches(4.6); hh = Inches(1.45)
    for i, (t, b, col) in enumerate(high):
        y = hy + (hh + Inches(0.1)) * i
        glass_card(s, hx, y, hw, hh, name=f"high_{i}", accent=col, accent_pos="left")
        add_text(s, hx + Inches(0.3), y + Inches(0.2), hw - Inches(0.5), Inches(0.45),
                 t, size=14, bold=True, color=col, font="Consolas")
        add_text(s, hx + Inches(0.3), y + Inches(0.7), hw - Inches(0.5), Inches(0.7),
                 b, size=11, color=SOFT, font="Calibri")

    page_number(s, prs, 27, TOTAL)
    add_morph_transition(s)
    set_speaker_notes(s, "Le code clé de tout le projet. 'Une seule annotation @Transactional. "
                         "Trois étapes. Atomicité garantie par Spring.'")
    return s


def s28_acid_visual(prs):
    s = blank(prs)
    section_marker(s, prs, "ATOMICITE", 5)
    add_text(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.7),
             "ACID, illustré",
             size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Calibri")

    acid = [
        ("A", "Atomicity",   "Tout ou rien. Le checkout est indivisible.",      CYAN),
        ("C", "Consistency", "Stocks et commandes restent coherents.",           VIOLET),
        ("I", "Isolation",   "Deux clients en simultane ne se marchent dessus.", PINK),
        ("D", "Durability",  "Une fois COMMIT, c'est inscrit pour toujours.",    LIME),
    ]
    cw = Inches(2.95); ch = Inches(4); spacing = Inches(0.15)
    sx = (prs.slide_width - (cw * 4 + spacing * 3)) // 2
    y = Inches(2.5)
    for i, (letter, name, body, col) in enumerate(acid):
        x = sx + (cw + spacing) * i
        glass_card(s, x, y, cw, ch, name=f"acid_{i}", accent=col, accent_pos="top")
        add_text(s, x, y + Inches(0.4), cw, Inches(2),
                 letter, size=140, bold=True, color=col,
                 align=PP_ALIGN.CENTER, font="Calibri", letter_spacing=-50)
        add_text(s, x, y + Inches(2.3), cw, Inches(0.5),
                 name, size=18, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, font="Calibri")
        add_text(s, x + Inches(0.2), y + Inches(3), cw - Inches(0.4), Inches(0.9),
                 body, size=11, color=SOFT, align=PP_ALIGN.CENTER, font="Calibri")

    page_number(s, prs, 28, TOTAL)
    add_morph_transition(s)
    set_speaker_notes(s, "Quatre lettres, quatre garanties. C'est la base de toutes les bases de données "
                         "transactionnelles depuis 40 ans. E-Store en bénéficie gratuitement.")
    return s


# ─── ACTE 7 — DEMO PARCOURS UTILISATEUR (slides 29-31) ───────────────
def s29_demo_intro(prs):
    s = blank(prs)
    section_marker(s, prs, "DEMO", 6)
    add_text(s, Inches(0.5), Inches(2.8), Inches(12.3), Inches(1.2),
             "Place à la démo",
             size=72, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Calibri")
    add_text(s, Inches(0.5), Inches(4.2), Inches(12.3), Inches(0.6),
             "8 étapes, 1 parcours, 1 utilisateur",
             size=22, color=CYAN, align=PP_ALIGN.CENTER, italic=True)

    # Compteur 8 etapes
    add_text(s, Inches(0.5), Inches(5.5), Inches(12.3), Inches(0.5),
             "▸ ▸ ▸ ▸ ▸ ▸ ▸ ▸",
             size=32, color=DIM, align=PP_ALIGN.CENTER, letter_spacing=600)

    page_number(s, prs, 29, TOTAL)
    add_morph_transition(s)
    set_speaker_notes(s, "Annonce de la démo. 'Maintenant la pratique. 8 étapes pour couvrir tout le système.'")
    return s


def _journey_slide(prs, visible_until):
    s = blank(prs)
    section_marker(s, prs, "DEMO", 6)
    add_text(s, Inches(0.5), Inches(1.0), Inches(12.3), Inches(0.6),
             "Parcours utilisateur",
             size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Calibri")

    steps = [
        ("01", "Catalogue",      "12 produits · filtre · recherche",     CYAN),
        ("02", "Fiche produit",  "Image · prix · stock · avis",           VIOLET),
        ("03", "AuthGuard",      "Tentative panier → /login",              PINK),
        ("04", "Login",          "user@estore.ma · session JWT",           LIME),
        ("05", "Add to cart",    "Badge passe à 2 (signal)",                AMBER),
        ("06", "Cart edit",      "Quantités · total recalculé",             CYAN),
        ("07", "Checkout",       "Order #1 · CONFIRMED · stock - 2",        VIOLET),
        ("08", "Review 5★",      "MongoDB · createdAt · index",             LIME),
    ]
    sw = Inches(1.5); sh = Inches(2.4); spacing = Inches(0.05)
    total_w = sw * 8 + spacing * 7
    sx = (prs.slide_width - total_w) // 2
    y = Inches(3)
    for i, (n, name, body, col) in enumerate(steps, 1):
        x = sx + (sw + spacing) * (i - 1)
        # Cercle au-dessus
        cy = Inches(2)
        cw_ = Inches(0.6)
        cx_ = x + (sw - cw_) // 2
        nc = s.shapes.add_shape(MSO_SHAPE.OVAL, cx_, cy, cw_, cw_)
        nc.name = f"journey_circle_{i}"
        if i <= visible_until:
            nc.fill.solid(); nc.fill.fore_color.rgb = col
            nc.line.color.rgb = WHITE; nc.line.width = Pt(2)
            tf = nc.text_frame
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            p.text = n
            for r in p.runs:
                r.font.bold = True; r.font.size = Pt(11); r.font.color.rgb = DEEP
        else:
            nc.fill.solid(); nc.fill.fore_color.rgb = INK
            nc.line.color.rgb = SUBTLE; nc.line.width = Pt(1)
            tf = nc.text_frame
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            p.text = "·"
            for r in p.runs:
                r.font.size = Pt(14); r.font.color.rgb = FAINT

        # Carte
        if i <= visible_until:
            glass_card(s, x, y, sw, sh, name=f"journey_card_{i}",
                       accent=col, accent_pos="top")
            add_text(s, x, y + Inches(0.3), sw, Inches(0.5),
                     name, size=12, bold=True, color=WHITE,
                     align=PP_ALIGN.CENTER, font="Calibri")
            add_text(s, x + Inches(0.1), y + Inches(0.85), sw - Inches(0.2), Inches(1.4),
                     body, size=9, color=SOFT,
                     align=PP_ALIGN.CENTER, font="Calibri")
        else:
            ph = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, sw, sh)
            ph.name = f"journey_card_{i}"
            ph.adjustments[0] = 0.06
            ph.fill.solid(); ph.fill.fore_color.rgb = NIGHT
            ph.line.color.rgb = SUBTLE; ph.line.width = Pt(0.5)

    # Ligne de connexion entre cercles
    line = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                   sx + Inches(0.3), Inches(2.3),
                                   sx + total_w - Inches(0.3), Inches(2.3))
    line.line.color.rgb = SUBTLE; line.line.width = Pt(1)

    return s


def s30_journey_half(prs):
    s = _journey_slide(prs, 4)
    page_number(s, prs, 30, TOTAL)
    add_morph_transition(s)
    set_speaker_notes(s, "Premières étapes du parcours (Morph). "
                         "'Catalogue, fiche produit, AuthGuard qui bloque, login.' "
                         "Bascule sur le navigateur pour faire la démo en parallèle.")
    return s


def s31_journey_full(prs):
    s = _journey_slide(prs, 8)
    page_number(s, prs, 31, TOTAL)
    add_morph_transition(s)
    set_speaker_notes(s, "Parcours complet (Morph). 'Login, panier, modification, validation, avis. "
                         "Le tout en moins de 60 secondes côté utilisateur.'")
    return s


# ─── ACTE 8 — METRICS DASHBOARD (slides 32-33) ───────────────────────
def s32_metrics_dashboard(prs):
    s = blank(prs)
    section_marker(s, prs, "RESULTATS", 7)
    add_text(s, Inches(0.5), Inches(1.0), Inches(12.3), Inches(0.6),
             "Tableau de bord",
             size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Calibri")
    add_text(s, Inches(0.5), Inches(1.6), Inches(12.3), Inches(0.4),
             "Les chiffres qui valident le projet",
             size=14, color=CYAN, align=PP_ALIGN.CENTER, italic=True)

    # 4 KPI en haut
    kpis = [
        ("9 / 9",   "Tests verts",         "100%",   1.0,   CYAN,   "kpi_tests"),
        ("13 ko",   "Bundle initial",      "léger",  0.13,  VIOLET, "kpi_bundle"),
        ("< 8 s",   "Démarrage backend",   "rapide", 0.8,   PINK,   "kpi_boot"),
        ("0",       "Vulnerabilité known", "safe",   1.0,   LIME,   "kpi_sec"),
    ]
    cw = Inches(3); ch = Inches(2); spacing = Inches(0.2)
    sx = (prs.slide_width - (cw * 4 + spacing * 3)) // 2
    y = Inches(2.4)
    for i, (val, lbl, sub, pct, col, nm) in enumerate(kpis):
        x = sx + (cw + spacing) * i
        glass_card(s, x, y, cw, ch, name=nm, accent=col, accent_pos="top")
        add_text(s, x, y + Inches(0.3), cw, Inches(0.8),
                 val, size=36, bold=True, color=col,
                 align=PP_ALIGN.CENTER, font="Calibri")
        add_text(s, x, y + Inches(1.1), cw, Inches(0.3),
                 lbl, size=11, color=WHITE, align=PP_ALIGN.CENTER, font="Calibri")
        add_text(s, x + Inches(0.3), y + Inches(1.45), cw - Inches(0.6), Inches(0.15),
                 "", size=8, color=DIM)
        progress_bar(s, x + Inches(0.3), y + Inches(1.65), cw - Inches(0.6), Inches(0.12),
                     pct, col, name=f"{nm}_bar")

    # Bandeau bas: tests detail
    by = Inches(5); bx = Inches(0.5); bw = prs.slide_width - Inches(1); bh = Inches(2.2)
    glass_card(s, bx, by, bw, bh, name="test_panel", accent=CYAN, accent_pos="left")
    add_text(s, bx + Inches(0.3), by + Inches(0.2), Inches(4), Inches(0.4),
             "TESTS UNITAIRES — DETAIL",
             size=11, color=CYAN, font="Consolas", letter_spacing=400)

    # 3 mini-cartes test
    test_x = bx + Inches(0.3); test_y = by + Inches(0.7); test_w = Inches(4); test_h = Inches(1.2)
    test_data = [
        ("ProductServiceTest", "4", "findById · search · pagination", VIOLET),
        ("CartServiceTest",    "2", "stock OK · stock insuffisant",    PINK),
        ("OrderServiceTest",   "3", "panier vide · checkout · trie",   LIME),
    ]
    for i, (name, n, body, col) in enumerate(test_data):
        x = test_x + (test_w + Inches(0.1)) * i
        ph = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 x, test_y, test_w, test_h)
        ph.adjustments[0] = 0.06
        ph.fill.solid(); ph.fill.fore_color.rgb = NIGHT
        ph.line.color.rgb = col; ph.line.width = Pt(1)
        add_text(s, x + Inches(0.2), test_y + Inches(0.15), Inches(2.5), Inches(0.4),
                 name, size=11, bold=True, color=col, font="Consolas")
        add_text(s, x + test_w - Inches(0.7), test_y + Inches(0.1), Inches(0.5), Inches(0.5),
                 n, size=22, bold=True, color=col,
                 align=PP_ALIGN.RIGHT, font="Calibri")
        add_text(s, x + Inches(0.2), test_y + Inches(0.65), test_w - Inches(0.3), Inches(0.5),
                 body, size=9, color=SOFT, font="Calibri")

    page_number(s, prs, 32, TOTAL)
    add_morph_transition(s)
    set_speaker_notes(s, "Le tableau de bord. Les KPI s'animent avec leurs barres de progression. "
                         "'9 tests, 100% verts. Bundle de 13 ko. Démarrage en moins de 8 secondes.'")
    return s


def s33_performance(prs):
    """Graphique d'evolution / performance."""
    s = blank(prs)
    section_marker(s, prs, "PERFORMANCE", 7)
    add_text(s, Inches(0.5), Inches(1.0), Inches(12.3), Inches(0.6),
             "Performance mesurée",
             size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Calibri")

    # Histogramme manuel (5 barres)
    bars = [
        ("Login",      0.24, "240 ms"),
        ("Catalogue",  0.18, "180 ms"),
        ("Add cart",   0.12, "120 ms"),
        ("Checkout",   0.41, "410 ms"),
        ("Reviews",    0.15, "150 ms"),
    ]
    chart_x = Inches(1); chart_y = Inches(2.5); chart_w = Inches(11); chart_h = Inches(4)
    glass_card(s, chart_x, chart_y, chart_w, chart_h, name="perf_chart")

    # Axe Y label
    add_text(s, chart_x + Inches(0.3), chart_y + Inches(0.3), Inches(3), Inches(0.3),
             "TEMPS DE REPONSE (MS) — P50",
             size=10, color=CYAN, font="Consolas", letter_spacing=300)

    bar_w = Inches(1.2); spacing = Inches(0.5); base_y = chart_y + chart_h - Inches(1)
    max_h = Inches(2.3)
    sx = chart_x + Inches(1.5)
    colors = [CYAN, VIOLET, PINK, AMBER, LIME]
    for i, ((name, ratio, label), col) in enumerate(zip(bars, colors)):
        x = sx + (bar_w + spacing) * i
        h = int(max_h * ratio * 2.4)  # Amplifie pour visibilite
        y = base_y - h
        # Barre
        bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  x, y, bar_w, h)
        bar.adjustments[0] = 0.1
        bar.fill.solid(); bar.fill.fore_color.rgb = col
        bar.line.fill.background()
        # Label valeur en haut
        add_text(s, x, y - Inches(0.4), bar_w, Inches(0.3),
                 label, size=11, bold=True, color=col,
                 align=PP_ALIGN.CENTER, font="Consolas")
        # Label en bas
        add_text(s, x, base_y + Inches(0.1), bar_w, Inches(0.3),
                 name, size=10, color=SOFT, align=PP_ALIGN.CENTER, font="Calibri")

    # Ligne de reference
    add_text(s, chart_x + Inches(0.3), base_y + Inches(0.7), chart_w - Inches(0.6), Inches(0.4),
             "▸ TOUTES LES OPERATIONS RESTENT SOUS LE SEUIL DE 500 MS",
             size=10, color=LIME, font="Consolas", letter_spacing=200,
             align=PP_ALIGN.CENTER)

    page_number(s, prs, 33, TOTAL)
    add_morph_transition(s)
    set_speaker_notes(s, "Graphique de performance. 'Toutes les opérations restent sous 500 ms. "
                         "Le checkout est le plus lourd à 410 ms — normal, c'est une transaction multi-table.'")
    return s


# ─── ACTE 9 — DIFFICULTE & SOLUTION (slides 34-35) ───────────────────
def s34_problem(prs):
    s = blank(prs)
    section_marker(s, prs, "INCIDENT", 8)
    add_text(s, Inches(0.5), Inches(1.0), Inches(12.3), Inches(0.6),
             "L'incident Docker",
             size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Calibri")

    # Console d'erreur
    cx = Inches(1.5); cy = Inches(2.5); cw = Inches(10); ch = Inches(2)
    bg = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, cy, cw, ch)
    bg.name = "console"
    bg.adjustments[0] = 0.06
    bg.fill.solid(); bg.fill.fore_color.rgb = RGBColor(0x10, 0x10, 0x18)
    bg.line.color.rgb = CORAL; bg.line.width = Pt(2)
    # Boutons macOS
    for i, col in enumerate([RGBColor(0xFF, 0x60, 0x57),
                              RGBColor(0xFF, 0xBD, 0x2E),
                              RGBColor(0x28, 0xC9, 0x40)]):
        d = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                cx + Inches(0.2) + Inches(0.25) * i,
                                cy + Inches(0.2),
                                Inches(0.18), Inches(0.18))
        d.fill.solid(); d.fill.fore_color.rgb = col
        d.line.fill.background()
    add_text(s, cx + Inches(0.4), cy + Inches(0.7), cw - Inches(0.8), Inches(1.1),
             "$ docker compose up -d\n\n"
             "ERROR: Docker Desktop installation failed.\n"
             "C:\\ProgramData\\DockerDesktop must be owned by an elevated account.",
             size=12, color=CORAL, font="Consolas")

    # Reaction
    add_text(s, Inches(0.5), Inches(5), Inches(12.3), Inches(0.5),
             "▸ J-1. La veille de la soutenance.",
             size=18, color=AMBER, align=PP_ALIGN.CENTER, italic=True)
    add_text(s, Inches(0.5), Inches(5.6), Inches(12.3), Inches(0.5),
             "Il fallait une solution. Maintenant.",
             size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    page_number(s, prs, 34, TOTAL)
    add_morph_transition(s)
    set_speaker_notes(s, "L'incident. 'À la veille de la soutenance, Docker Desktop refuse de s'installer. "
                         "Permission Windows. Pas de temps pour résoudre proprement.'")
    return s


def s35_solution(prs):
    s = blank(prs)
    section_marker(s, prs, "RESOLUTION", 8)
    add_text(s, Inches(0.5), Inches(1.0), Inches(12.3), Inches(0.6),
             "La solution",
             size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Calibri")

    # Console qui se transforme en solution (Morph)
    cx = Inches(1.5); cy = Inches(2.5); cw = Inches(10); ch = Inches(2)
    bg = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, cy, cw, ch)
    bg.name = "console"
    bg.adjustments[0] = 0.06
    bg.fill.solid(); bg.fill.fore_color.rgb = RGBColor(0x10, 0x10, 0x18)
    bg.line.color.rgb = LIME; bg.line.width = Pt(2)
    for i, col in enumerate([RGBColor(0xFF, 0x60, 0x57),
                              RGBColor(0xFF, 0xBD, 0x2E),
                              RGBColor(0x28, 0xC9, 0x40)]):
        d = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                cx + Inches(0.2) + Inches(0.25) * i,
                                cy + Inches(0.2),
                                Inches(0.18), Inches(0.18))
        d.fill.solid(); d.fill.fore_color.rgb = col
        d.line.fill.background()
    add_text(s, cx + Inches(0.4), cy + Inches(0.7), cw - Inches(0.8), Inches(1.1),
             "$ ./mvnw spring-boot:run\n\n"
             "  Started EstoreApplication in 7.823 seconds\n"
             "  ✓ H2 in-memory · ✓ DataSeeder · ✓ JWT ready",
             size=12, color=LIME, font="Consolas")

    # Explication
    add_text(s, Inches(0.5), Inches(5), Inches(12.3), Inches(0.5),
             "▸ Le profil DEV existait depuis le début.",
             size=18, color=CYAN, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(5.6), Inches(12.3), Inches(0.5),
             "Architecture résiliente. Plan B intégré.",
             size=22, bold=True, color=LIME, align=PP_ALIGN.CENTER)

    # Lecon
    add_text(s, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.4),
             "// LA RESILIENCE EST UNE PROPRIETE D'ARCHITECTURE, PAS UNE REACTION D'URGENCE",
             size=10, color=DIM, align=PP_ALIGN.CENTER,
             font="Consolas", letter_spacing=300)

    page_number(s, prs, 35, TOTAL)
    add_morph_transition(s)
    set_speaker_notes(s, "Morph spectaculaire : la console rouge devient verte. "
                         "'Le profil DEV existait depuis le début. H2 in-memory, zéro dépendance externe. "
                         "Plan B intégré. C'est ça, l'architecture résiliente.'")
    return s


# ─── ACTE 10 — ROADMAP & FUTUR (slides 36-37) ────────────────────────
def s36_roadmap(prs):
    s = blank(prs)
    section_marker(s, prs, "ROADMAP", 9)
    add_text(s, Inches(0.5), Inches(1.0), Inches(12.3), Inches(0.6),
             "Et après ?",
             size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Calibri")

    # 4 milestones
    milestones = [
        ("Q1", "Paiement",        "Stripe · CMI · 3DS",         CYAN),
        ("Q2", "Notifications",   "Email · SMS · Push",          VIOLET),
        ("Q3", "Recommandations", "ML collaborative · embedding", PINK),
        ("Q4", "Cloud-native",    "Kubernetes · AWS · CI/CD",     LIME),
    ]
    cw = Inches(2.95); ch = Inches(3.5); spacing = Inches(0.15)
    sx = (prs.slide_width - (cw * 4 + spacing * 3)) // 2
    y = Inches(2.5)
    for i, (q, name, body, col) in enumerate(milestones):
        x = sx + (cw + spacing) * i
        # Card
        glass_card(s, x, y, cw, ch, name=f"ms_{i}", accent=col, accent_pos="top")
        # Quartier
        add_text(s, x, y + Inches(0.3), cw, Inches(0.5),
                 q, size=14, color=col, align=PP_ALIGN.CENTER,
                 font="Consolas", letter_spacing=400, bold=True)
        # Big icon
        icons = ["💳", "🔔", "🤖", "☁️"]
        add_text(s, x, y + Inches(0.9), cw, Inches(1),
                 icons[i], size=64, color=col, align=PP_ALIGN.CENTER)
        # Title
        add_text(s, x + Inches(0.2), y + Inches(2.2), cw - Inches(0.4), Inches(0.5),
                 name, size=20, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, font="Calibri")
        # Body
        add_text(s, x + Inches(0.2), y + Inches(2.8), cw - Inches(0.4), Inches(0.6),
                 body, size=11, color=SOFT,
                 align=PP_ALIGN.CENTER, font="Calibri", italic=True)

    # Ligne directrice en bas (timeline)
    line = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                   Inches(0.5), Inches(6.3),
                                   prs.slide_width - Inches(0.5), Inches(6.3))
    line.line.color.rgb = CYAN; line.line.width = Pt(2)

    add_text(s, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.4),
             "FROM MVP TO CLOUD-NATIVE PLATFORM",
             size=11, color=DIM, align=PP_ALIGN.CENTER,
             font="Consolas", letter_spacing=400)

    page_number(s, prs, 36, TOTAL)
    add_morph_transition(s)
    set_speaker_notes(s, "Roadmap en 4 trimestres. 'On a un MVP solide. La suite : paiement, "
                         "notifications, recommandations ML, déploiement cloud.' Vision long terme.")
    return s


def s37_vision(prs):
    s = blank(prs)
    section_marker(s, prs, "VISION", 9)

    add_text(s, Inches(0.5), Inches(2.5), Inches(12.3), Inches(1.5),
             "E-Store n'est pas la fin.",
             size=58, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Calibri")
    add_text(s, Inches(0.5), Inches(3.6), Inches(12.3), Inches(1),
             "C'est la fondation.",
             size=58, bold=True, color=CYAN, align=PP_ALIGN.CENTER, font="Calibri")

    # Trois piliers de la vision
    pillars = [
        ("PERFORMANT", CYAN),
        ("SCALABLE",   VIOLET),
        ("OUVERT",     PINK),
    ]
    cw = Inches(3); ch = Inches(0.8); spacing = Inches(0.4)
    total_w = cw * 3 + spacing * 2
    sx = (prs.slide_width - total_w) // 2
    y = Inches(5.5)
    for i, (kw, col) in enumerate(pillars):
        x = sx + (cw + spacing) * i
        chip(s, x, y, cw, ch, kw, name=f"pillar_{i}", color=col, size=14)

    page_number(s, prs, 37, TOTAL)
    add_morph_transition(s)
    set_speaker_notes(s, "Vision claire. 'E-Store est une fondation, pas une fin.' "
                         "Trois mots-clés pour résumer : performant, scalable, ouvert.")
    return s


# ─── ACTE 11 — CLOSE (slide 38) ──────────────────────────────────────
def s38_thanks(prs):
    s = blank(prs)
    section_marker(s, prs, "Q & R", 10)

    add_text(s, Inches(0.5), Inches(2), Inches(12.3), Inches(1.5),
             "Merci.",
             size=120, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
             font="Calibri", letter_spacing=-100)

    add_text(s, Inches(0.5), Inches(4.2), Inches(12.3), Inches(0.7),
             "Place à vos questions.",
             size=28, color=CYAN, align=PP_ALIGN.CENTER, italic=True)

    # Chips contacts
    contacts = [
        ("AKRAM BELMOUSSA",       CYAN),
        ("NOUHAILA BEN SOUMANE",  VIOLET),
        ("PR. OMAR ZAHOUR",       PINK),
    ]
    cw = Inches(3.6); ch = Inches(0.7); spacing = Inches(0.3)
    total_w = cw * 3 + spacing * 2
    sx = (prs.slide_width - total_w) // 2
    y = Inches(5.6)
    for i, (name, col) in enumerate(contacts):
        x = sx + (cw + spacing) * i
        chip(s, x, y, cw, ch, name, name=f"contact_{i}",
             color=col, size=10, font="Consolas")

    add_text(s, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.4),
             "// FSBM · UNIVERSITE HASSAN II · 2025-2026",
             size=10, color=DIM, align=PP_ALIGN.CENTER,
             font="Consolas", letter_spacing=600)

    page_number(s, prs, 38, TOTAL)
    add_morph_transition(s)
    set_speaker_notes(s, "Slide finale. Sourire, regarder le jury, tendre la main pour inviter aux questions. "
                         "'Merci pour votre attention. Nous sommes prêts à répondre à vos questions.'")
    return s


# =====================================================================
#  BUILDER
# =====================================================================
def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    builders = [
        # ACTE 1 - OUVERTURE
        s01_black, s02_logo_zoom, s03_title_reveal, s04_tagline, s05_authors,
        # ACTE 2 - VISION
        s06_vision, s07_challenge, s08_metrics_preview, s09_promise,
        # ACTE 3 - ARCHITECTURE
        s10_arch_intro, s11_arch_layer1, s12_arch_layer2, s13_arch_layer3,
        s14_domains, s15_domain_zoom,
        # ACTE 4 - STACK
        s16_stack_galaxy, s17_frontend_stack, s18_backend_stack,
        s19_data_hybrid, s20_devops,
        # ACTE 5 - SECURITE JWT
        s21_jwt_step1, s22_jwt_step3, s23_jwt_step5, s24_jwt_complete,
        s25_security_summary,
        # ACTE 6 - ATOMICITE
        s26_atomic_intro, s27_checkout_code, s28_acid_visual,
        # ACTE 7 - DEMO
        s29_demo_intro, s30_journey_half, s31_journey_full,
        # ACTE 8 - METRICS
        s32_metrics_dashboard, s33_performance,
        # ACTE 9 - INCIDENT
        s34_problem, s35_solution,
        # ACTE 10 - VISION
        s36_roadmap, s37_vision,
        # ACTE 11 - CLOSE
        s38_thanks,
    ]
    for fn in builders:
        fn(prs)

    out = os.path.join(OUT_DIR, "E-Store-Keynote.pptx")
    prs.save(out)
    return out


if __name__ == "__main__":
    print("=" * 70)
    print("E-STORE | KEYNOTE EDITION")
    print("Generation premium avec transitions Morph")
    print("=" * 70)
    out = build()
    size_kb = os.path.getsize(out) // 1024
    print(f"  [OK] {os.path.basename(out):50s} ({size_kb} ko)")
    print("=" * 70)
    print(f"38 slides chainees avec Morph dans : {out}")
    print()
    print("Pour optimiser le rendu :")
    print("  1. Ouvrir dans PowerPoint 2016 ou plus recent")
    print("  2. Verifier que les transitions Morph sont actives (onglet Transitions)")
    print("  3. Mode presentation : touche F5")
    print("  4. Apple Keynote / Google Slides : Morph -> 'Magic Move' (a re-mapper)")
